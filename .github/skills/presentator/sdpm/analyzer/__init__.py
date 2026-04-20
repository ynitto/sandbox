# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Template analysis: extract layout info, theme colors, fonts, and color usage."""
from pathlib import Path

from sdpm.builder import PPTXBuilder
from sdpm.utils.io import write_json, read_json

from pptx import Presentation


def analyze_template(template_path: Path):
    """Analyze template and return structured result.

    Returns dict with layouts (name + notes), theme colors, fonts, color usage,
    table styles, slide size.
    Color usage is loaded from cache if available; caller must populate cache via
    cache_color_usage() if missing.
    """
    prs = Presentation(str(template_path))
    emu = 6350

    # Build layout→notes mapping from slides
    layout_notes = {}
    for slide in prs.slides:
        name = slide.slide_layout.name or "(unnamed)"
        if name not in layout_notes:
            notes = ""
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            layout_notes[name] = notes

    layouts = []
    for name, notes in layout_notes.items():
        entry = {"name": name}
        if notes:
            entry["notes"] = notes
        layouts.append(entry)

    theme_colors = extract_theme_colors(template_path)
    fonts = extract_fonts(template_path)
    color_usage = _load_color_usage_cache(template_path)
    slide_size = {
        "width": int(prs.slide_width / emu),
        "height": int(prs.slide_height / emu),
    }

    return {
        "slide_size": slide_size,
        "layouts": layouts,
        "theme_colors": theme_colors,
        "color_usage": color_usage,
        "fonts": fonts,
    }


def get_layout_placeholders(template_path: Path, layout_name: str):
    """Get placeholder details and notes for a layout directly from pptx."""
    prs = Presentation(str(template_path))
    emu = 6350

    # Find layout
    layout_obj = None
    for master in prs.slide_masters:
        for sl in master.slide_layouts:
            if sl.name == layout_name:
                layout_obj = sl
                break
        if layout_obj:
            break
    if not layout_obj:
        return None

    # Get notes from first slide using this layout
    notes = ""
    for slide in prs.slides:
        if slide.slide_layout.name == layout_name:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            break

    placeholders = []
    for ph in layout_obj.placeholders:
        if ph.placeholder_format.type in (13, 15, 16):
            continue
        info = {
            "idx": ph.placeholder_format.idx,
            "x": int(ph.left / emu),
            "y": int(ph.top / emu),
            "width": int(ph.width / emu),
            "height": int(ph.height / emu),
        }
        # Font size from first defRPr in placeholder XML
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        for defRPr in ph._element.iter(f'{ns_a}defRPr'):
            sz = defRPr.get('sz')
            if sz:
                info["fontSize"] = int(sz) / 100  # hundredths of pt → pt
                break
        desc = ph.text.strip().replace('\x0b', ' / ') if ph.text else ""
        if desc:
            info["description"] = desc
        placeholders.append(info)

    return {"name": layout_name, "notes": notes, "placeholders": placeholders}


def extract_fonts(template_path: Path) -> dict:
    """Extract default fonts from template theme XML."""
    import zipfile
    from lxml import etree

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    with zipfile.ZipFile(str(template_path)) as z:
        for name in z.namelist():
            if "theme" in name and name.endswith(".xml"):
                root = etree.fromstring(z.read(name))
                minor = root.find(".//a:minorFont", ns)
                if minor is not None:
                    latin = minor.find("a:latin", ns)
                    ea = minor.find("a:ea", ns)
                    return {
                        "halfwidth": latin.get("typeface") if latin is not None else None,
                        "fullwidth": ea.get("typeface") if ea is not None else None,
                    }
    return {"halfwidth": None, "fullwidth": None}


def extract_theme_colors(template_path: Path):
    """Extract theme colors from template (clrMap-aware)."""
    colors, _ = PPTXBuilder._extract_theme_colors(template_path)
    result = {}
    role_map = {"text": "text", "background": "background", "text2": "text2", "background2": "background2"}
    for key, role in role_map.items():
        if colors.get(key):
            result[role] = colors[key]
    for i in range(1, 7):
        key = f"accent{i}"
        if colors.get(key):
            result[key] = colors[key]
    return result


def _color_usage_cache_path(template_path: Path) -> Path:
    project_root = Path(__file__).resolve().parent.parent.parent
    return project_root / ".cache" / "templates" / template_path.stem / "color-usage.json"


def _load_color_usage_cache(template_path: Path):
    """Load color usage from cache if fresh, else return empty list."""
    cache_file = _color_usage_cache_path(template_path)
    if cache_file.exists():
        template_mtime = template_path.stat().st_mtime
        if cache_file.stat().st_mtime >= template_mtime:
            return read_json(cache_file)
    return []


def cache_color_usage(template_path: Path, usage: list):
    """Save color usage data to cache."""
    cache_file = _color_usage_cache_path(template_path)
    cache_file.parent.mkdir(parents=True, exist_ok=True)
    write_json(cache_file, usage, suffix="\n")


def cache_preview_pngs(template_path: Path, preview_dir: Path):
    """Copy preview PNGs to cache directory."""
    import shutil
    cache_dir = _color_usage_cache_path(template_path).parent
    cache_dir.mkdir(parents=True, exist_ok=True)
    for png in sorted(preview_dir.glob("page*.png")):
        shutil.copy2(png, cache_dir / png.name)


def extract_color_usage_from_pngs(preview_dir: Path):
    """Extract color usage from preview PNGs via pixel analysis."""
    from PIL import Image
    from collections import Counter

    pngs = sorted(preview_dir.glob("page*.png"))
    if not pngs:
        return []

    colors = Counter()
    for png in pngs:
        img = Image.open(png).convert("RGB")
        img_q = img.quantize(colors=32, method=Image.Quantize.MEDIANCUT).convert("RGB")
        w, h = img_q.size
        for y in range(0, h, 2):
            for x in range(0, w, 2):
                colors[img_q.getpixel((x, y))] += 1

    total = sum(colors.values())
    result = []
    for (r, g, b), count in colors.most_common(20):
        pct = count / total * 100
        if pct < 0.2:
            break
        result.append({
            "color": f"#{r:02X}{g:02X}{b:02X}",
            "percentage": round(pct, 1),
        })
    return result
