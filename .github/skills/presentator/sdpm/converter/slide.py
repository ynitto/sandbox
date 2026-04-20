# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Slide extraction."""
import sys
import defusedxml
defusedxml.defuse_stdlib()

import xml.etree.ElementTree as ET




from .constants import _NS, EMU_PER_PX, _extract_autofit_props, _position_diff
from .color import _resolve_scheme_color
from .text import _extract_styled_text
from .elements import (extract_textbox_element, extract_picture_element, _dispatch_shape)

def detect_layout(slide):
    """Detect slide layout type."""
    # Check layout name first
    try:
        name = slide.slide_layout.name.lower()
        if 'thank' in name:
            return "thankyou"
        if 'agenda' in name:
            return "agenda"
        if 'section' in name:
            if 'sub' in name:
                return "subsection"
            return "section"
        if 'left line' in name:
            return "content"
    except Exception:
        pass

    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type == 3:  # SUBTITLE
            return "title"
        if ph_type == 4:  # CENTER_TITLE subtitle
            return "title"

    has_title = False
    has_content = False
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type == 1:
                has_title = True
            elif ph_type in (2, 7):  # BODY, OBJECT
                has_content = True

    # Distinguish content from title_only by checking for left line decoration
    if has_title and not has_content:
        try:
            name = slide.slide_layout.name.lower()
            if 'content' in name or 'bulleted' in name or 'subtitle' in name:
                return "content"
        except Exception:
            pass
        # Check for left accent line (content layout marker)
        for shape in slide.shapes:
            if not shape.is_placeholder and shape.left < 50 * 6350:
                if shape.width < 20 * 6350 and shape.height > 500 * 6350:
                    return "content"
        return "title_only"

    if has_title and has_content:
        return "content"
    return "title_only"

def _resolve_inherited_styles(elements, slide, theme_colors, color_mapping):
    """Resolve inherited text color, fontSize, fontFamily from shape's lstStyle/defRPr."""
    if not theme_colors or not color_mapping:
        return
    for elem in elements:
        if elem.get("type") not in ("textbox", "shape"):
            continue
        if not elem.get("text") and not elem.get("paragraphs") and not elem.get("items"):
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if round(shape.left / EMU_PER_PX) == elem.get("x") and round(shape.top / EMU_PER_PX) == elem.get("y"):
                elem.update(_extract_autofit_props(shape))
                defRPr = shape.text_frame._txBody.find(f'.//{{{_NS["a"]}}}lstStyle//{{{_NS["a"]}}}defRPr')
                if defRPr is not None:
                    if not elem.get("_lstStyle"):
                        sz = defRPr.get('sz')
                        if sz and not elem.get("fontSize"):
                            all_sizes = [r.font.size for p in shape.text_frame.paragraphs for r in p.runs if r.font.size]
                            if not all_sizes:
                                elem["fontSize"] = int(sz) // 100
                    if not elem.get("fontColor") and not elem.get("_phIdx"):
                        # Skip for placeholders (theme handles color) and inline-colored text
                        all_runs_have_color = all(
                            r.font.color.type is not None
                            for p in shape.text_frame.paragraphs for r in p.runs if r.text.strip()
                        )
                        if not all_runs_have_color:
                            scheme = defRPr.find(f'.//{{{_NS["a"]}}}schemeClr')
                            if scheme is not None:
                                resolved = _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                                if resolved:
                                    elem["fontColor"] = resolved
                if not elem.get("fontFamily"):
                    run_fonts = set()
                    all_have_font = True
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.name:
                                run_fonts.add(r.font.name)
                            elif r.text.strip():
                                all_have_font = False
                    if len(run_fonts) == 1 and all_have_font:
                        elem["fontFamily"] = run_fonts.pop()
                break

def extract_slide(slide, theme_colors=None, color_mapping=None, theme_styles=None, master_idx=0, output_dir=None, slide_idx=0, pptx_path=None, use_layout_names=False, builder_text_color=None):
    """Extract slide content to dict."""
    if use_layout_names:
        layout_name = slide.slide_layout.name if slide.slide_layout.name else f"layout-{master_idx+1:02d}"
    else:
        layout_name = detect_layout(slide)
    slide_dict = {
        "layout": layout_name,
        "masterIndex": master_idx
    }
    
    # Extract slide background (if different from layout)
    try:
        bg = slide.background._element.find(f'{{{_NS["p"]}}}bg')
        if bg is not None:
            bgPr = bg.find(f'{{{_NS["p"]}}}bgPr')
            if bgPr is not None:
                solid = bgPr.find(f'{{{_NS["a"]}}}solidFill')
                if solid is not None:
                    srgb = solid.find(f'{{{_NS["a"]}}}srgbClr')
                    scheme = solid.find(f'{{{_NS["a"]}}}schemeClr')
                    if srgb is not None:
                        slide_dict["background"] = f"#{srgb.get('val')}"
                    elif scheme is not None:
                        from .color import _resolve_color_with_transforms
                        resolved = _resolve_color_with_transforms(scheme, theme_colors, color_mapping)
                        if resolved:
                            slide_dict["background"] = resolved
    except Exception:
        pass
    
    # Extract placeholders by idx
    placeholders = {}
    for shape in slide.shapes:
        if not shape.is_placeholder or not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        idx = shape.placeholder_format.idx
        ph_type = shape.placeholder_format.type
        # Skip system placeholders
        if ph_type in (10, 11, 12, 13, 14, 15, 16):  # DATE, SLIDE_NUMBER, FOOTER, HEADER, etc.
            continue
        str_idx = str(idx)
        if str_idx in placeholders:
            continue
        text = shape.text
        # Styled text for title
        if shape == slide.shapes.title:
            default_tc = None
            if color_mapping:
                tx1 = color_mapping.get('tx1', 'dk1')
                if theme_colors and tx1 in theme_colors:
                    default_tc = theme_colors[tx1]
            styled = _extract_styled_text(shape.text_frame.paragraphs[0].runs, theme_colors, color_mapping, default_text_color=default_tc, paragraph=shape.text_frame.paragraphs[0]) if shape.text_frame.paragraphs else text
            text = styled if styled != shape.text else shape.text
        val = text
        # Check for explicit font size (all runs same size)
        all_runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        all_sizes = [r.font.size for r in all_runs]
        common_size = all_sizes[0] if all_sizes and all(s is not None for s in all_sizes) and len(set(all_sizes)) == 1 else None
        if common_size:
            from pptx.util import Pt
            val = {"text": text, "fontSize": round(common_size / Pt(1))}
        # Preserve position if different from layout
        try:
            for lph in slide.slide_layout.placeholders:
                if lph.placeholder_format.idx == idx:
                    diff = _position_diff(shape, lph)
                    if diff:
                        if isinstance(val, str):
                            val = {"text": val}
                        val.update(diff)
                    break
        except Exception:
            pass
        # Preserve alignment if explicit
        try:
            align = shape.text_frame.paragraphs[0].alignment
            if align is not None:
                align_map = {1: "left", 2: "center", 3: "right", 4: "justify"}
                a = align_map.get(int(align))
                if a:
                    if isinstance(val, str):
                        val = {"text": val}
                    val["align"] = a
        except Exception:
            pass
        placeholders[str_idx] = val
    if placeholders:
        slide_dict["placeholders"] = placeholders
    
    # Extract speaker notes
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_dict["notes"] = notes_text
    except Exception:
        pass
    
    # Extract content from placeholders as textboxes (to preserve position)
    elements = []
    
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (2, 7, 13):
            if shape.has_text_frame and shape.text.strip():
                elem = extract_textbox_element(shape, theme_colors, color_mapping, theme_styles, is_placeholder=True, builder_text_color=builder_text_color)
                if elem and (elem.get("text", "").strip() or elem.get("paragraphs")):
                    elem["_phIdx"] = shape.placeholder_format.idx
                    # Preserve lstStyle from layout placeholder or shape itself
                    layout_ph = None
                    try:
                        for lph in slide.slide_layout.placeholders:
                            if lph.placeholder_format.idx == shape.placeholder_format.idx:
                                layout_ph = lph
                                break
                    except Exception:
                        pass
                    src = layout_ph if layout_ph is not None else shape
                    lstStyle = src._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
                    if lstStyle is not None and len(lstStyle) > 0:
                        # Merge spcBef/spcAft from master bodyStyle if missing
                        try:
                            master = slide.slide_layout.slide_master
                            txStyles = master.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}txStyles')
                            bodyStyle = txStyles.find(f'{{{_NS["p"]}}}bodyStyle') if txStyles is not None else None
                            if bodyStyle is not None:
                                for lvl in lstStyle:
                                    tag = lvl.tag.split('}')[1]
                                    master_lvl = bodyStyle.find(f'{{{_NS["a"]}}}{tag}')
                                    if master_lvl is not None:
                                        for spc_tag in ('spcBef', 'spcAft', 'lnSpc'):
                                            if lvl.find(f'{{{_NS["a"]}}}{spc_tag}') is None:
                                                master_spc = master_lvl.find(f'{{{_NS["a"]}}}{spc_tag}')
                                                if master_spc is not None:
                                                    from copy import deepcopy
                                                    lvl.append(deepcopy(master_spc))
                        except Exception:
                            pass
                        elem["_lstStyle"] = ET.tostring(lstStyle, encoding='unicode')
                    # Resolve color if all runs inherit
                    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
                    if runs and any(r.font.color.type is None for r in runs):
                        # Skip fontColor for _phIdx elements (theme handles it)
                        pass
                    # Apply spcAft from lstStyle to paragraphs (textbox doesn't inherit lstStyle spacing)
                    if elem.get("paragraphs") and elem.get("_lstStyle"):
                        try:
                            _lst = ET.fromstring(elem["_lstStyle"])
                            _lvl1 = _lst.find(f'{{{_NS["a"]}}}lvl1pPr')
                            if _lvl1 is not None:
                                _spcAft = _lvl1.find(f'{{{_NS["a"]}}}spcAft/{{{_NS["a"]}}}spcPts')
                                if _spcAft is not None:
                                    sa_val = int(_spcAft.get('val'))
                                    for p in elem["paragraphs"]:
                                        if isinstance(p, dict) and p.get("spaceAfter") is None:
                                            p["spaceAfter"] = sa_val
                                _lnSpc = _lvl1.find(f'{{{_NS["a"]}}}lnSpc/{{{_NS["a"]}}}spcPct')
                                if _lnSpc is not None:
                                    ls_val = int(_lnSpc.get('val'))
                                    for p in elem["paragraphs"]:
                                        if isinstance(p, dict) and p.get("lineSpacingPct") is None:
                                            p["lineSpacingPct"] = ls_val
                        except Exception:
                            pass
                    elem.update(_extract_autofit_props(shape))
                    # Resolve bodyPr margins from master placeholder
                    if not elem.get("marginLeft") and not elem.get("marginRight"):
                        try:
                            for mph in slide.slide_layout.slide_master.placeholders:
                                if mph.placeholder_format.type in (2, 7):  # BODY, OBJECT
                                    mbp = mph._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                                    if mbp is not None:
                                        lIns = mbp.get('lIns')
                                        rIns = mbp.get('rIns')
                                        if lIns is not None:
                                            elem["marginLeft"] = round(int(lIns) / EMU_PER_PX)
                                        if rIns is not None:
                                            elem["marginRight"] = round(int(rIns) / EMU_PER_PX)
                                    break
                        except Exception:
                            pass
                    elements.append(elem)
    
    # Extract placeholder images (PICTURE type or pic element)
    img_counter = 0
    for shape in slide.shapes:
        if shape.is_placeholder and (shape.placeholder_format.type == 18 or shape._element.tag.endswith('}pic')):
            try:
                img = shape.image
                ext = img.content_type.split('/')[-1].replace('jpeg', 'jpg')
                img_name = f"slide{slide_idx+1}_img{img_counter}.{ext}"
                if output_dir:
                    img_dir = output_dir / "images"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    (img_dir / img_name).write_bytes(img.blob)
                elem = {
                    "type": "image",
                    "src": f"images/{img_name}",
                    "x": round(shape.left / EMU_PER_PX),
                    "y": round(shape.top / EMU_PER_PX),
                    "width": round(shape.width / EMU_PER_PX),
                    "height": round(shape.height / EMU_PER_PX),
                }
                elements.append(elem)
                img_counter += 1
            except Exception:
                pass
    
    # Extract non-placeholder elements
    # img_counter continues from placeholder images
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        
        try:
            elem, img_counter = _dispatch_shape(shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter, builder_text_color=builder_text_color, pptx_path=pptx_path)
            if elem:
                # Skip tiny invisible connector stubs (e.g. w=20, h=0 decorative connectors)
                if elem.get("type") == "line":
                    dx = abs(elem.get("x2", 0) - elem.get("x1", 0))
                    dy = abs(elem.get("y2", 0) - elem.get("y1", 0))
                    if dx <= 30 and dy <= 30:
                        continue
                elements.append(elem)
        except Exception as e:
            print(f"Warning: Failed to extract shape {shape.name}: {e}", file=sys.stderr)
    
    # Extract mc:AlternateContent fallback images (3D models, etc.)
    _mc_ns = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
    for ac in slide.element.find(f'.//{{{_NS["p"]}}}spTree').findall(f'{{{_mc_ns}}}AlternateContent'):
        fb = ac.find(f'{{{_mc_ns}}}Fallback')
        if fb is None:
            continue
        pic = fb.find(f'{{{_NS["p"]}}}pic')
        if pic is None:
            continue
        try:
            from pptx.shapes.picture import Picture
            pic_shape = Picture(pic, slide.part)
            elem = extract_picture_element(pic_shape, output_dir, slide_idx, img_counter, theme_colors, color_mapping)
            if elem:
                img_counter += 1
                elements.append(elem)
        except Exception:
            pass

    # Resolve inherited text color and fontSize from shape's lstStyle/defRPr
    _resolve_inherited_styles(elements, slide, theme_colors, color_mapping)

    if elements:
        slide_dict["elements"] = elements
    
    return slide_dict
