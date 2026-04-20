# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Placeholder handling: title, content, section fills."""
from pptx.dml.color import RGBColor
from sdpm.utils.text import normalize_spacing


class PlaceholderMixin:
    """Mixin providing placeholder fill methods."""

    @staticmethod
    def _resolve_placeholder(val):
        """Split placeholder value into (text, style_dict)."""
        if isinstance(val, dict):
            return val.get("text", ""), val
        return val, {}
    
    def _apply_placeholder_style(self, placeholder, text, style):
        """Apply Styled Text + style properties to a placeholder."""
        from pptx.enum.text import PP_ALIGN
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        self._apply_styled_text(p, normalize_spacing(text),
                                no_default_color=not style.get("fontColor"),
                                no_default_font=True)
        # Apply bold/italic/fontColor on all runs
        for run in p.runs:
            if style.get("bold"):
                run.font.bold = True
            if style.get("italic"):
                run.font.italic = True
            fc = style.get("fontColor")
            if fc:
                run.font.color.rgb = RGBColor.from_string(fc.lstrip('#'))
        # Alignment
        align = style.get("align")
        if align:
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                           "right": PP_ALIGN.RIGHT}.get(align)
    
    def _fill_placeholder(self, placeholder, val):
        """Resolve and apply a placeholder value (string or object)."""
        text, style = self._resolve_placeholder(val)
        self._apply_placeholder_style(placeholder, text, style)
    
    def _apply_placeholder_position(self, ph, val):
        """Apply custom position/size to placeholder if specified in val."""
        _, style = self._resolve_placeholder(val)
        abs_keys = ("_x", "_y", "_width", "_height")
        off_keys = ("offsetX", "offsetY", "offsetWidth", "offsetHeight")
        has_absolute = any(style.get(k) is not None for k in abs_keys)
        has_offset = any(style.get(k) is not None for k in off_keys)
        if not has_absolute and not has_offset:
            return
        from pptx.oxml.ns import qn
        from lxml import etree
        # Get layout defaults
        lx, ly, lw, lh = 0, 0, 0, 0
        try:
            for lph in ph.part.slide_layout.placeholders:
                if lph.placeholder_format.idx == ph.placeholder_format.idx:
                    lx, ly, lw, lh = lph.left, lph.top, lph.width, lph.height
                    break
        except Exception:
            pass
        if has_absolute:
            x = int(style["_x"] * self.EMU_PER_PX) if style.get("_x") is not None else lx
            y = int(style["_y"] * self.EMU_PER_PX) if style.get("_y") is not None else ly
            cx = int(style["_width"] * self.EMU_PER_PX) if style.get("_width") is not None else lw
            cy = int(style["_height"] * self.EMU_PER_PX) if style.get("_height") is not None else lh
        else:
            x = lx + int(style.get("offsetX", 0) * self.EMU_PER_PX)
            y = ly + int(style.get("offsetY", 0) * self.EMU_PER_PX)
            cx = lw + int(style.get("offsetWidth", 0) * self.EMU_PER_PX)
            cy = lh + int(style.get("offsetHeight", 0) * self.EMU_PER_PX)
        spPr = ph._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ph._element, qn('p:spPr'))
        xfrm = spPr.find(qn('a:xfrm'))
        if xfrm is None:
            xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        off = xfrm.find(qn('a:off'))
        ext = xfrm.find(qn('a:ext'))
        if off is None:
            off = etree.SubElement(xfrm, qn('a:off'))
        if ext is None:
            ext = etree.SubElement(xfrm, qn('a:ext'))
        off.set('x', str(x))
        off.set('y', str(y))
        ext.set('cx', str(cx))
        ext.set('cy', str(cy))

    def _fill_title(self, slide, d):
        self._fill_placeholders(slide, d)
    
    def _fill_content(self, slide, d):
        self._fill_placeholders(slide, d)
    
    def _set_bullet(self, paragraph, char=None):
        """Set bullet point formatting from master's bodyStyle definition."""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        level = paragraph.level or 0
        style = self._list_styles.get(level, {})
        
        pPr = paragraph._element.get_or_add_pPr()
        if style.get('marL'):
            pPr.set('marL', style['marL'])
        if style.get('indent'):
            pPr.set('indent', style['indent'])
        
        bu_font = style.get('buFont', 'Arial')
        bu_char = char or style.get('buChar', '•')
        font_elem = etree.SubElement(pPr, qn('a:buFont'))
        font_elem.set('typeface', bu_font)
        char_elem = etree.SubElement(pPr, qn('a:buChar'))
        char_elem.set('char', bu_char)
    
    def _set_numbering(self, paragraph, numbering_type='arabicPeriod'):
        """Set numbering formatting, using master's marL/indent for indentation."""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        level = paragraph.level or 0
        style = self._list_styles.get(level, {})
        
        pPr = paragraph._element.get_or_add_pPr()
        if style.get('marL'):
            pPr.set('marL', style['marL'])
        if style.get('indent'):
            pPr.set('indent', style['indent'])
        
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
        buAutoNum.set('type', numbering_type)
    
    def _fill_section(self, slide, d):
        self._fill_placeholders(slide, d)
    
    def _fill_title_only(self, slide, d):
        self._set_or_remove_title(slide, d)

    def _fill_placeholders(self, slide, d):
        """Fill placeholders from idx-based placeholders dict."""
        ph_dict = d.get("placeholders", {})
        # Title removal: empty string at idx 0 removes the title placeholder
        title_idx = "0"
        if title_idx in ph_dict:
            text, _ = self._resolve_placeholder(ph_dict[title_idx])
            if not text and slide.shapes.title:
                sp = slide.shapes.title._element
                sp.getparent().remove(sp)
        # Fill all placeholders by idx
        for ph in slide.placeholders:
            idx = str(ph.placeholder_format.idx)
            if idx in ph_dict:
                val = ph_dict[idx]
                text, _ = self._resolve_placeholder(val)
                if not text and idx == title_idx:
                    continue  # already removed above
                self._fill_placeholder(ph, val)
                self._apply_placeholder_position(ph, val)
        # Items into first BODY/OBJECT placeholder
        if "items" in d:
            for ph in slide.placeholders:
                if ph.placeholder_format.type in (2, 7):
                    tf = ph.text_frame
                    for i, item in enumerate(d["items"]):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.level = 0
                        self._set_bullet(p)
                        self._apply_styled_text(p, item)
                    break
    
    def _set_or_remove_title(self, slide, d):
        """Set title text or remove title placeholder if empty."""
        ph_dict = d.get("placeholders", {})
        val = ph_dict.get("0")
        if val is None:
            return
        text, style = self._resolve_placeholder(val)
        if text and slide.shapes.title:
            self._apply_placeholder_style(slide.shapes.title, text, style)
        elif not text and slide.shapes.title:
            sp = slide.shapes.title._element
            sp.getparent().remove(sp)
    
