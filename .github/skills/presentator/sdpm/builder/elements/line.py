# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Line and polyline elements."""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Pt
from sdpm.schema.defaults import ELEMENT_DEFAULTS

_DEFAULTS = ELEMENT_DEFAULTS["line"]


class LineMixin:
    """Mixin providing line element methods."""

    def _add_polyline(self, slide, elem, points):
        """Add polyline as freeform shape with arrow support."""
        from lxml import etree
        from pptx.oxml.ns import qn
    
        default_color = self.theme_colors["text"]
        color = elem.get("color", default_color)
        line_width = elem.get("lineWidth", _DEFAULTS["lineWidth"])
        hex_c = color.lstrip("#")
    
        # Bounding box
        xs = [p[0] for p in points]
        ys = [p[1] for p in points]
        min_x, max_x = min(xs), max(xs)
        min_y, max_y = min(ys), max(ys)
        bw = max(max_x - min_x, 1)
        bh = max(max_y - min_y, 1)
    
        x_emu = self._px_to_emu(min_x)
        y_emu = self._px_to_emu(min_y)
        w_emu = self._px_to_emu(bw)
        h_emu = self._px_to_emu(bh)
    
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_emu, y_emu, w_emu, h_emu)
    
        # Build custom geometry polyline
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        cust_geom = etree.SubElement(etree.Element('dummy'), f'{{{ns_a}}}custGeom')
        etree.SubElement(cust_geom, f'{{{ns_a}}}avLst')
        etree.SubElement(cust_geom, f'{{{ns_a}}}gdLst')
        etree.SubElement(cust_geom, f'{{{ns_a}}}ahLst')
        etree.SubElement(cust_geom, f'{{{ns_a}}}cxnLst')
        rect = etree.SubElement(cust_geom, f'{{{ns_a}}}rect')
        rect.set('l', 'l')
        rect.set('t', 't')
        rect.set('r', 'r')
        rect.set('b', 'b')
        path_lst = etree.SubElement(cust_geom, f'{{{ns_a}}}pathLst')
        path_el = etree.SubElement(path_lst, f'{{{ns_a}}}path')
        path_el.set('w', str(w_emu))
        path_el.set('h', str(h_emu))
    
        for i, p in enumerate(points):
            px = round((p[0] - min_x) / bw * w_emu) if bw > 1 else 0
            py = round((p[1] - min_y) / bh * h_emu) if bh > 1 else 0
            if i == 0:
                mv = etree.SubElement(path_el, f'{{{ns_a}}}moveTo')
                pt = etree.SubElement(mv, f'{{{ns_a}}}pt')
            else:
                ln = etree.SubElement(path_el, f'{{{ns_a}}}lnTo')
                pt = etree.SubElement(ln, f'{{{ns_a}}}pt')
            pt.set('x', str(px))
            pt.set('y', str(py))
    
        # Replace preset geometry
        sp_pr = shape._element.spPr
        for prst in sp_pr.findall(qn('a:prstGeom')):
            sp_pr.remove(prst)
        xfrm = sp_pr.find(qn('a:xfrm'))
        if xfrm is not None:
            xfrm.addnext(cust_geom)
    
        # Remove default style
        for style in shape._element.findall(qn('p:style')):
            shape._element.remove(style)
    
        # No fill
        shape.fill.background()
    
        # Line style
        shape.line.color.rgb = RGBColor(int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16))
        shape.line.width = Pt(line_width)
    
        # Arrow heads (polyline has no flip — no swap needed)
        ln = sp_pr.find(f'{{{ns_a}}}ln')
        if ln is None:
            ln = etree.SubElement(sp_pr, f'{{{ns_a}}}ln')
        arrow_start = elem.get("arrowStart")
        arrow_end = elem.get("arrowEnd")
        if arrow_end:
            te = etree.SubElement(ln, f'{{{ns_a}}}tailEnd')
            te.set('type', arrow_end)
        if arrow_start:
            he = etree.SubElement(ln, f'{{{ns_a}}}headEnd')
            he.set('type', arrow_start)
    
    def _add_line(self, slide, elem):
        """Add line/connector to slide."""
        points = elem.get("points")
        if points and len(points) >= 2:
            return self._add_polyline(slide, elem, points)

        # Require x1/y1/x2/y2 format (x/y/width/height is not supported — use linter to catch)
        if "x1" not in elem:
            return  # No coordinates — skip silently
        x1_emu = self._px_to_emu(elem["x1"])
        y1_emu = self._px_to_emu(elem["y1"])
        x2_emu = self._px_to_emu(elem["x2"])
        y2_emu = self._px_to_emu(elem["y2"])

        # Determine connector type
        connector_type_str = elem.get("connectorType", _DEFAULTS["connectorType"])
        connector_type_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curved": MSO_CONNECTOR.CURVE,
        }
        connector_type = connector_type_map.get(connector_type_str, MSO_CONNECTOR.STRAIGHT)

        # Add connector (line)
        connector = slide.shapes.add_connector(
            connector_type,
            x1_emu, y1_emu,
            x2_emu, y2_emu
        )

        # V-H-V elbow: rotate connector so it starts vertically
        elbow_start = elem.get("elbowStart", "horizontal")
        is_vhv = elbow_start == "vertical" and connector_type_str == "elbow"
        if is_vhv:
            from lxml import etree
            dx = x2_emu - x1_emu
            dy = y2_emu - y1_emu
            # Swap width/height for rotated bounding box
            cx = abs(dy)
            cy = abs(dx)
            mid_x = (x1_emu + x2_emu) // 2
            mid_y = (y1_emu + y2_emu) // 2
            off_x = mid_x - cx // 2
            off_y = mid_y - cy // 2
            xfrm = connector._element.spPr.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                from pptx.oxml.ns import qn
                xfrm.find(qn('a:off')).set('x', str(off_x))
                xfrm.find(qn('a:off')).set('y', str(off_y))
                xfrm.find(qn('a:ext')).set('cx', str(cx))
                xfrm.find(qn('a:ext')).set('cy', str(cy))
                # Clear auto-set flips from python-pptx, apply correct combo
                for attr in ('flipH', 'flipV'):
                    if attr in xfrm.attrib:
                        del xfrm.attrib[attr]
                connector.rotation = 270
                if dy > 0:
                    xfrm.set('flipH', '1')
                if dx < 0:
                    xfrm.set('flipV', '1')

        # AWS recommended defaults (theme-aware)
        default_color = self.theme_colors["text"]
        default_line_width = 1.25

        # Apply exact preset geometry if available
        preset = elem.get("preset")
        if preset:
            try:
                from lxml import etree
                sp_pr = connector._element.spPr
                prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                if prst_geom is not None:
                    prst_geom.set('prst', preset)
            except Exception:
                pass

        # Apply elbow adjustments (waypoints)
        adjustments = elem.get("adjustments")
        if adjustments and connector_type_str in ("elbow", "curved"):
            try:
                from lxml import etree
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                sp_pr = connector._element.spPr
                prst_geom = sp_pr.find(f'.//{{{ns}}}prstGeom')
                if prst_geom is not None:
                    # Use explicit preset if given, otherwise auto-select
                    if not preset:
                        if len(adjustments) >= 3:
                            prst_geom.set('prst', 'bentConnector5')
                        elif len(adjustments) >= 2:
                            prst_geom.set('prst', 'bentConnector4')
                        else:
                            prst_geom.set('prst', 'bentConnector3')
                    avLst = prst_geom.find(f'{{{ns}}}avLst')
                    if avLst is None:
                        avLst = etree.SubElement(prst_geom, f'{{{ns}}}avLst')
                    for i, val in enumerate(adjustments):
                        gd = etree.SubElement(avLst, f'{{{ns}}}gd')
                        gd.set('name', f'adj{i + 1}')
                        gd.set('fmla', f'val {int(val * 100000)}')
            except Exception:
                pass

        # Resolve arrow properties: arrowStart/arrowEnd
        arrow_start = elem.get("arrowStart")
        arrow_end = elem.get("arrowEnd")

        if arrow_start or arrow_end:
            from lxml import etree
            ln = connector._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is None:
                ln = etree.SubElement(connector._element.spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')

            if arrow_start:
                head_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                head_elem.set('type', arrow_start)

            if arrow_end:
                tail_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                tail_elem.set('type', arrow_end)

        # Apply line color or gradient
        line_gradient = elem.get("lineGradient")
        color = elem.get("color", default_color)
        line_width = elem.get("lineWidth", default_line_width)

        if line_gradient:
            # Apply line gradient via shared helper
            try:
                from sdpm.builder.formatting import _build_grad_fill_element
                connector.line.width = Pt(line_width)  # creates _ln element
                ln = connector.line._ln
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                for child in list(ln):
                    if child.tag in (f'{ns}solidFill', f'{ns}noFill', f'{ns}gradFill'):
                        ln.remove(child)
                grad = _build_grad_fill_element(line_gradient)
                ln.insert(0, grad)
            except Exception:
                # Fallback to solid color
                stops = line_gradient.get("stops", [])
                if stops:
                    color_hex = stops[0].get("color", "#FFFFFF")
                    hex_color = color_hex.lstrip("#")
                    connector.line.fill.solid()
                    connector.line.color.rgb = RGBColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16)
                    )
                    connector.line.width = Pt(line_width)
        elif color and color != "none":
            connector.line.fill.solid()
            hex_color = color.lstrip("#")
            connector.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            connector.line.width = Pt(line_width)

        # Apply dash style
        dash_style = elem.get("dashStyle")
        if dash_style:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            dash_map = {
                "solid": MSO_LINE_DASH_STYLE.SOLID,
                "dash": MSO_LINE_DASH_STYLE.DASH,
                "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
                "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
                "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
            }
            if dash_style in dash_map:
                connector.line.dash_style = dash_map[dash_style]
            elif dash_style in ("sysDash", "sysDot", "sysDashDot", "sysDashDotDot"):
                from lxml import etree
                from pptx.oxml.ns import qn
                ln = connector.line._ln
                prstDash = etree.SubElement(ln, qn('a:prstDash'))
                prstDash.set('val', dash_style)
    

