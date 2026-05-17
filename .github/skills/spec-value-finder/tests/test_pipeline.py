"""spec-value-finder のパイプライン検証。

Excel/Word を実際に生成し、extract → find → fill が一気通貫で動くことを確認する。
Neo4j も外部サーバも不要なため、CI でそのまま実行できる。
"""
import json
from pathlib import Path

import openpyxl
import pytest

import docx

from comparer import format_comparison, read_manual_values, run_comparison
from extract import extract_file, find_files, to_markdown
from filler import fill
from finder import find_candidates
from mapping import validate_mapping, write_draft
from models import normalize
import yaml


# ---------------------------------------------------------------------------
# フィクスチャ生成
# ---------------------------------------------------------------------------

def _make_source_xlsx(path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Network"
    ws.append(["分類", "項目", "最大値", "単位"])
    ws.append(["MTU設定", "MTU上限", 1500, "bytes"])
    ws.append(["MTU設定", "ジャンボフレーム", 9000, "bytes"])
    wb.save(path)


def _make_source_docx(path):
    d = docx.Document()
    d.add_heading("ネットワーク仕様", level=1)
    d.add_paragraph("本機器の最大転送単位（ＭＴＵ）は 1500 バイトである。")
    d.add_heading("電源", level=1)
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "項目"
    t.cell(0, 1).text = "値"
    t.cell(1, 0).text = "定格電圧"
    t.cell(1, 1).text = "100V"
    d.save(path)


def _make_template_xlsx(path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["項目", "値", "出典", "確信度", "要確認"])
    ws.append(["MTU上限", "", "", "", ""])
    ws.append(["定格電圧", "", "", "", ""])
    wb.save(path)


# ---------------------------------------------------------------------------
# extract
# ---------------------------------------------------------------------------

def test_extract_excel_keeps_breadcrumb(tmp_path):
    src = tmp_path / "HW仕様書.xlsx"
    _make_source_xlsx(src)
    doc = extract_file(src)
    assert doc.fmt == "excel"
    cells = [c for b in doc.blocks for row in b.rows for c in row]
    mtu = next(c for c in cells if c.text == "1500")
    # breadcrumb にシート名と列見出しが含まれる
    assert "Network" in mtu.path
    assert "最大値" in mtu.path


def test_extract_word_picks_up_nontabular(tmp_path):
    src = tmp_path / "仕様書.docx"
    _make_source_docx(src)
    doc = extract_file(src)
    assert doc.fmt == "word"
    texts = [b for b in doc.blocks if getattr(b, "kind", "") == "text"]
    para = next(b for b in texts if "最大転送単位" in b.text)
    # 見出し階層が breadcrumb に入る
    assert "ネットワーク仕様" in para.path
    tables = [b for b in doc.blocks if getattr(b, "kind", "") == "table"]
    assert tables and any(c.text == "100V" for r in tables[0].rows for c in r)


def test_to_markdown_renders(tmp_path):
    src = tmp_path / "HW仕様書.xlsx"
    _make_source_xlsx(src)
    md = to_markdown(extract_file(src))
    assert "| MTU上限 |" in md.replace(" ", " ")
    assert "Network" in md


# ---------------------------------------------------------------------------
# find_files（部分一致走査）
# ---------------------------------------------------------------------------

def test_find_files_partial_match(tmp_path):
    _make_source_xlsx(tmp_path / "HW仕様書_v2.xlsx")
    _make_source_xlsx(tmp_path / "別資料.xlsx")
    hits = find_files(tmp_path, "HW仕様")
    assert [p.name for p in hits] == ["HW仕様書_v2.xlsx"]
    assert len(find_files(tmp_path, "")) == 2


# ---------------------------------------------------------------------------
# find（候補抽出）
# ---------------------------------------------------------------------------

def test_find_candidates_excel(tmp_path):
    src = tmp_path / "HW仕様書.xlsx"
    _make_source_xlsx(src)
    docs = [extract_file(src)]
    items = [{"target": "MTU上限", "keywords": ["MTU上限", "MTU"],
              "section_hint": "Network"}]
    result = find_candidates(docs, items)
    cands = result["items"][0]["candidates"]
    assert cands, "候補が見つかること"
    top = cands[0]
    # 行文脈に値 1500 が含まれる
    assert "1500" in top["row_values"]
    assert "MTU上限" in top["text"]


def test_find_candidates_word_nontabular(tmp_path):
    src = tmp_path / "仕様書.docx"
    _make_source_docx(src)
    docs = [extract_file(src)]
    # 全角ＭＴＵ の本文を半角キーワードで拾えること（NFKC 正規化）
    items = [{"target": "MTU", "keywords": ["MTU"]}]
    result = find_candidates(docs, items)
    cands = result["items"][0]["candidates"]
    assert any("最大転送単位" in c["text"] for c in cands)


# ---------------------------------------------------------------------------
# validate
# ---------------------------------------------------------------------------

def test_validate_detects_missing_keywords():
    data = {"version": 1, "source": {"folder": "./s"},
            "items": [{"target": "X", "keywords": []}]}
    errors, _ = validate_mapping(data)
    assert any("keywords" in e for e in errors)


def test_validate_passes_valid_mapping():
    data = {"version": 1, "source": {"folder": "./s"},
            "items": [{"target": "MTU上限", "keywords": ["MTU"],
                       "section_hint": "Network"}]}
    errors, _ = validate_mapping(data)
    assert errors == []


def test_write_draft_is_parseable(tmp_path):
    out = tmp_path / "mapping.yaml"
    write_draft(out, source_text="MTU は 1500")
    data = yaml.safe_load(out.read_text(encoding="utf-8"))
    assert data["version"] == 1
    assert data["items"] == []


# ---------------------------------------------------------------------------
# fill
# ---------------------------------------------------------------------------

def test_fill_excel_writes_and_appends(tmp_path):
    template = tmp_path / "記入例.xlsx"
    _make_template_xlsx(template)
    findings = [
        {"target": "MTU上限", "value": "1500", "source": "HW.xlsx :: C2",
         "confidence": "high", "needs_review": False},
        {"target": "電源電圧", "value": "100V", "source": "HW.xlsx :: B2",
         "confidence": "low", "needs_review": True},
    ]
    out = tmp_path / "結果.xlsx"
    report = fill(template, findings, out)
    assert report["matched"] == 1
    assert report["appended"] == 1

    wb = openpyxl.load_workbook(out)
    ws = wb.active
    rows = {r[0].value: r for r in ws.iter_rows()}
    assert rows["MTU上限"][1].value == "1500"
    assert rows["MTU上限"][3].value == "高"
    # 末尾追記された行も値と要確認フラグを持つ
    assert rows["電源電圧"][1].value == "100V"
    assert rows["電源電圧"][4].value == "要確認"


def test_fill_word_replaces_placeholder(tmp_path):
    template = tmp_path / "記入例.docx"
    d = docx.Document()
    d.add_paragraph("MTU上限: {{MTU上限}}")
    d.save(template)
    findings = [{"target": "MTU上限", "value": "1500"}]
    out = tmp_path / "結果.docx"
    report = fill(template, findings, out)
    assert report["replaced"] == 1
    assert "1500" in docx.Document(str(out)).paragraphs[0].text
    assert "{{" not in docx.Document(str(out)).paragraphs[0].text


# ---------------------------------------------------------------------------
# PDF / PowerPoint / Markdown / txt
# ---------------------------------------------------------------------------

def _make_pdf(path, text):
    """テキスト層を持つ最小構成の PDF を生成する（オフセットを計算して組み立てる）。"""
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
    ]
    stream = ("BT /F1 18 Tf 72 700 Td (" + text + ") Tj ET").encode("latin-1")
    objs.append(b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n"
                + stream + b"\nendstream")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, 1):
        offsets.append(len(pdf))
        pdf += str(i).encode() + b" 0 obj\n" + body + b"\nendobj\n"
    xref_pos = len(pdf)
    pdf += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n"
    pdf += b"0000000000 65535 f \n"
    for off in offsets:
        pdf += ("%010d 00000 n \n" % off).encode()
    pdf += b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
    pdf += b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    Path(path).write_bytes(bytes(pdf))


def test_extract_pdf_text_layer(tmp_path):
    src = tmp_path / "HW仕様書.pdf"
    _make_pdf(src, "MTU 1500 bytes")
    doc = extract_file(src)
    assert doc.fmt == "pdf"
    texts = [b for b in doc.blocks if getattr(b, "kind", "") == "text"]
    assert any("MTU" in b.text and "1500" in b.text for b in texts)
    assert all(b.locator == "ページ1" for b in texts)


def test_extract_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    src = tmp_path / "提案.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "ネットワーク仕様"
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    tb.text_frame.text = "MTU上限は 1500 bytes"
    gt = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(5), Inches(1))
    gt.table.cell(0, 0).text = "項目"
    gt.table.cell(0, 1).text = "値"
    gt.table.cell(1, 0).text = "定格電圧"
    gt.table.cell(1, 1).text = "100V"
    prs.save(src)

    doc = extract_file(src)
    assert doc.fmt == "powerpoint"
    texts = [b for b in doc.blocks if getattr(b, "kind", "") == "text"]
    tables = [b for b in doc.blocks if getattr(b, "kind", "") == "table"]
    assert any("MTU上限" in b.text for b in texts)
    assert all(b.locator == "スライド1" for b in texts)
    assert tables and any(c.text == "100V" for r in tables[0].rows for c in r)


def test_extract_markdown_headings_and_table(tmp_path):
    src = tmp_path / "仕様.md"
    src.write_text(
        "# ネットワーク仕様\n\n"
        "## MTU設定\n\n"
        "本機器の最大転送単位は次の通り。\n\n"
        "| 項目 | 値 |\n| --- | --- |\n| MTU上限 | 1500 |\n",
        encoding="utf-8")
    doc = extract_file(src)
    assert doc.fmt == "markdown"
    texts = [b for b in doc.blocks if getattr(b, "kind", "") == "text"]
    tables = [b for b in doc.blocks if getattr(b, "kind", "") == "table"]
    para = next(b for b in texts if "最大転送単位" in b.text)
    assert para.path == ["ネットワーク仕様", "MTU設定"]
    assert tables and any(c.text == "1500" for r in tables[0].rows for c in r)


def test_extract_txt(tmp_path):
    src = tmp_path / "メモ.txt"
    src.write_text("MTU上限: 1500 bytes\n\n定格電圧: 100V\n", encoding="utf-8")
    doc = extract_file(src)
    assert doc.fmt == "text"
    texts = [b for b in doc.blocks if getattr(b, "kind", "") == "text"]
    assert len(texts) == 2
    assert any("1500" in b.text for b in texts)


def test_find_across_formats(tmp_path):
    specs = tmp_path / "specs"
    specs.mkdir()
    (specs / "spec.md").write_text(
        "# ネットワーク\n\n最大転送単位（MTU）は 1500 bytes。\n", encoding="utf-8")
    _make_pdf(specs / "spec.pdf", "MTU max value is 1500")
    docs = [extract_file(p) for p in find_files(specs, "")]
    items = [{"target": "MTU上限", "keywords": ["MTU"]}]
    result = find_candidates(docs, items)
    cands = result["items"][0]["candidates"]
    assert len(cands) >= 2
    assert any("ページ1" in c["location"] for c in cands)


def test_unsupported_format_raises(tmp_path):
    legacy = tmp_path / "old.xls"
    legacy.write_text("dummy", encoding="utf-8")
    with pytest.raises(ValueError, match="旧形式"):
        extract_file(legacy)


# ---------------------------------------------------------------------------
# compare（手動転記 × スキル抽出の突き合わせ）
# ---------------------------------------------------------------------------

def _make_manual_xlsx(path, rows):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["項目", "値", "出典"])
    for item, val in rows:
        ws.append([item, val, ""])
    wb.save(path)


def read_manual_values_dict(rows, tmp_path):
    """rows から手動転記 xlsx を作って read_manual_values で読み戻すテスト補助。"""
    path = tmp_path / "手動転記.xlsx"
    _make_manual_xlsx(path, rows)
    return read_manual_values(path)


def test_read_manual_values_from_excel(tmp_path):
    src = tmp_path / "手動転記.xlsx"
    _make_manual_xlsx(src, [("MTU上限", "1500"), ("定格電圧", "100V")])
    manual = read_manual_values(src)
    assert manual[normalize("MTU上限")]["value"] == "1500"
    assert manual[normalize("定格電圧")]["value"] == "100V"


def test_read_manual_values_from_markdown(tmp_path):
    src = tmp_path / "手動転記.md"
    src.write_text("# 記入シート\n\n| 項目 | 値 |\n| --- | --- |\n| MTU上限 | 1500 |\n",
                   encoding="utf-8")
    manual = read_manual_values(src)
    assert manual[normalize("MTU上限")]["value"] == "1500"


def test_compare_verdicts(tmp_path):
    manual = read_manual_values_dict([
        ("MTU上限", "1500"),       # 一致
        ("ジャンボフレーム", "9000"),  # 不一致
        ("定格電圧", ""),            # 手動空
        ("予備項目", "ABC"),         # スキル側になし
    ], tmp_path)
    findings = [
        {"target": "MTU上限", "value": "1500", "confidence": "high"},
        {"target": "ジャンボフレーム", "value": "1500", "confidence": "high"},
        {"target": "定格電圧", "value": "100V", "confidence": "high"},
        {"target": "消費電力", "value": "50W", "confidence": "high"},  # 手動になし
    ]
    result = run_comparison(manual, findings)
    by = {it["target"]: it for it in result["items"]}
    assert by["MTU上限"]["verdict"] == "match"
    assert by["ジャンボフレーム"]["verdict"] == "mismatch"
    assert by["定格電圧"]["verdict"] == "missing_manual"
    assert by["消費電力"]["verdict"] == "missing_manual"
    assert by["予備項目"]["verdict"] == "missing_skill"
    assert result["summary"]["mismatch"] == 1
    assert "判定結果" in format_comparison(result)


def test_compare_numeric_match_flags_review(tmp_path):
    manual = read_manual_values_dict([("MTU上限", "1,500")], tmp_path)
    findings = [{"target": "MTU上限", "value": "1500 bytes", "confidence": "high"}]
    result = run_comparison(manual, findings)
    it = result["items"][0]
    assert it["verdict"] == "match"
    assert it["needs_review"] is True  # 数値一致だが表記差


def test_compare_low_confidence_match_flags_review(tmp_path):
    manual = read_manual_values_dict([("MTU上限", "1500")], tmp_path)
    findings = [{"target": "MTU上限", "value": "1500", "confidence": "low"}]
    result = run_comparison(manual, findings)
    assert result["items"][0]["verdict"] == "match"
    assert result["items"][0]["needs_review"] is True


def test_compare_uses_mapping_order(tmp_path):
    manual = read_manual_values_dict([("B項目", "2"), ("A項目", "1")], tmp_path)
    findings = [{"target": "A項目", "value": "1"}, {"target": "B項目", "value": "2"}]
    mapping_items = [{"target": "A項目"}, {"target": "B項目"}]
    result = run_comparison(manual, findings, mapping_items)
    assert [it["target"] for it in result["items"]] == ["A項目", "B項目"]


# ---------------------------------------------------------------------------
# 一気通貫
# ---------------------------------------------------------------------------

def test_end_to_end(tmp_path):
    specs = tmp_path / "specs"
    specs.mkdir()
    _make_source_xlsx(specs / "HW仕様書.xlsx")

    files = find_files(specs, "HW仕様")
    docs = [extract_file(p) for p in files]
    items = [{"target": "MTU上限", "keywords": ["MTU上限"], "section_hint": "Network"}]
    result = find_candidates(docs, items)
    assert result["items"][0]["candidates"]

    # Claude が候補を吟味して findings.json を作る想定 → ここでは値を直接確定
    findings = [{"target": "MTU上限", "value": "1500",
                 "source": result["items"][0]["candidates"][0]["location"],
                 "confidence": "high", "needs_review": False}]
    template = tmp_path / "記入例.xlsx"
    _make_template_xlsx(template)
    out = tmp_path / "結果.xlsx"
    fill(template, findings, out)
    wb = openpyxl.load_workbook(out)
    assert wb.active["B2"].value == "1500"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-v"]))
