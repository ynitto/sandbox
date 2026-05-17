"""元仕様書を構造化抽出する。

- Excel : openpyxl。結合セルは値を全構成セルへ展開し、表構造の breadcrumb を保つ。
- Word  : python-docx。見出し階層を breadcrumb 化し、表＋本文段落を拾う。
- PowerPoint: python-pptx。スライドタイトルを文脈に、図形テキスト＋表を拾う。
- PDF   : pypdfium2。透明テキスト層をページ単位で抽出する（表構造解析はしない）。
- Markdown/txt: 標準ライブラリのみ。見出し・段落・パイプ表を素直に拾う。

Table Transformer も Neo4j も使わない（純Python・サーバ/GPU不要）。
"""
from __future__ import annotations

import re
from pathlib import Path

from models import Cell, ExtractedDoc, Table, TextBlock

EXCEL_SUFFIXES = {".xlsx", ".xlsm"}
WORD_SUFFIXES = {".docx"}
PPTX_SUFFIXES = {".pptx"}
PDF_SUFFIXES = {".pdf"}
MD_SUFFIXES = {".md", ".markdown"}
TXT_SUFFIXES = {".txt"}
SUPPORTED = (EXCEL_SUFFIXES | WORD_SUFFIXES | PPTX_SUFFIXES
             | PDF_SUFFIXES | MD_SUFFIXES | TXT_SUFFIXES)


def find_files(folder: Path, name_match: str = "") -> list[Path]:
    """folder 配下を再帰的に走査し、ファイル名に name_match を部分一致で含む
    Excel/Word を返す（大文字小文字無視・name_match 空なら全件）。"""
    folder = Path(folder).expanduser()
    if not folder.is_dir():
        raise NotADirectoryError(f"フォルダが見つかりません: {folder}")
    needle = name_match.lower()
    hits: list[Path] = []
    for p in sorted(folder.rglob("*")):
        if not p.is_file() or p.suffix.lower() not in SUPPORTED:
            continue
        if p.name.startswith("~$"):  # Office の一時ロックファイル
            continue
        if needle and needle not in p.name.lower():
            continue
        hits.append(p)
    return hits


# ---------------------------------------------------------------------------
# breadcrumb 推定（結合セル carry-forward）
# ---------------------------------------------------------------------------

def _infer_paths(table: Table, context: list[str]) -> None:
    rows = table.rows
    if not rows:
        return
    header = next((r for r in rows if r and r[0].is_header), rows[0])
    col_headers = {c.col: c.text for c in header}
    data_rows = [r for r in rows if r and not r[0].is_header]

    for c in header:
        c.path = [*context, c.text] if c.text else list(context)

    if not data_rows:
        return

    n = len(data_rows)
    col_cells: dict[int, list[Cell]] = {}
    for r in data_rows:
        for c in r:
            col_cells.setdefault(c.col, []).append(c)

    # 30%超が空の列は階層列とみなす（結合セル未展開のレイアウト対策）
    hierarchy = {ci for ci, cells in col_cells.items()
                 if n > 1 and sum(1 for c in cells if not c.text.strip()) / n > 0.3}
    carry: dict[int, str] = {}
    sorted_h = sorted(hierarchy)

    for row in data_rows:
        for cell in row:
            col = cell.col
            val = cell.text.strip()
            if col in hierarchy:
                if val:
                    carry[col] = val
                elif col in carry:
                    cell.text = carry[col]
            path = list(context)
            for hc in sorted_h:
                if hc != col and carry.get(hc):
                    path.append(carry[hc])
            col_name = col_headers.get(col, "")
            if col_name:
                path.append(col_name)
            cell.path = [p for p in path if p]


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def extract_excel(path: Path) -> ExtractedDoc:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    blocks: list = []
    for name in wb.sheetnames:
        ws = wb[name]
        merged: dict[tuple[int, int], object] = {}
        for rng in ws.merged_cells.ranges:
            top_left = ws.cell(row=rng.min_row, column=rng.min_col).value
            for r in range(rng.min_row, rng.max_row + 1):
                for c in range(rng.min_col, rng.max_col + 1):
                    merged[(r, c)] = top_left

        grid: list[list[Cell]] = []
        for row in ws.iter_rows():
            cells: list[Cell] = []
            for cell in row:
                v = merged.get((cell.row, cell.column), cell.value)
                text = "" if v is None else str(v).strip()
                cells.append(Cell(text=text, row=cell.row, col=cell.column))
            if any(c.text for c in cells):
                grid.append(cells)

        if not grid:
            continue
        for c in grid[0]:
            c.is_header = True
        table = Table(container=name, rows=grid)
        _infer_paths(table, [name])
        blocks.append(table)

    wb.close()
    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="excel", blocks=blocks)


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"(?:heading|見出し)\s*(\d+)", re.IGNORECASE)


def _heading_level(style_name: str) -> int:
    """見出しスタイルならレベル(1始まり)、本文なら 0 を返す。"""
    m = _HEADING_RE.search(style_name or "")
    if m:
        return int(m.group(1))
    if (style_name or "").strip().lower() in ("title", "表題"):
        return 1
    return 0


def extract_word(path: Path) -> ExtractedDoc:
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph

    doc = docx.Document(str(path))
    blocks: list = []
    heading_stack: list[tuple[int, str]] = []  # (level, text)
    table_no = 0

    def current_path() -> list[str]:
        return [t for _, t in heading_stack]

    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = DocxParagraph(child, doc)
            text = (para.text or "").strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            level = _heading_level(style)
            if level:
                heading_stack[:] = [h for h in heading_stack if h[0] < level]
                blocks.append(TextBlock(text=text, style=style, path=current_path()))
                heading_stack.append((level, text))
            else:
                blocks.append(TextBlock(text=text, style=style or "Normal",
                                        path=current_path()))
        elif child.tag == qn("w:tbl"):
            table_no += 1
            dtbl = DocxTable(child, doc)
            grid: list[list[Cell]] = []
            for r_idx, drow in enumerate(dtbl.rows):
                cells = [Cell(text=(dc.text or "").strip(), row=r_idx, col=c_idx,
                              is_header=(r_idx == 0))
                         for c_idx, dc in enumerate(drow.cells)]
                if any(c.text for c in cells):
                    grid.append(cells)
            if not grid:
                continue
            container = f"表#{table_no}"
            table = Table(container=container, rows=grid)
            _infer_paths(table, current_path() or [container])
            blocks.append(table)

    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="word", blocks=blocks)


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def extract_pptx(path: Path) -> ExtractedDoc:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list = []
    table_no = 0
    for n, slide in enumerate(prs.slides, 1):
        loc = f"スライド{n}"
        title_shape = slide.shapes.title
        title = ""
        if title_shape is not None and title_shape.has_text_frame:
            title = (title_shape.text or "").strip()
        context = [title] if title else [loc]
        if title:
            blocks.append(TextBlock(text=title, style="Heading 1",
                                    path=[], locator=loc))
        for shape in slide.shapes:
            if shape.has_table:
                table_no += 1
                grid: list[list[Cell]] = []
                for r_idx, row in enumerate(shape.table.rows):
                    cells = [Cell(text=(c.text or "").strip(), row=r_idx, col=c_idx,
                                  is_header=(r_idx == 0))
                             for c_idx, c in enumerate(row.cells)]
                    if any(c.text for c in cells):
                        grid.append(cells)
                if grid:
                    table = Table(container=f"{loc} 表#{table_no}", rows=grid)
                    _infer_paths(table, context)
                    blocks.append(table)
            elif shape.has_text_frame and shape is not title_shape:
                text = (shape.text or "").strip()
                if text:
                    blocks.append(TextBlock(text=text, style="Normal",
                                            path=[title] if title else [],
                                            locator=loc))
    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="powerpoint", blocks=blocks)


# ---------------------------------------------------------------------------
# PDF（透明テキスト層）
# ---------------------------------------------------------------------------

def _split_paragraphs(text: str) -> list[str]:
    return [c.strip() for c in re.split(r"\n\s*\n", text) if c.strip()]


def extract_pdf(path: Path) -> ExtractedDoc:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(str(path))
    blocks: list = []
    try:
        for n in range(len(pdf)):
            page = pdf[n]
            textpage = page.get_textpage()
            text = textpage.get_text_range() or ""
            textpage.close()
            page.close()
            loc = f"ページ{n + 1}"
            for chunk in _split_paragraphs(text):
                blocks.append(TextBlock(text=chunk, style="Normal",
                                        path=[loc], locator=loc))
    finally:
        pdf.close()
    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="pdf", blocks=blocks)


# ---------------------------------------------------------------------------
# Markdown / txt
# ---------------------------------------------------------------------------

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")


def _is_md_table_row(line: str) -> bool:
    s = line.strip()
    return s.startswith("|") and s.count("|") >= 2


def _is_md_separator(line: str) -> bool:
    s = line.strip().strip("|")
    return bool(s) and "-" in s and all(ch in " :-|" for ch in s)


def _parse_md_row(line: str) -> list[str]:
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _md_table(lines: list[str], context: list[str]) -> Table:
    grid: list[list[Cell]] = []
    for r_idx, line in enumerate(lines):
        values = _parse_md_row(line)
        cells = [Cell(text=v, row=r_idx, col=c_idx, is_header=(r_idx == 0))
                 for c_idx, v in enumerate(values)]
        grid.append(cells)
    table = Table(container="表", rows=grid)
    _infer_paths(table, context or ["表"])
    return table


def extract_markdown(path: Path) -> ExtractedDoc:
    lines = Path(path).read_text(encoding="utf-8", errors="replace").splitlines()
    blocks: list = []
    heading_stack: list[tuple[int, str]] = []
    para_buf: list[str] = []

    def heading_path() -> list[str]:
        return [t for _, t in heading_stack]

    def flush() -> None:
        nonlocal para_buf
        joined = " ".join(s.strip() for s in para_buf).strip()
        if joined:
            blocks.append(TextBlock(text=joined, style="Normal", path=heading_path()))
        para_buf = []

    i = 0
    while i < len(lines):
        line = lines[i]
        m = _MD_HEADING_RE.match(line)
        if m:
            flush()
            level, title = len(m.group(1)), m.group(2).strip()
            heading_stack[:] = [h for h in heading_stack if h[0] < level]
            blocks.append(TextBlock(text=title, style=f"Heading {level}",
                                    path=heading_path()))
            heading_stack.append((level, title))
        elif (_is_md_table_row(line) and i + 1 < len(lines)
              and _is_md_separator(lines[i + 1])):
            flush()
            tbl = [line]
            j = i + 2
            while j < len(lines) and _is_md_table_row(lines[j]):
                tbl.append(lines[j])
                j += 1
            blocks.append(_md_table(tbl, heading_path()))
            i = j
            continue
        elif not line.strip():
            flush()
        else:
            para_buf.append(line)
        i += 1
    flush()
    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="markdown", blocks=blocks)


def extract_txt(path: Path) -> ExtractedDoc:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    blocks = [TextBlock(text=chunk, style="Normal")
              for chunk in _split_paragraphs(text)]
    return ExtractedDoc(source=str(Path(path).resolve()), filename=Path(path).name,
                        fmt="text", blocks=blocks)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_LEGACY_HINT = {".xls": ".xlsx", ".doc": ".docx", ".ppt": ".pptx"}


def extract_file(path: Path | str) -> ExtractedDoc:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix in EXCEL_SUFFIXES:
        return extract_excel(path)
    if suffix in WORD_SUFFIXES:
        return extract_word(path)
    if suffix in PPTX_SUFFIXES:
        return extract_pptx(path)
    if suffix in PDF_SUFFIXES:
        return extract_pdf(path)
    if suffix in MD_SUFFIXES:
        return extract_markdown(path)
    if suffix in TXT_SUFFIXES:
        return extract_txt(path)
    if suffix in _LEGACY_HINT:
        raise ValueError(
            f"旧形式は非対応です（{suffix}）。{_LEGACY_HINT[suffix]} に変換してください: {path}")
    raise ValueError(f"対応していないファイル形式です: {suffix}")


# ---------------------------------------------------------------------------
# Markdown 化（Claude が全文を読むための中間表現）
# ---------------------------------------------------------------------------

def _table_to_markdown(table: Table) -> str:
    if not table.rows:
        return ""
    width = max(len(r) for r in table.rows)
    lines: list[str] = []
    sep_done = False
    for row in table.rows:
        texts = [c.text for c in row] + [""] * (width - len(row))
        escaped = [t.replace("|", "\\|").replace("\n", " ") for t in texts]
        lines.append("| " + " | ".join(escaped) + " |")
        if not sep_done:
            lines.append("|" + "|".join([" --- "] * width) + "|")
            sep_done = True
    return "\n".join(lines)


def to_markdown(doc: ExtractedDoc) -> str:
    parts = [f"# {doc.filename}", ""]
    last_loc = ""
    for block in doc.blocks:
        loc = getattr(block, "locator", "") or ""
        if loc and loc != last_loc:
            parts.append(f"## {loc}")
            parts.append("")
        last_loc = loc or last_loc
        if isinstance(block, Table):
            parts.append(f"### {block.container}")
            parts.append(_table_to_markdown(block))
            parts.append("")
        else:
            level = _heading_level(block.style)
            if level:
                parts.append(f"{'#' * min(level + 1, 6)} {block.text}")
            else:
                parts.append(block.text)
            parts.append("")
    return "\n".join(parts).rstrip() + "\n"
