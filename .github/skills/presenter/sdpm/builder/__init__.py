# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""PPTX Builder - Core PPTX generation from JSON.

PPTXBuilder is composed from mixin classes:
- FormattingMixin: fill, line, opacity, unit conversion, text styling
- PlaceholderMixin: title, content, section placeholder fills
- TextboxMixin, ShapeMixin, ImageMixin, LineMixin, ChartMixin, TableMixin, GroupMixin
"""
from pathlib import Path

import json

from sdpm.utils.io import read_json

from pptx import Presentation
from pptx.util import Emu, Pt  # noqa: F401
from pptx.dml.color import RGBColor  # noqa: F401
from pptx.enum.shapes import MSO_SHAPE  # noqa: F401
from pptx.enum.shapes import MSO_CONNECTOR  # noqa: F401

from sdpm.utils.text import (  # noqa: F401
    is_fullwidth, normalize_spacing, parse_styled_text,
    _parse_non_link_styles, _expand_styled_newlines,
)
from sdpm.utils.svg import (  # noqa: F401
    _recolor_svg, get_svg_dimensions, generate_qr_svg, add_svg_to_slide,
)
from sdpm.utils.image import resolve_image_path, apply_image_effects  # noqa: F401
from sdpm.utils.effects import apply_effects  # noqa: F401
from sdpm.assets import (  # noqa: F401
    ICON_DIR, ICON_LOCAL_DIR, resolve_icon_path, check_icon_exists,
    resolve_asset_path, check_asset_exists,
    _assets_not_installed_error,
)

from sdpm.builder.constants import ARCH_GROUP_DEFS  # noqa: F401
from sdpm.builder.formatting import FormattingMixin
from sdpm.builder.placeholder import PlaceholderMixin
from sdpm.builder.elements.textbox import TextboxMixin
from sdpm.builder.elements.shape import ShapeMixin
from sdpm.builder.elements.image import ImageMixin
from sdpm.builder.elements.line import LineMixin
from sdpm.builder.elements.chart import ChartMixin
from sdpm.builder.elements.table import TableMixin
from sdpm.builder.elements.group import GroupMixin

# Re-export effect presets for backward compatibility
from sdpm.utils.effects import (  # noqa: E402, F401
    SHADOW_PRESETS, GLOW_PRESETS, REFLECTION_PRESETS,
    BEVEL_PRESETS, ROTATION3D_PRESETS,
)


class PPTXBuilder(
    FormattingMixin,
    PlaceholderMixin,
    TextboxMixin,
    ShapeMixin,
    ImageMixin,
    LineMixin,
    ChartMixin,
    TableMixin,
    GroupMixin,
):
    """Generate PowerPoint presentations from JSON definitions."""

    # Default slide dimensions (1920x1080 basis)
    SLIDE_WIDTH = 12192000
    SLIDE_HEIGHT = 6858000
    EMU_PER_PX = 6350

    def __init__(self, template_path: Path, custom_template: bool = False,
                 fonts: dict = None, base_dir: Path = None, keep_empty_placeholders: bool = False,
                 default_text_color: str = None):
        """Initialize PPTXBuilder.

        Args:
            template_path: Path to the PowerPoint template file.
            custom_template: Whether the template is user-provided (vs built-in default).
            fonts: Font configuration dict with "fullwidth" and "halfwidth" keys.
                   Required. Set via presentation.json "fonts" field.
            base_dir: Base directory for resolving relative paths in slide JSON.
            default_text_color: Default color for text and icons (e.g. "#FFFFFF").
                   Required. Set via presentation.json "defaultTextColor" field.

        Raises:
            ValueError: If fonts or default_text_color is not provided.
        """
        if fonts is None:
            raise ValueError(
                "'fonts' is required. Set 'fonts' in presentation.json "
                '(e.g. {"fullwidth": "メイリオ", "halfwidth": "Calibri"}). '
                "Run analyze-template to detect fonts from your template."
            )
        if default_text_color is None:
            raise ValueError(
                "'defaultTextColor' is required. Set 'defaultTextColor' in presentation.json "
                '(e.g. "defaultTextColor": "#FFFFFF"). '
                "This color is used as the default for text and icon colors."
            )
        self.prs = Presentation(str(template_path))
        self.theme_colors, self.is_dark = self._extract_theme_colors(template_path)
        self.theme_colors["text"] = default_text_color
        self.master_idx = 0
        self.EMU_PER_PX = int(self.prs.slide_width) / 1920
        self.custom_template = custom_template
        self.fonts = fonts
        self.keep_empty_placeholders = keep_empty_placeholders
        self.layouts = self._build_layout_map()
        self._base_dir = base_dir if base_dir is not None else Path(".")
        self._list_styles = self._load_list_styles()
        self._clear_slides()

    @staticmethod
    def _extract_theme_colors(template_path):
        """Extract theme colors from template's clrScheme + clrMap."""
        import zipfile
        from lxml import etree
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

        # Find theme XML linked to slide master 1
        scheme = {}
        with zipfile.ZipFile(str(template_path)) as z:
            # Determine which theme file slideMaster1 references
            theme_target = 'ppt/theme/theme1.xml'  # fallback
            for name in z.namelist():
                if 'slideMaster1.xml.rels' in name:
                    rels = etree.fromstring(z.read(name))
                    for rel in rels:
                        if 'theme' in rel.get('Target', ''):
                            target = rel.get('Target').replace('..', 'ppt')
                            if not target.startswith('ppt/'):
                                target = 'ppt/theme/' + target.split('/')[-1]
                            theme_target = target
                            break
                    break

            if theme_target in z.namelist():
                tree = etree.fromstring(z.read(theme_target))
            else:
                # Fallback: first theme file
                tree = None
                for name in sorted(z.namelist()):
                    if 'theme' in name and name.endswith('.xml'):
                        tree = etree.fromstring(z.read(name))
                        break

            if tree is not None:
                cs = tree.find(f'.//{{{ns_a}}}clrScheme')
                if cs is not None:
                    for child in cs:
                        tag = child.tag.split('}')[1]
                        val_el = child[0] if len(child) > 0 else None
                        if val_el is not None:
                            hex_val = val_el.get('lastClr') or val_el.get('val') or '000000'
                            try:
                                int(hex_val, 16)
                            except ValueError:
                                hex_val = val_el.get('lastClr') or '000000'
                            scheme[tag] = hex_val

        # Parse clrMap from slide master
        prs = Presentation(str(template_path))
        master = prs.slide_masters[0]
        clr_map = master.element.find(f'.//{{{ns_p}}}clrMap')
        bg1_ref = clr_map.get('bg1', 'lt1') if clr_map is not None else 'lt1'
        tx1_ref = clr_map.get('tx1', 'dk1') if clr_map is not None else 'dk1'
        bg2_ref = clr_map.get('bg2', 'lt2') if clr_map is not None else 'lt2'
        tx2_ref = clr_map.get('tx2', 'dk2') if clr_map is not None else 'dk2'

        colors = {
            "text": f"#{scheme.get(tx1_ref, '000000')}",
            "background": f"#{scheme.get(bg1_ref, 'FFFFFF')}",
            "text2": f"#{scheme.get(tx2_ref, '000000')}",
            "background2": f"#{scheme.get(bg2_ref, 'FFFFFF')}",
        }
        for i in range(1, 7):
            colors[f"accent{i}"] = f"#{scheme.get(f'accent{i}', '4A90D9')}"

        # is_dark: text color luminance > 128 means dark background
        tx = colors["text"].lstrip("#")
        r, g, b = int(tx[:2], 16), int(tx[2:4], 16), int(tx[4:6], 16)
        is_dark = (0.299 * r + 0.587 * g + 0.114 * b) > 128

        return colors, is_dark

    def _build_layout_map(self):
        """Build layout name → index mapping from template's slide layouts."""
        layouts = {}
        for i, layout in enumerate(self.prs.slide_masters[self.master_idx].slide_layouts):
            name = layout.name if layout.name else f"layout-{i+1:02d}"
            layouts[name] = i
        return layouts

    def _load_list_styles(self):
        """Load bullet/numbering definitions from slide master's bodyStyle."""
        from pptx.oxml.ns import qn
        master = self.prs.slide_masters[self.master_idx]
        tx_styles = master._element.find(qn('p:txStyles'))
        if tx_styles is None:
            return {}
        body_style = tx_styles.find(qn('p:bodyStyle'))
        if body_style is None:
            return {}
        styles = {}
        for i, lvl in enumerate(body_style):
            bu_char = lvl.find(qn('a:buChar'))
            bu_font = lvl.find(qn('a:buFont'))
            bu_auto_num = lvl.find(qn('a:buAutoNum'))
            entry = {
                'marL': lvl.get('marL'),
                'indent': lvl.get('indent'),
            }
            if bu_char is not None:
                entry['buChar'] = bu_char.get('char')
            if bu_font is not None:
                entry['buFont'] = bu_font.get('typeface')
            if bu_auto_num is not None:
                entry['buAutoNum'] = bu_auto_num.get('type')
            styles[i] = entry
        return styles

    def _clear_slides(self):
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def _get_layout(self, layout_name: str):
        if layout_name not in self.layouts:
            raise ValueError(f"Unknown layout: {layout_name}. Available: {list(self.layouts.keys())}")
        layout_idx = self.layouts[layout_name]
        return self.prs.slide_masters[self.master_idx].slide_layouts[layout_idx]

    def add_slide(self, slide_def: dict):
        layout_name = slide_def.get("layout")
        if not layout_name:
            raise ValueError(f"'layout' is required. Available layouts: {list(self.layouts.keys())}")
        master_idx = slide_def.get("masterIndex", self.master_idx)

        # Search layout in specified master, fallback to default master
        layout = None
        for mi in ([master_idx, self.master_idx] if master_idx != self.master_idx else [self.master_idx]):
            for i, sl in enumerate(self.prs.slide_masters[mi].slide_layouts):
                if (sl.name or f"layout-{i+1:02d}") == layout_name:
                    layout = sl
                    break
            if layout:
                break
        if layout is None:
            raise ValueError(f"Unknown layout: '{layout_name}'. Available layouts: {list(self.layouts.keys())}")
        slide = self.prs.slides.add_slide(layout)

        self._fill_placeholders(slide, slide_def)

        # Per-slide theme override (save/restore around element processing)
        saved_text = self.theme_colors["text"]
        saved_bg = self.theme_colors["background"]
        saved_is_dark = self.is_dark

        bg_color = slide_def.get("background")
        if bg_color:
            from pptx.dml.color import RGBColor
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip('#'))
            self.theme_colors["background"] = bg_color
            r, g, b = int(bg_color.lstrip('#')[:2], 16), int(bg_color.lstrip('#')[2:4], 16), int(bg_color.lstrip('#')[4:6], 16)
            self.is_dark = (0.299 * r + 0.587 * g + 0.114 * b) < 128

        slide_text_color = slide_def.get("defaultTextColor")
        if slide_text_color:
            self.theme_colors["text"] = slide_text_color

        # Process elements
        elements = slide_def.get("elements", [])
        expanded = []
        for elem in elements:
            if elem.get("type") == "include":
                src = elem.get("src", "")
                include_path = Path(src) if Path(src).is_absolute() else self._base_dir / src
                if include_path.exists():
                    try:
                        inc_data = read_json(include_path)
                    except (json.JSONDecodeError, ValueError) as e:
                        raise ValueError(f"Invalid JSON in include file {src}: {e}")
                    inc_elements = inc_data if isinstance(inc_data, list) else inc_data.get("elements", [])
                    expanded.extend(inc_elements)
            else:
                expanded.append(elem)
        for elem in expanded:
            elem_type = elem.get("type")
            if elem_type == "group":
                self._add_group(slide, elem)
            elif elem_type == "arch-group":
                self._add_arch_group(slide, elem)
            elif elem_type == "table":
                self._add_table(slide, elem)
            elif elem_type == "textbox":
                ph_idx = elem.get("_phIdx")
                if ph_idx is not None:
                    for ph in slide.placeholders:
                        if ph.placeholder_format.idx == ph_idx:
                            paras = elem.get("paragraphs")
                            text = elem.get("text", "")
                            if paras:
                                # Write each paragraph as separate <a:p> to preserve inherited bullets
                                tf = ph.text_frame
                                for pi, para_item in enumerate(paras):
                                    para_text = para_item.get("text", "") if isinstance(para_item, dict) else str(para_item)
                                    if pi == 0:
                                        p = tf.paragraphs[0]
                                    else:
                                        p = tf.add_paragraph()
                                    self._apply_styled_text(p, para_text, no_default_color=True, no_default_font=True)
                                    # Restore bullet and indent from JSON
                                    if isinstance(para_item, dict):
                                        if para_item.get("bullet") is True:
                                            from pptx.oxml.ns import qn as _qn
                                            from lxml import etree as _et
                                            pPr = p._element.get_or_add_pPr()
                                            if para_item.get("marL") is not None:
                                                pPr.set('marL', str(para_item["marL"]))
                                            if para_item.get("indent") is not None:
                                                pPr.set('indent', str(para_item["indent"]))
                                            # Remove buNone if present, add buChar
                                            for bn in pPr.findall(_qn('a:buNone')):
                                                pPr.remove(bn)
                                            if para_item.get("buFont"):
                                                bf = _et.SubElement(pPr, _qn('a:buFont'))
                                                bf.set('typeface', para_item["buFont"])
                                            bc = _et.SubElement(pPr, _qn('a:buChar'))
                                            bc.set('char', '•')
                            elif text.strip():
                                val = text
                                if elem.get("align"):
                                    val = {"text": text, "align": elem["align"]}
                                self._fill_placeholder(ph, val)
                            # Apply position from elem x/y/width/height
                            pos_val = {}
                            if elem.get("x") is not None:
                                pos_val["_x"] = elem["x"]
                            if elem.get("y") is not None:
                                pos_val["_y"] = elem["y"]
                            if elem.get("width") is not None:
                                pos_val["_width"] = elem["width"]
                            if elem.get("height") is not None:
                                pos_val["_height"] = elem["height"]
                            if pos_val:
                                self._apply_placeholder_position(ph, pos_val)
                            break
                else:
                    self._add_textbox(slide, elem)
            elif elem_type == "image":
                self._add_image(slide, elem)
            elif elem_type == "shape":
                self._add_shape(slide, elem)
            elif elem_type == "freeform":
                self._add_freeform_shape(slide, elem)
            elif elem_type == "line":
                self._add_line(slide, elem)
            elif elem_type == "chart":
                self._add_chart(slide, elem)
            elif elem_type == "video":
                self._add_video(slide, elem)

        if "notes" in slide_def:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.clear()
            p = notes_frame.paragraphs[0]
            self._apply_styled_text(p, slide_def["notes"])

        # Remove empty placeholders
        if not self.keep_empty_placeholders:
            for ph in list(slide.placeholders):
                if ph.has_text_frame and not ph.text_frame.text.strip():
                    slide.shapes._spTree.remove(ph._element)
                elif not ph.has_text_frame and ph.placeholder_format.type == 18:  # PICTURE
                    slide.shapes._spTree.remove(ph._element)

        # Restore per-slide theme override
        self.theme_colors["text"] = saved_text
        self.theme_colors["background"] = saved_bg
        self.is_dark = saved_is_dark

        if slide_def.get("hidden"):
            slide._element.set("show", "0")

        return slide

    def save(self, output_path: Path):
        self.prs.save(str(output_path))


def resolve_override(slide_def, id_map, visited=None):
    """Resolve override chain and merge elements."""
    if visited is None:
        visited = set()

    override_id = slide_def.get("override")
    if not override_id:
        return slide_def

    if override_id in visited:
        chain = " -> ".join(visited) + f" -> {override_id}"
        raise ValueError(f"Circular override detected: {chain}")

    if override_id not in id_map:
        raise ValueError(f"Override target '{override_id}' not found")

    visited.add(override_id)
    base = resolve_override(id_map[override_id], id_map, visited)

    base_elements = list(base.get("elements", []))
    override_elements = slide_def.get("elements", [])

    result = dict(slide_def)
    result["elements"] = base_elements + override_elements
    result.pop("override", None)

    return result


def validate_icons_in_json(data: dict) -> list:
    """Validate all icons in JSON and return list of missing icons."""
    missing = []
    slides = data.get("slides", [])

    def check_elements(elements):
        if not elements:
            return
        for elem in elements:
            if elem.get("type") == "image":
                src = elem.get("src") or elem.get("path", "")
                if src.startswith("icons:") or src.startswith("assets:"):
                    if not check_asset_exists(src):
                        missing.append(src)
            elif elem.get("type") == "group":
                check_elements(elem.get("elements", []))

    for slide in slides:
        check_elements(slide.get("elements", []))

    return list(set(missing))
