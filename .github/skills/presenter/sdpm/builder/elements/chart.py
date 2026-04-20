# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Chart element."""
from pptx.dml.color import RGBColor
from pptx.util import Pt


class ChartMixin:
    """Mixin providing chart element methods."""

    def _add_chart(self, slide, elem):
        """Add a native PPTX chart element."""
        from lxml import etree
        from pptx.oxml.ns import qn
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE, XL_TICK_MARK
        from pptx.chart.data import CategoryChartData
    
        chart_type = elem.get("chartType", "bar")
        categories = elem.get("categories", [])
        series_defs = elem.get("series", [])
        stacked = elem.get("stacked", False)
        horizontal = elem.get("horizontal", False)
        show_legend = elem.get("legend", True)
        show_data_labels = elem.get("dataLabels", False)
        number_format = elem.get("numberFormat")
        markers = elem.get("markers", True)
        smooth = elem.get("smooth", False)
        title = elem.get("title")
        style = elem.get("style", {})
    
        # Style overrides (all optional, sensible defaults applied)
        style_gridline_color = style.get("gridlineColor")
        style_gridline_width = style.get("gridlineWidth")
        style_gridline_dash = style.get("gridlineDash")  # solid|dash|dot
        style_axis_color = style.get("axisColor")
        style_axis_width = style.get("axisWidth")
        style_font_color = style.get("fontColor")
        style_font_size = style.get("fontSize")
        style_legend_position = style.get("legendPosition")  # bottom|right|left|top
        style_gap_width = style.get("gapWidth")
        style_line_width = style.get("lineWidth")
        style_marker_size = style.get("markerSize")
    
        # Determine XL_CHART_TYPE
        type_map = {
            ("bar", False, False): XL_CHART_TYPE.COLUMN_CLUSTERED,
            ("bar", True, False): XL_CHART_TYPE.COLUMN_STACKED,
            ("bar", False, True): XL_CHART_TYPE.BAR_CLUSTERED,
            ("bar", True, True): XL_CHART_TYPE.BAR_STACKED,
            ("line", False, False): XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE,
            ("line", True, False): XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE,
            ("pie", False, False): XL_CHART_TYPE.PIE,
            ("donut", False, False): XL_CHART_TYPE.DOUGHNUT,
        }
        xl_type = type_map.get((chart_type, stacked, horizontal))
        if xl_type is None:
            xl_type = type_map.get((chart_type, False, False), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
        # Build chart data
        chart_data = CategoryChartData()
        # Generate dummy categories if missing (e.g. donut charts)
        if not categories and series_defs:
            max_vals = max(len(s.get("values", [])) for s in series_defs)
            categories = [f"Cat{i+1}" for i in range(max_vals)]
        chart_data.categories = categories
        for s in series_defs:
            chart_data.add_series(s.get("name", ""), s.get("values", []))
    
        # Add chart to slide
        x = self._px_to_emu(elem.get("x", 192))
        y = self._px_to_emu(elem.get("y", 216))
        w = self._px_to_emu(elem.get("width", 1536))
        h = self._px_to_emu(elem.get("height", 700))
        graphic_frame = slide.shapes.add_chart(xl_type, x, y, w, h, chart_data)
        chart = graphic_frame.chart

        # Restore original chart styling for perfect roundtrip
        if elem.get("_chartXml"):
            try:
                from lxml import etree
                from pptx.oxml.ns import qn
                orig = etree.fromstring(elem["_chartXml"])
                ns_c = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

                # Get python-pptx's Excel relationship ID
                xlsx_rid = None
                for rel_key, rel in chart.part.rels.items():
                    if 'oleObject' in rel.reltype or 'package' in rel.reltype:
                        xlsx_rid = rel_key
                        break

                # Update externalData r:id to match python-pptx's rels
                ext_data = orig.find('c:externalData', ns_c)
                if ext_data is not None and xlsx_rid:
                    ext_data.set(f'{{{ns_c["r"]}}}id', xlsx_rid)
                elif ext_data is not None:
                    orig.remove(ext_data)

                # Remove userShapes ref (will re-add with correct rId)
                user_shapes_el = orig.find('c:userShapes', ns_c)
                if user_shapes_el is not None:
                    orig.remove(user_shapes_el)

                # Replace entire chart XML
                chart.part._element = orig

                # Add chartUserShapes drawing
                if elem.get("_chartUserShapes"):
                    from pptx.opc.package import Part
                    from pptx.opc.packuri import PackURI
                    drawing_xml = elem["_chartUserShapes"].encode('utf-8')
                    chart_num = chart.part.partname.split('chart')[-1].replace('.xml', '')
                    drawing_pn = PackURI(f'/ppt/drawings/drawing{chart_num}.xml')
                    drawing_part = Part(
                        drawing_pn,
                        'application/vnd.openxmlformats-officedocument.drawingml.chartshapes+xml',
                        chart.part.package,
                        drawing_xml,
                    )
                    us_rid = chart.part.relate_to(drawing_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/chartUserShapes')
                    # Add userShapes element with correct rId
                    us_el = etree.SubElement(orig, qn('c:userShapes'))
                    us_el.set(qn('r:id'), us_rid)

                return
            except Exception as _chart_err:
                import sys
                print(f"Warning: chart XML restore failed: {_chart_err}", file=sys.stderr)
                pass  # Fall through to normal chart building
    
        # Theme-aware design tokens (from template theme colors)
        tc = self.theme_colors
        text_hex = tc["text"].lstrip("#")
        bg_hex = tc["background"].lstrip("#")
        text_color_str = style_font_color.lstrip("#") if style_font_color else text_hex
        text_color = RGBColor.from_string(text_color_str)
        # Muted text: blend text toward background
        tr, tg, tb = int(text_hex[:2],16), int(text_hex[2:4],16), int(text_hex[4:6],16)
        br, bgg, bb = int(bg_hex[:2],16), int(bg_hex[2:4],16), int(bg_hex[4:6],16)
        mr, mg, mb = (tr+br)//2, (tg+bgg)//2, (tb+bb)//2
        muted_str = f"{mr:02X}{mg:02X}{mb:02X}"
        text_muted = RGBColor.from_string(muted_str)
        # Gridline/axis: subtle color near background
        default_gridline = f"#{(br+tr*1)//2:02X}{(bgg+tg*1)//2:02X}{(bb+tb*1)//2:02X}"
        gridline_color = (style_gridline_color or default_gridline).lstrip("#")
        axis_line_color = (style_axis_color or default_gridline).lstrip("#")
        gridline_width = str(int((style_gridline_width or 0.25) * 12700))  # pt to EMU
        gridline_dash = style_gridline_dash or "dash"
        axis_width = str(int((style_axis_width or 0.5) * 12700))
        tick_font_size = style_font_size or 10
        legend_font_size = style_font_size or 10
        # Default series colors (neutral palette, overridable per-series via JSON)
        default_series_colors = [tc[f"accent{i}"].lstrip("#") for i in range(1, 7)]
    
        # Title
        chart.has_title = bool(title)
        if title:
            tf = chart.chart_title.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(elem.get("titleFontSize", 14))
            title_fc = elem.get("titleFontColor")
            if title_fc:
                tf.paragraphs[0].font.color.rgb = RGBColor.from_string(title_fc.lstrip("#"))
            elif not elem.get("valueAxis"):
                tf.paragraphs[0].font.color.rgb = text_color
            tf.paragraphs[0].font.name = self.fonts["halfwidth"]
    
        # Transparent background (chartSpace + plotArea)
        for target in [chart._chartSpace]:
            sp_pr = target.find(qn('c:spPr'))
            if sp_pr is None:
                sp_pr = etree.SubElement(target, qn('c:spPr'))
            for child in list(sp_pr):
                sp_pr.remove(child)
            etree.SubElement(sp_pr, qn('a:noFill'))
    
        plot_area = chart._chartSpace.find(f'.//{qn("c:plotArea")}')
        if plot_area is not None:
            pa_sp_pr = plot_area.find(qn('c:spPr'))
            if pa_sp_pr is None:
                pa_sp_pr = etree.SubElement(plot_area, qn('c:spPr'))
            for child in list(pa_sp_pr):
                pa_sp_pr.remove(child)
            etree.SubElement(pa_sp_pr, qn('a:noFill'))
    
        # --- Axis & gridline styling (bar/line only) ---
        if chart_type not in ("pie", "donut"):
            val_axis_def = elem.get("valueAxis", {})
            cat_axis_def = elem.get("categoryAxis", {})
            
            # Value axis scaling
            if val_axis_def.get("min") is not None:
                chart.value_axis.minimum_scale = val_axis_def["min"]
            if val_axis_def.get("max") is not None:
                chart.value_axis.maximum_scale = val_axis_def["max"]
            if val_axis_def.get("majorUnit") is not None:
                chart.value_axis.major_unit = val_axis_def["majorUnit"]
            
            # Gridlines
            val_axis_el = chart._chartSpace.find(f'.//{qn("c:valAx")}')
            if val_axis_def.get("gridlines") is False:
                # Remove gridlines
                if val_axis_el is not None:
                    for gl in val_axis_el.findall(qn('c:majorGridlines')):
                        val_axis_el.remove(gl)
            elif val_axis_el is not None:
                major_gl = val_axis_el.find(qn('c:majorGridlines'))
                if major_gl is None:
                    major_gl = etree.SubElement(val_axis_el, qn('c:majorGridlines'))
                gl_sp_pr = major_gl.find(qn('c:spPr'))
                if gl_sp_pr is None:
                    gl_sp_pr = etree.SubElement(major_gl, qn('c:spPr'))
                for child in list(gl_sp_pr):
                    gl_sp_pr.remove(child)
                ln = etree.SubElement(gl_sp_pr, qn('a:ln'))
                ln.set('w', gridline_width)
                solid = etree.SubElement(ln, qn('a:solidFill'))
                srgb = etree.SubElement(solid, qn('a:srgbClr'))
                srgb.set('val', gridline_color)
                prstDash = etree.SubElement(ln, qn('a:prstDash'))
                prstDash.set('val', gridline_dash)
    
            # Axis lines: subtle solid (skip for roundtrip)
            if not val_axis_def and not cat_axis_def:
                for ax_tag in ['c:valAx', 'c:catAx']:
                    ax_el = chart._chartSpace.find(f'.//{qn(ax_tag)}')
                    if ax_el is None:
                        continue
                    ax_sp_pr = ax_el.find(qn('c:spPr'))
                    if ax_sp_pr is None:
                        ax_sp_pr = etree.SubElement(ax_el, qn('c:spPr'))
                    for old_ln in ax_sp_pr.findall(qn('a:ln')):
                        ax_sp_pr.remove(old_ln)
                    ln = etree.SubElement(ax_sp_pr, qn('a:ln'))
                    ln.set('w', axis_width)
                    solid = etree.SubElement(ln, qn('a:solidFill'))
                    srgb = etree.SubElement(solid, qn('a:srgbClr'))
                    srgb.set('val', axis_line_color)
            else:
                # Roundtrip: apply axis line settings
                if val_axis_def.get("line") == "none":
                    ax_el = chart._chartSpace.find(f'.//{qn("c:valAx")}')
                    if ax_el is not None:
                        ax_sp_pr = ax_el.find(qn('c:spPr'))
                        if ax_sp_pr is None:
                            ax_sp_pr = etree.SubElement(ax_el, qn('c:spPr'))
                        for old_ln in ax_sp_pr.findall(qn('a:ln')):
                            ax_sp_pr.remove(old_ln)
                        ln = etree.SubElement(ax_sp_pr, qn('a:ln'))
                        etree.SubElement(ln, qn('a:noFill'))
    
            # Tick labels: brand font, muted color (skip if roundtrip with axis defs)
            chart_font_sz = elem.get("_chartFontSize")
            if not val_axis_def and not cat_axis_def:
                for axis in [chart.value_axis, chart.category_axis]:
                    axis.has_title = False
                    axis.tick_labels.font.size = Pt(tick_font_size)
                    axis.tick_labels.font.color.rgb = text_muted
                    axis.tick_labels.font.name = self.fonts["halfwidth"]
            elif chart_font_sz:
                for axis in [chart.value_axis, chart.category_axis]:
                    axis.has_title = False
                    axis.tick_labels.font.size = Pt(chart_font_sz / 100)
            
            # Tick marks
            if val_axis_def.get("tickMark") == "none":
                chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE
            if cat_axis_def.get("tickMark") == "none":
                chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE
    
            # Remove minor gridlines
            if val_axis_el is not None:
                minor_gl = val_axis_el.find(qn('c:minorGridlines'))
                if minor_gl is not None:
                    val_axis_el.remove(minor_gl)
    
        # --- Bar-specific: gap width & overlap ---
        if chart_type == "bar":
            x_chart = plot_area.find(qn('c:barChart'))
            if x_chart is not None:
                # Gap width: tighter grouping
                gap_w = x_chart.find(qn('c:gapWidth'))
                if gap_w is None:
                    gap_w = etree.SubElement(x_chart, qn('c:gapWidth'))
                gap_w.set('val', str(elem.get("gapWidth") or style_gap_width or ('80' if not stacked else '60')))
                # Overlap for clustered
                if not stacked:
                    overlap = x_chart.find(qn('c:overlap'))
                    if overlap is None:
                        overlap = etree.SubElement(x_chart, qn('c:overlap'))
                    overlap.set('val', str(elem.get("overlap", -10)))
    
        # --- Line-specific: smooth curves ---
        if chart_type == "line" and smooth:
            for line_chart in plot_area.findall(qn('c:lineChart')):
                for ser in line_chart.findall(qn('c:ser')):
                    sm = ser.find(qn('c:smooth'))
                    if sm is None:
                        sm = etree.SubElement(ser, qn('c:smooth'))
                    sm.set('val', '1')
    
        # --- Legend ---
        # Hide legend for single-series bar (categories shown on axis)
        if chart_type == "bar" and len(series_defs) <= 1:
            show_legend = elem.get("legend", False)
        chart.has_legend = show_legend
        if show_legend:
            legend_pos_map = {
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "right": XL_LEGEND_POSITION.RIGHT,
                "left": XL_LEGEND_POSITION.LEFT,
                "top": XL_LEGEND_POSITION.TOP,
            }
            chart.legend.position = legend_pos_map.get(style_legend_position, XL_LEGEND_POSITION.BOTTOM)
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(legend_font_size)
            chart.legend.font.color.rgb = text_muted
            chart.legend.font.name = self.fonts["halfwidth"]
    
        # --- Series styling ---
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            color_str = series_defs[i].get("color") if i < len(series_defs) else None
            if not color_str:
                color_str = default_series_colors[i % len(default_series_colors)]
            else:
                color_str = color_str.lstrip("#")
            rgb = RGBColor.from_string(color_str)
    
            if chart_type in ("pie", "donut"):
                for j, point in enumerate(series.points):
                    pc = series_defs[0].get("colors", []) if series_defs else []
                    c = pc[j].lstrip("#") if j < len(pc) else default_series_colors[j % len(default_series_colors)]
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(c)
            elif chart_type == "line":
                series.format.line.color.rgb = rgb
                series.format.line.width = Pt(style_line_width or 2.5)
                if markers:
                    series.marker.style = XL_MARKER_STYLE.CIRCLE
                    series.marker.size = style_marker_size or 8
                    series.marker.format.fill.solid()
                    series.marker.format.fill.fore_color.rgb = rgb
                    # White marker border for polish
                    series.marker.format.line.color.rgb = rgb
                    series.marker.format.line.width = Pt(1.5)
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = rgb
                # Subtle border on bars for definition
                series.format.line.color.rgb = rgb
                series.format.line.width = Pt(0)
                # Per-point colors
                point_colors = series_defs[i].get("pointColors") if i < len(series_defs) else None
                point_xml = series_defs[i].get("_pointXml") if i < len(series_defs) else None
                if point_colors:
                    for pt_idx_str, pt_color in point_colors.items():
                        pt_idx = int(pt_idx_str)
                        try:
                            series.points[pt_idx].format.fill.solid()
                            series.points[pt_idx].format.fill.fore_color.rgb = RGBColor.from_string(pt_color.lstrip("#"))
                        except Exception:
                            pass
                if point_xml:
                    # Apply raw XML for gradient points
                    ser_el = series._element
                    for pt_idx_str, xml_str in point_xml.items():
                        pt_idx = int(pt_idx_str)
                        try:
                            dPt = etree.SubElement(ser_el, qn('c:dPt'))
                            idx_el = etree.SubElement(dPt, qn('c:idx'))
                            idx_el.set('val', str(pt_idx))
                            spPr = etree.fromstring(xml_str)
                            dPt.append(spPr)
                        except Exception:
                            pass
    
        # --- Data labels ---
        if show_data_labels:
            plot.has_data_labels = True
            dl = plot.data_labels
            val_axis_def = elem.get("valueAxis", {})
            cat_axis_def = elem.get("categoryAxis", {})
            if not val_axis_def and not cat_axis_def:
                dl.font.size = Pt(style_font_size or 10)
                dl.font.color.rgb = text_color
                dl.font.name = self.fonts["halfwidth"]
            elif chart_font_sz:
                dl.font.size = Pt(chart_font_sz / 100)
            if chart_type in ("pie", "donut"):
                nf = number_format or "0%"
                dl.number_format = nf
                dl.show_percentage = ("%" in nf)
                dl.show_value = ("%" not in nf)
                dl.show_category_name = False
            elif number_format:
                dl.number_format = number_format
    
        # Donut hole size
        if chart_type == "donut" and elem.get("holeSize") is not None:
            from pptx.oxml.ns import qn
            doughnut = chart._chartSpace.find('.//' + qn('c:doughnutChart'))
            if doughnut is not None:
                hs = doughnut.find(qn('c:holeSize'))
                if hs is not None:
                    hs.set('val', str(elem["holeSize"]))
    

