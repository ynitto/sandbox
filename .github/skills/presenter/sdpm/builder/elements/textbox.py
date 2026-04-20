# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Textbox element."""
from pptx.dml.color import RGBColor
from pptx.util import Pt
from sdpm.schema.defaults import ELEMENT_DEFAULTS
from sdpm.utils.effects import apply_effects
from sdpm.utils.text import _expand_styled_newlines

_DEFAULTS = ELEMENT_DEFAULTS["textbox"]


class TextboxMixin:
    """Mixin providing textbox element methods."""

    def _add_textbox(self, slide, elem):
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        
        align = elem.get("align", _DEFAULTS["align"])
        x_pct = elem.get("x", 3)
        y_pct = elem.get("y", 25)
        width_pct = elem.get("width", 94)
        auto_width = elem.get("autoWidth", False)
        
        # Convert percentage to EMU (x is always left edge)
        x_emu = self._px_to_emu(x_pct)
        y_emu = self._px_to_emu(y_pct)
        width_emu = self._px_to_emu(width_pct)
        height_px = elem.get("height")
        height_emu = self._px_to_emu(height_px) if height_px else self._px_to_emu(10)
        
        textbox = slide.shapes.add_textbox(x_emu, y_emu, width_emu, height_emu)
        
        # Inject lstStyle XML if provided (preserves placeholder formatting)
        lst_style_xml = elem.get("_lstStyle")
        if lst_style_xml:
            from lxml import etree
            from pptx.oxml.ns import qn
            new_lst = etree.fromstring(lst_style_xml)
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            # If paragraphs have bullets, remove buNone from lstStyle
            has_bullets = any(
                (p.get('list') if isinstance(p, dict) else False)
                for p in (elem.get('paragraphs') or [])
            )
            if has_bullets:
                for lvl in new_lst:
                    buNone = lvl.find(f'{{{ns}}}buNone')
                    if buNone is not None:
                        lvl.remove(buNone)
            txBody = textbox.text_frame._txBody
            old_lst = txBody.find(qn('a:lstStyle'))
            if old_lst is not None:
                txBody.replace(old_lst, new_lst)
            else:
                txBody.insert(1, new_lst)  # after bodyPr
        
        tf = textbox.text_frame
        if elem.get("_noAutofit"):
            # Clean bodyPr to match original (no autofit, no wrap override)
            from pptx.oxml.ns import qn
            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            for attr in list(bodyPr.attrib):
                if 'wrap' in attr:
                    del bodyPr.attrib[attr]
            for child in list(bodyPr):
                tag = child.tag.split('}')[1]
                if tag in ('spAutoFit', 'noAutofit', 'normAutofit'):
                    bodyPr.remove(child)
        else:
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        if not elem.get("_noAutofit"):
            tf.word_wrap = not auto_width
        
        # Apply rotation
        rotation = elem.get("rotation", _DEFAULTS["rotation"])
        if rotation != 0:
            textbox.rotation = rotation
        
        # Apply flip via XML
        flip_h = elem.get("flipH", False)
        flip_v = elem.get("flipV", False)
        
        if flip_h or flip_v:
            try:
                xfrm = textbox._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    if flip_h:
                        xfrm.set('flipH', '1')
                    if flip_v:
                        xfrm.set('flipV', '1')
            except Exception:
                pass
        
        # Apply fill color
        fill_color = elem.get("fill")
        
        if fill_color == "none":
            textbox.fill.background()
        elif fill_color:
            textbox.fill.solid()
            hex_color = fill_color.lstrip("#")
            textbox.fill.fore_color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        else:
            textbox.fill.background()
        
        # Apply fill opacity
        opacity = elem.get("opacity")
        if opacity is not None and 0 <= opacity < 1:
            self._set_fill_opacity(textbox, opacity)
        
        # Set paragraph alignment
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "just": PP_ALIGN.JUSTIFY}
        align = elem.get("align")
        has_lst_style = bool(elem.get("_lstStyle"))
        
        # Check if paragraphs array exists (multiple paragraphs without bullets)
        paragraphs = elem.get("paragraphs")
        
        if paragraphs:
            # Add multiple paragraphs (with or without bullets/numbering)
            default_font_size = None if has_lst_style else elem.get("fontSize")
            font_family = elem.get("fontFamily")
            fc = elem.get("fontColor")
            default_color = RGBColor.from_string(fc.lstrip('#')) if fc else None
            for i, para_item in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check if it's a dict with bullet/numbering info or just a string
                if isinstance(para_item, dict):
                    para_text = para_item.get("text", "")
                    list_def = para_item.get("list")
                else:
                    para_text = para_item
                    list_def = None
                
                # Apply list (bullet or numbering)
                if list_def and isinstance(list_def, dict):
                    list_type = list_def.get("type", "disc")
                    p.level = list_def.get("level", 0)
                    if list_type == "disc":
                        self._set_bullet(p)
                    else:
                        self._set_numbering(p, list_type)
                
                # Apply text
                if para_text:
                    para_font_size = para_item.get("fontSize", default_font_size) if isinstance(para_item, dict) else default_font_size
                    self._apply_styled_text(p, para_text, default_color=default_color, default_font_size=para_font_size, font_family=font_family, no_default_font=has_lst_style)
                
                # Apply line spacing percentage (must be before spcBef/spcAft in XML)
                if isinstance(para_item, dict) and para_item.get("lineSpacingPct"):
                    from lxml import etree
                    from pptx.oxml.ns import qn as _qn
                    pPr = p._element.get_or_add_pPr()
                    lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
                    spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
                    spcPct.set('val', str(para_item["lineSpacingPct"]))
                
                # Apply space after
                if isinstance(para_item, dict) and para_item.get("spaceAfter") is not None:
                    p.space_after = Pt(para_item["spaceAfter"] / 100)
                if isinstance(para_item, dict) and para_item.get("spaceBefore") is not None:
                    p.space_before = Pt(para_item["spaceBefore"] / 100)
                
                # Apply marL/indent
                if isinstance(para_item, dict) and para_item.get("marL") is not None:
                    pPr = p._element.get_or_add_pPr()
                    pPr.set('marL', str(para_item["marL"]))
                if isinstance(para_item, dict) and para_item.get("indent") is not None:
                    pPr = p._element.get_or_add_pPr()
                    pPr.set('indent', str(para_item["indent"]))
                
                # Apply paragraph level
                if isinstance(para_item, dict) and para_item.get("level"):
                    pPr = p._element.get_or_add_pPr()
                    pPr.set('lvl', str(para_item["level"]))
                
                
                para_align = para_item.get("align") if isinstance(para_item, dict) else None
                effective_align = para_align or align
                if effective_align or not has_lst_style:
                    p.alignment = align_map.get(effective_align, PP_ALIGN.LEFT)
                
                # Add buNone if lstStyle has bullet defs but paragraph has no bullet
                if has_lst_style and not (isinstance(para_item, dict) and para_item.get("list")):
                    from lxml import etree as _et
                    from pptx.oxml.ns import qn as _qn
                    pPr = p._element.get_or_add_pPr()
                    if pPr.find(_qn('a:buChar')) is None and pPr.find(_qn('a:buAutoNum')) is None:
                        _et.SubElement(pPr, _qn('a:buNone'))
        else:
            # Single text — expand styled tags spanning newlines, then split
            text = elem.get("text", "")
            
            text = text.replace("\\n", "\n")
            text = _expand_styled_newlines(text)
            default_font_size = None if has_lst_style else elem.get("fontSize")
            font_family = elem.get("fontFamily")
            fc = elem.get("fontColor")
            default_color = RGBColor.from_string(fc.lstrip('#')) if fc else None
            lines = text.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                self._apply_styled_text(p, line, default_color=default_color, default_font_size=default_font_size, font_family=font_family, no_default_font=has_lst_style)
                if align or not has_lst_style:
                    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
                if has_lst_style:
                    from lxml import etree as _et
                    from pptx.oxml.ns import qn as _qn
                    pPr = p._element.get_or_add_pPr()
                    if pPr.find(_qn('a:buChar')) is None and pPr.find(_qn('a:buAutoNum')) is None:
                        _et.SubElement(pPr, _qn('a:buNone'))
                # Apply indent/marL override
                if elem.get("indent") is not None or elem.get("marL") is not None:
                    from pptx.oxml.ns import qn as _qn
                    pPr = p._element.get_or_add_pPr()
                    if elem.get("indent") is not None:
                        pPr.set('indent', str(elem["indent"]))
                    if elem.get("marL") is not None:
                        pPr.set('marL', str(elem["marL"]))
        
        # Apply line spacing (single-text) - must be before buNone in pPr
        lnSpcPct = elem.get("lineSpacingPct")
        if lnSpcPct:
            from lxml import etree as _et
            from pptx.oxml.ns import qn as _qn
            for p in tf.paragraphs:
                pPr = p._element.get_or_add_pPr()
                lnSpc = _et.Element(_qn('a:lnSpc'))
                spcPct = _et.SubElement(lnSpc, _qn('a:spcPct'))
                spcPct.set('val', str(lnSpcPct))
                pPr.insert(0, lnSpc)
        
        # Apply character spacing
        spc = elem.get("_spc")
        if spc is not None:
            from pptx.oxml.ns import qn as _qn
            for p in tf.paragraphs:
                for r_el in p._element.findall(_qn('a:r')):
                    rPr = r_el.find(_qn('a:rPr'))
                    if rPr is None:
                        from lxml import etree as _et
                        rPr = _et.SubElement(r_el, _qn('a:rPr'))
                    rPr.set('spc', str(spc))
        
        # Apply line (border) to textbox
        line_color = elem.get("line")
        line_width = elem.get("lineWidth", 1)
        
        if line_color == "none" or line_color is None:
            textbox.line.fill.background()
        elif line_color:
            textbox.line.fill.solid()
            hex_color = line_color.lstrip("#")
            textbox.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            textbox.line.width = Pt(line_width)
        
        # Apply text frame margins (px input → EMU)
        margin_left = elem.get("marginLeft")
        if margin_left is not None:
            tf.margin_left = self._px_to_emu(margin_left)
        margin_top = elem.get("marginTop")
        if margin_top is not None:
            tf.margin_top = self._px_to_emu(margin_top)
        margin_right = elem.get("marginRight")
        if margin_right is not None:
            tf.margin_right = self._px_to_emu(margin_right)
        margin_bottom = elem.get("marginBottom")
        if margin_bottom is not None:
            tf.margin_bottom = self._px_to_emu(margin_bottom)
        
        # Apply vertical align
        va = elem.get("verticalAlign")
        if va is not None:
            _va_map = {"top": 1, "middle": 3, "bottom": 4}
            tf.vertical_anchor = _va_map.get(va, 1)
        
        # Text direction
        if elem.get("textDirection"):
            from pptx.oxml.ns import qn
            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('vert', elem["textDirection"])
        
        # Apply text gradient
        text_gradient = elem.get("textGradient")
        if text_gradient:
            self._apply_text_gradient(textbox, text_gradient)
        
        # Apply per-run text gradient
        grad_runs = elem.get("_textGradientRuns")
        if grad_runs:
            self._apply_text_gradient_runs(textbox, grad_runs)
        
        # Override cap=all and bold from lstStyle
        if elem.get("_capNone") or elem.get("_boldOff"):
            from pptx.oxml.ns import qn as _qn
            for r in textbox._element.findall('.//' + _qn('a:r')):
                rPr = r.find(_qn('a:rPr'))
                if rPr is None:
                    from lxml import etree
                    rPr = etree.SubElement(r, _qn('a:rPr'))
                    r.insert(0, rPr)
                if elem.get("_capNone"):
                    rPr.set('cap', 'none')
                if elem.get("_boldOff"):
                    rPr.set('b', '0')
        
        # Apply visual effects
        apply_effects(textbox._element, elem, self.EMU_PER_PX)
    

