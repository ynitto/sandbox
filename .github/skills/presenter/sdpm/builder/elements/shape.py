# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Shape and freeform elements."""
import sys
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from sdpm.schema.defaults import ELEMENT_DEFAULTS
from sdpm.utils.effects import apply_effects
from sdpm.utils.text import _expand_styled_newlines

_DEFAULTS = ELEMENT_DEFAULTS["shape"]


class ShapeMixin:
    """Mixin providing shape element methods."""

    def _add_shape(self, slide, elem):
        """Add shape to slide.
        
        Supported shapes:
        - rectangle, rounded_rectangle
        - oval, circle
        - arrow_right, arrow_left, arrow_up, arrow_down
        - triangle, diamond, pentagon, hexagon
        
        Args:
            shape: Shape type (required)
            x, y: Position in percentage (required)
            width, height: Size in percentage (required)
            fill: Fill color (#RRGGBB or "none")
            line: Line color (#RRGGBB or "none")
            lineWidth: Line width in pt (default: 1)
            text: Text inside shape (optional)
            fontSize: Font size for text (default: 14)
        """
        shape_type = elem.get("shape")
        if not shape_type:
            return
        
        # Map shape names to MSO_SHAPE constants
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "circle": MSO_SHAPE.OVAL,  # Circle is oval with equal width/height
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "arrow_up": MSO_SHAPE.UP_ARROW,
            "arrow_down": MSO_SHAPE.DOWN_ARROW,
            "arrow_circular": MSO_SHAPE.CIRCULAR_ARROW,
            "arrow_left_right": MSO_SHAPE.LEFT_RIGHT_ARROW,
            "arrow_up_down": MSO_SHAPE.UP_DOWN_ARROW,
            "arrow_curved_right": MSO_SHAPE.CURVED_RIGHT_ARROW,
            "arrow_curved_left": MSO_SHAPE.CURVED_LEFT_ARROW,
            "arrow_curved_up": MSO_SHAPE.CURVED_UP_ARROW,
            "arrow_curved_down": MSO_SHAPE.CURVED_DOWN_ARROW,
            "arrow_circular_left": MSO_SHAPE.LEFT_CIRCULAR_ARROW,
            "arrow_circular_left_right": MSO_SHAPE.LEFT_RIGHT_CIRCULAR_ARROW,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "cross": MSO_SHAPE.CROSS,
            "trapezoid": MSO_SHAPE.TRAPEZOID,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "chevron": MSO_SHAPE.CHEVRON,
            "donut": MSO_SHAPE.DONUT,
            "arc": MSO_SHAPE.ARC,
            "block_arc": MSO_SHAPE.BLOCK_ARC,
            "chord": MSO_SHAPE.CHORD,
            "pie": MSO_SHAPE.PIE,
            "pie_wedge": MSO_SHAPE.PIE_WEDGE,
            "cloud": MSO_SHAPE.CLOUD,
            "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
            "star_5_point": MSO_SHAPE.STAR_5_POINT,
            "no_symbol": MSO_SHAPE.NO_SYMBOL,
            "callout_rectangle": MSO_SHAPE.RECTANGULAR_CALLOUT,
            "callout_rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
            "callout_oval": MSO_SHAPE.OVAL_CALLOUT,
            "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
            "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
            "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
            "left_brace": MSO_SHAPE.LEFT_BRACE,
            "right_brace": MSO_SHAPE.RIGHT_BRACE,
            "left_bracket": MSO_SHAPE.LEFT_BRACKET,
            "right_bracket": MSO_SHAPE.RIGHT_BRACKET,
            "cylinder": MSO_SHAPE.CAN,
            "math_not_equal": MSO_SHAPE.MATH_NOT_EQUAL,
        }
        
        mso_shape = shape_map.get(shape_type)
        raw_prst = None
        if not mso_shape:
            # Use rectangle as base, then override prstGeom with raw preset name
            mso_shape = MSO_SHAPE.RECTANGLE
            raw_prst = shape_type
        
        x_pct = elem.get("x", 10)
        y_pct = elem.get("y", 10)
        width_pct = elem.get("width", 20)
        height_pct = elem.get("height", 10)
        
        x_emu = self._px_to_emu(x_pct)
        y_emu = self._px_to_emu(y_pct)
        width_emu = self._px_to_emu(width_pct)
        height_emu = self._px_to_emu(height_pct)
        
        # For circle, use minimum of width/height
        if shape_type == "circle":
            size = min(width_emu, height_emu)
            width_emu = height_emu = size
        
        shape = slide.shapes.add_shape(mso_shape, x_emu, y_emu, width_emu, height_emu)
        
        # Override prstGeom for unknown shape types
        if raw_prst:
            from pptx.oxml.ns import qn
            sp_pr = shape._element.find(qn('p:spPr'))
            for pg in sp_pr.findall(qn('a:prstGeom')):
                pg.set('prst', raw_prst)
        
        # Apply rotation
        rotation = elem.get("rotation", _DEFAULTS["rotation"])
        if rotation != 0:
            shape.rotation = rotation
        
        # Apply adjustments (shape control points)
        adjustments = elem.get("adjustments")
        if adjustments and hasattr(shape, 'adjustments'):
            try:
                if shape_type in ("block_arc", "arc"):
                    adjustments = self._convert_arc_adjustments(shape_type, list(adjustments))
                for i, adj_val in enumerate(adjustments):
                    if i < len(shape.adjustments):
                        shape.adjustments[i] = adj_val
            except Exception:
                pass
        
        # Apply fill color or gradient
        fill_color = elem.get("fill")
        gradient = elem.get("gradient")
        
        if gradient:
            self._apply_gradient_fill(shape, gradient)
        elif fill_color == "none" or fill_color is None:
            # Transparent fill
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.fill.background()  # Set to transparent
        elif fill_color:
            shape.fill.solid()
            hex_color = fill_color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            # Apply fill opacity (0.0 = fully transparent, 1.0 = fully opaque)
            opacity = elem.get("opacity")
            if opacity is not None and 0 <= opacity < 1:
                self._set_fill_opacity(shape, opacity)
        
        # Apply pattern fill
        pattern_fill = elem.get("patternFill")
        if pattern_fill:
            from lxml import etree
            from pptx.oxml.ns import qn
            sp_pr = shape._element.spPr
            # Remove existing fill
            for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:pattFill'):
                for el in sp_pr.findall(qn(tag)):
                    sp_pr.remove(el)
            patt = etree.SubElement(sp_pr, qn('a:pattFill'))
            patt.set('prst', pattern_fill.get("pattern", "dkDnDiag"))
            fg = etree.SubElement(patt, qn('a:fgClr'))
            fg_srgb = etree.SubElement(fg, qn('a:srgbClr'))
            fg_srgb.set('val', pattern_fill.get("fgColor", "#FFFFFF").lstrip("#"))
            bg = etree.SubElement(patt, qn('a:bgClr'))
            bg_srgb = etree.SubElement(bg, qn('a:srgbClr'))
            bg_srgb.set('val', pattern_fill.get("bgColor", "#000000").lstrip("#"))
        
        # Apply line color or gradient
        line_gradient = elem.get("lineGradient")
        line_color = elem.get("line")
        line_width = elem.get("lineWidth", 1)
        
        if line_gradient:
            fwd = {"lineGradient": line_gradient}
            if elem.get("lineWidth") is not None:
                fwd["lineWidth"] = line_width
            self._apply_shape_formatting(shape, fwd)
        elif line_color == "none" or line_color is None:
            shape.line.fill.background()
        elif line_color:
            shape.line.fill.solid()
            hex_color = line_color.lstrip("#")
            shape.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            shape.line.width = Pt(line_width)
        
        # Apply dash style to shape line
        self._apply_dash_style(shape, elem)
        
        # Add text if specified
        text = elem.get("text")
        items = elem.get("items")
        paragraphs = elem.get("paragraphs")
        
        if (text or items or paragraphs) and shape.has_text_frame:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf = shape.text_frame
            
            # Inject lstStyle XML if provided
            lst_style_xml = elem.get("_lstStyle")
            if lst_style_xml:
                from lxml import etree
                from pptx.oxml.ns import qn
                txBody = tf._txBody
                old_lst = txBody.find(qn('a:lstStyle'))
                new_lst = etree.fromstring(lst_style_xml)
                if old_lst is not None:
                    txBody.replace(old_lst, new_lst)
                else:
                    txBody.insert(1, new_lst)
            
            tf.word_wrap = True
            tf.auto_size = None if elem.get("_noAutofit") else (MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT if elem.get("_spAutoFit") else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
            tf.clear()
            
            # Apply margins (px input → EMU)
            if elem.get("marginLeft") is not None:
                tf.margin_left = self._px_to_emu(elem["marginLeft"])
            if elem.get("marginTop") is not None:
                tf.margin_top = self._px_to_emu(elem["marginTop"])
            if elem.get("marginRight") is not None:
                tf.margin_right = self._px_to_emu(elem["marginRight"])
            if elem.get("marginBottom") is not None:
                tf.margin_bottom = self._px_to_emu(elem["marginBottom"])
            
            # Apply vertical align (default: middle)
            _va_map = {"top": 1, "middle": 3, "bottom": 4}
            va = elem.get("verticalAlign", "middle")
            tf.vertical_anchor = _va_map.get(va, 3)
            
            # Text direction (vert270, vert, etc.)
            if elem.get("textDirection") or elem.get("autoWidth"):
                bodyPr = tf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                if bodyPr is not None:
                    if elem.get("textDirection"):
                        bodyPr.set('vert', elem["textDirection"])
                    if elem.get("autoWidth"):
                        bodyPr.set('wrap', 'none')
            
            font_size = elem.get("fontSize", 14)
            
            if items:
                # Add bullet points
                for i, item in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.level = 0
                    self._set_bullet(p)
                    item_text = item.get("text", item) if isinstance(item, dict) else item
                    self._apply_styled_text(p, item_text, default_font_size=font_size)
                    
                    if isinstance(item, dict) and item.get("spaceAfter") is not None:
                        p.space_after = Pt(item["spaceAfter"] / 100)
                    
                    # Apply text alignment
                    align = elem.get("align", "center")
                    if align == "center":
                        p.alignment = 2
                    elif align == "right":
                        p.alignment = 3
                    elif align == "left":
                        p.alignment = 1
                    else:
                        p.alignment = 2  # Default center
            elif paragraphs:
                # Paragraphs array with per-paragraph styles
                align_map = {"left": 1, "center": 2, "right": 3}
                for i, para_def in enumerate(paragraphs):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    if isinstance(para_def, str):
                        para_def = {"text": para_def}
                    para_text = para_def.get("text", "")
                    para_fs = para_def.get("fontSize")  # None = inherit from shape
                    # Bullet
                    bu = para_def.get("bullet")
                    if bu:
                        p.level = 0
                        self._set_bullet(p, char=bu if isinstance(bu, str) else None)
                    # Text
                    if para_text:
                        self._apply_styled_text(p, para_text, default_font_size=para_fs)
                    # endParaRPr font size
                    end_fs = para_def.get("endFontSize", para_fs if not para_text else None)
                    if end_fs:
                        from lxml import etree
                        from pptx.oxml.ns import qn as _qn
                        endPr = p._element.find(_qn('a:endParaRPr'))
                        if endPr is None:
                            endPr = etree.SubElement(p._element, _qn('a:endParaRPr'))
                        endPr.set('sz', str(int(end_fs * 100)))
                    # Alignment
                    pa = para_def.get("align")
                    if pa:
                        p.alignment = align_map.get(pa, 1)
            else:
                # Single text
                text = text.replace("\\n", "\n")
                text = _expand_styled_newlines(text)
                lines = text.split("\n")
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    self._apply_styled_text(p, line, default_font_size=font_size)
                    # Set endParaRPr font size for empty paragraphs
                    if not line.strip():
                        from lxml import etree
                        from pptx.oxml.ns import qn as _qn
                        endPr = p._element.find(_qn('a:endParaRPr'))
                        if endPr is None:
                            endPr = etree.SubElement(p._element, _qn('a:endParaRPr'))
                        endPr.set('sz', str(int(font_size * 100)))
                
                # Apply text alignment
                align = elem.get("align", "center")
                align_val = {"center": 2, "right": 3, "left": 1}.get(align, 2)
                for p in tf.paragraphs:
                    p.alignment = align_val
                
                # Apply indent/marL
                if elem.get("indent") is not None or elem.get("marL") is not None:
                    from pptx.oxml.ns import qn as _qn
                    for p in tf.paragraphs:
                        pPr = p._element.get_or_add_pPr()
                        if elem.get("indent") is not None:
                            pPr.set('indent', str(elem["indent"]))
                        if elem.get("marL") is not None:
                            pPr.set('marL', str(elem["marL"]))
        
        # Add hyperlink to entire shape if specified
        link = elem.get("link")
        if link:
            shape.click_action.hyperlink.address = link
        
        # Apply flip
        if elem.get("flipH") or elem.get("flipV"):
            from pptx.oxml.ns import qn as _qn
            xfrm = shape._element.spPr.find(_qn('a:xfrm'))
            if xfrm is not None:
                if elem.get("flipH"):
                    xfrm.set('flipH', '1')
                if elem.get("flipV"):
                    xfrm.set('flipV', '1')
        
        # Apply arrow heads
        arrow_start = elem.get("arrowStart")
        arrow_end = elem.get("arrowEnd")
        if arrow_start or arrow_end:
            from lxml import etree as _et
            from pptx.oxml.ns import qn as _qn
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            ln = shape._element.spPr.find(_qn('a:ln'))
            if ln is None:
                ln = _et.SubElement(shape._element.spPr, _qn('a:ln'))
            if arrow_end:
                te = _et.SubElement(ln, f'{{{ns_a}}}tailEnd')
                te.set('type', arrow_end)
            if arrow_start:
                he = _et.SubElement(ln, f'{{{ns_a}}}headEnd')
                he.set('type', arrow_start)
        
        # Apply character spacing
        spc = elem.get("_spc")
        if spc is not None and shape.has_text_frame:
            from pptx.oxml.ns import qn as _qn
            for p in shape.text_frame.paragraphs:
                for r_el in p._element.findall(_qn('a:r')):
                    rPr = r_el.find(_qn('a:rPr'))
                    if rPr is None:
                        from lxml import etree as _et
                        rPr = _et.SubElement(r_el, _qn('a:rPr'))
                    rPr.set('spc', str(spc))
        
        # Apply visual effects
        apply_effects(shape._element, elem, self.EMU_PER_PX)

    @staticmethod
    def _apply_dash_style(shape, elem):
        """Apply dash style to shape line."""
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        dash_style = elem.get("dashStyle")
        if not dash_style:
            return
        dash_map = {
            "solid": MSO_LINE_DASH_STYLE.SOLID,
            "dash": MSO_LINE_DASH_STYLE.DASH,
            "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
            "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
            "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
            "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
            "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
            "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
        }
        if dash_style in dash_map:
            shape.line.dash_style = dash_map[dash_style]
        elif dash_style in ("sysDash", "sysDot", "sysDashDot", "sysDashDotDot"):
            from lxml import etree
            from pptx.oxml.ns import qn
            ln = shape.line._ln
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', dash_style)

    @staticmethod
    def _convert_arc_adjustments(shape_type, adj):
        """Convert arc/block_arc adjustments from intuitive to PowerPoint format.

        Input:  [start_angle, sweep, clockwise?, thickness?]
          - start_angle: degrees (0=right/3h, 90=bottom/6h, 180=left/9h, 270=top/12h)
          - sweep: degrees or "N%" string (e.g. "73%" → 262.8°)
          - clockwise: bool (default True). block_arc: 4th param, arc: 3rd param
          - thickness: float 0-1 (block_arc only, default 0.25)

        PowerPoint draws block_arc/arc counter-clockwise from adj[0] to adj[1].
        To get clockwise, we swap start and end.
        """
        if len(adj) < 2:
            return adj
        start = adj[0]
        sweep = adj[1]
        # Parse "N%" string
        if isinstance(sweep, str) and sweep.endswith("%"):
            sweep = float(sweep.rstrip("%")) / 100 * 360
        # Determine clockwise flag and thickness based on shape type
        if shape_type == "block_arc":
            thickness = adj[2] if len(adj) >= 3 and not isinstance(adj[2], bool) else 0.25
            cw_idx = 3 if len(adj) >= 3 and not isinstance(adj[2], bool) else 2
            clockwise = adj[cw_idx] if len(adj) > cw_idx and isinstance(adj[cw_idx], bool) else True
        else:  # arc
            thickness = None
            clockwise = adj[2] if len(adj) >= 3 and isinstance(adj[2], bool) else True
        # Convert: OOXML blockArc/arc uses adj1=startAngle, adj2=endAngle,
        # drawing clockwise from adj1 to adj2.
        # Our API: [startAngle, sweepDegrees_or_percent, thickness?, clockwise?]
        # Convert sweep to endAngle: adj2 = start + sweep (cw) or start - sweep (ccw)
        if clockwise:
            ppt_start = start
            ppt_end = (start + sweep) % 360
        else:
            ppt_start = (start - sweep) % 360
            ppt_end = start
        # python-pptx adjustment unit: degrees * 0.6 (60000 OOXML units / 100000 scale)
        result = [ppt_start * 0.6, ppt_end * 0.6]
        if thickness is not None:
            result.append(thickness)
        return result

    @staticmethod
    def _build_path_command(path_el, cmd, ns_a, px_to_emu):
        """Build a single OOXML path command element."""
        from lxml import etree
        c = cmd.get("cmd", "")
        if c == "M":
            mv = etree.SubElement(path_el, f'{{{ns_a}}}moveTo')
            pt = etree.SubElement(mv, f'{{{ns_a}}}pt')
            pt.set('x', str(px_to_emu(cmd["x"])))
            pt.set('y', str(px_to_emu(cmd["y"])))
        elif c == "L":
            ln = etree.SubElement(path_el, f'{{{ns_a}}}lnTo')
            pt = etree.SubElement(ln, f'{{{ns_a}}}pt')
            pt.set('x', str(px_to_emu(cmd["x"])))
            pt.set('y', str(px_to_emu(cmd["y"])))
        elif c == "C":
            cb = etree.SubElement(path_el, f'{{{ns_a}}}cubicBezTo')
            for p in cmd["pts"]:
                pt = etree.SubElement(cb, f'{{{ns_a}}}pt')
                pt.set('x', str(px_to_emu(p[0])))
                pt.set('y', str(px_to_emu(p[1])))
        elif c == "Q":
            qb = etree.SubElement(path_el, f'{{{ns_a}}}quadBezTo')
            for p in cmd["pts"]:
                pt = etree.SubElement(qb, f'{{{ns_a}}}pt')
                pt.set('x', str(px_to_emu(p[0])))
                pt.set('y', str(px_to_emu(p[1])))
        elif c == "A":
            arc = etree.SubElement(path_el, f'{{{ns_a}}}arcTo')
            arc.set('wR', str(px_to_emu(cmd["wR"])))
            arc.set('hR', str(px_to_emu(cmd["hR"])))
            arc.set('stAng', str(round(cmd["stAng"] * 60000)))
            arc.set('swAng', str(round(cmd["swAng"] * 60000)))
        elif c == "Z":
            etree.SubElement(path_el, f'{{{ns_a}}}close')

    def _add_freeform_shape(self, slide, elem):
        """Add freeform shape with custom geometry."""
        from lxml import etree
        from pptx.oxml.ns import qn
        
        x_emu = elem.get("_xEmu") or self._px_to_emu(elem.get("x", 0))
        y_emu = elem.get("_yEmu") or self._px_to_emu(elem.get("y", 0))
        width_emu = elem.get("_widthEmu") or self._px_to_emu(elem.get("width", 10))
        height_emu = elem.get("_heightEmu") or self._px_to_emu(elem.get("height", 10))
        
        # Create rectangle as placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_emu, y_emu, width_emu, height_emu)
        
        # Build custom geometry from path array or raw XML
        custom_geom = elem.get("customGeometry")
        
        # Resolve paths: "paths" (multi) or "path" (single, backward compat)
        paths_def = elem.get("paths")
        path_cmds = elem.get("path")
        if paths_def:
            paths = []
            for p in paths_def:
                if isinstance(p, list):
                    paths.append({"commands": p})
                else:
                    paths.append(p)
        elif path_cmds:
            paths = [{"commands": path_cmds}]
        else:
            paths = []
        
        if paths:
            try:
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                cust_geom = etree.SubElement(etree.Element('dummy'), f'{{{ns_a}}}custGeom')
                etree.SubElement(cust_geom, f'{{{ns_a}}}avLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}gdLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}ahLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}cxnLst')
                rect = etree.SubElement(cust_geom, f'{{{ns_a}}}rect')
                rect.set('l', 'l')
                rect.set('t', 't')
                rect.set('r', 'r')
                rect.set('b', 'b')

                # Use raw pathLst XML if available (lossless roundtrip)
                raw_path_xml = elem.get("_pathLstXml")
                if raw_path_xml:
                    path_lst = etree.fromstring(raw_path_xml)
                    cust_geom.append(path_lst)
                else:
                    path_lst = etree.SubElement(cust_geom, f'{{{ns_a}}}pathLst')
                    # px → EMU, no scaling — path coords = shape interior px
                    def px_to_emu(v):
                        return round(v * self.EMU_PER_PX)

                    for path_def in paths:
                        cmds = path_def["commands"]
                        path_el = etree.SubElement(path_lst, f'{{{ns_a}}}path')
                        path_el.set('w', str(width_emu))
                        path_el.set('h', str(height_emu))
                        # Path fill mode
                        fill_mode = path_def.get("fill")
                        if fill_mode is None:
                            has_close = any(
                                c.get("cmd") == "Z" for c in cmds
                            )
                            fill_mode = "norm" if has_close else "none"
                        if fill_mode != "norm":
                            path_el.set('fill', fill_mode)
                        for cmd in cmds:
                            self._build_path_command(
                                path_el, cmd, ns_a, px_to_emu
                            )
                
                sp_pr = shape._element.spPr
                for prst_geom in sp_pr.findall(qn('a:prstGeom')):
                    sp_pr.remove(prst_geom)
                xfrm = sp_pr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(cust_geom)
                
                parent = shape._element
                for style in parent.findall(qn('p:style')):
                    parent.remove(style)
            except Exception as e:
                print(f"Warning: Failed to build freeform path: {e}", file=sys.stderr)
        elif custom_geom:
            try:
                custom_geom_clean = custom_geom.replace('ns0:', 'a:').replace('xmlns:ns0=', 'xmlns:a=')
                cust_geom_elem = etree.fromstring(custom_geom_clean)
                sp_pr = shape._element.spPr
                
                for prst_geom in sp_pr.findall(qn('a:prstGeom')):
                    sp_pr.remove(prst_geom)
                
                xfrm = sp_pr.find(qn('a:xfrm'))
                if xfrm is not None:
                    children = list(sp_pr)
                    xfrm_idx = children.index(xfrm)
                    sp_pr.insert(xfrm_idx + 1, cust_geom_elem)
                
                # Remove style element
                parent = shape._element
                for style in parent.findall(qn('p:style')):
                    parent.remove(style)
            except Exception as e:
                print(f"Warning: Failed to apply freeform geometry: {e}", file=sys.stderr)
        
        # Apply formatting (fill, line, gradient)
        self._apply_shape_formatting(shape, elem)
        
        # Apply dash style
        self._apply_dash_style(shape, elem)
        
        # Apply arrow heads
        if elem.get("headEnd") or elem.get("tailEnd"):
            from lxml import etree
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            ln = shape._element.spPr.find(qn('a:ln'))
            if ln is None:
                ln = etree.SubElement(shape._element.spPr, qn('a:ln'))
            if elem.get("tailEnd"):
                te = etree.SubElement(ln, f'{{{ns_a}}}tailEnd')
                te.set('type', elem["tailEnd"])
            if elem.get("headEnd"):
                he = etree.SubElement(ln, f'{{{ns_a}}}headEnd')
                he.set('type', elem["headEnd"])
        
        # Apply rotation
        rotation = elem.get("rotation", _DEFAULTS["rotation"])
        if rotation != 0:
            shape.rotation = rotation
        
        # Apply flip
        if elem.get("flipH") or elem.get("flipV"):
            sp_pr = shape._element.spPr
            xfrm = sp_pr.find(qn('a:xfrm'))
            if xfrm is not None:
                if elem.get("flipH"):
                    xfrm.set('flipH', '1')
                if elem.get("flipV"):
                    xfrm.set('flipV', '1')
        # Text in freeform
        text = elem.get("text")
        items = elem.get("items")
        if (text or items) and shape.has_text_frame:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = None if elem.get("_noAutofit") else (MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT if elem.get("_spAutoFit") else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
            tf.clear()
            if elem.get("marginLeft") is not None:
                tf.margin_left = self._px_to_emu(elem["marginLeft"])
            if elem.get("marginTop") is not None:
                tf.margin_top = self._px_to_emu(elem["marginTop"])
            if elem.get("marginRight") is not None:
                tf.margin_right = self._px_to_emu(elem["marginRight"])
            if elem.get("marginBottom") is not None:
                tf.margin_bottom = self._px_to_emu(elem["marginBottom"])
            va = elem.get("verticalAlign")
            if va is not None:
                _va_map = {"top": 1, "middle": 3, "bottom": 4}
                tf.vertical_anchor = _va_map.get(va, 3)
            font_size = elem.get("fontSize", 14)
            if items:
                for i, item in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.level = 0
                    self._set_bullet(p)
                    item_text = item.get("text", item) if isinstance(item, dict) else item
                    self._apply_styled_text(p, item_text, default_font_size=font_size)
            else:
                text = text.replace("\\n", "\n")
                text = _expand_styled_newlines(text)
                for i, line in enumerate(text.split("\n")):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    self._apply_styled_text(p, line, default_font_size=font_size)
                    if not line.strip():
                        from lxml import etree
                        from pptx.oxml.ns import qn as _qn
                        endPr = p._element.find(_qn('a:endParaRPr'))
                        if endPr is None:
                            endPr = etree.SubElement(p._element, _qn('a:endParaRPr'))
                        endPr.set('sz', str(int(font_size * 100)))
            align_val = {"center": 2, "right": 3, "left": 1}.get(elem.get("align", "center"), 2)
            for p in tf.paragraphs:
                p.alignment = align_val
        
        # Apply visual effects
        apply_effects(shape._element, elem, self.EMU_PER_PX)
    

