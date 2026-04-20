# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Image element."""
import sys
from pathlib import Path
from pptx.util import Emu
from sdpm.schema.defaults import ELEMENT_DEFAULTS
from sdpm.utils.image import resolve_image_path, apply_image_effects
from sdpm.utils.effects import apply_effects
from sdpm.utils.svg import _recolor_svg, get_svg_dimensions, generate_qr_svg, add_svg_to_slide
from sdpm.utils.text import _expand_styled_newlines
from sdpm.assets import is_recolor_protected

_DEFAULTS = ELEMENT_DEFAULTS["image"]

class ImageMixin:
    """Mixin providing image element methods."""

    def _add_image(self, slide, elem):
        """Add image element to slide.
        
        src: icons:NAME or file path (supports ~)
        """
        from pptx.enum.text import PP_ALIGN
        
        src = elem.get("src") or elem.get("path", "")
        x_pct = elem.get("x", 0)
        y_pct = elem.get("y", 0)
        width_pct = elem.get("width")
        height_pct = elem.get("height")
        label = elem.get("label")
        label_pos = elem.get("labelPosition", "bottom")
        label_size = elem.get("labelSize", 11)
        link = elem.get("link")
        rotation = elem.get("rotation", _DEFAULTS["rotation"])
        icon_color = elem.get("iconColor")
        
        if not src:
            return
        
        # QR code generation
        if src.startswith("qr:"):
            qr_url = src[3:]
            qr_size = int(width_pct or height_pct or 200)
            svg_bytes = generate_qr_svg(
                qr_url, size=qr_size,
                color=elem.get("color"),
                gradient=elem.get("gradient"),
                theme="dark" if self.is_dark else "light",
            )
            x = self._px_to_emu(x_pct)
            y = self._px_to_emu(y_pct)
            w = self._px_to_emu(width_pct or 200)
            h = self._px_to_emu(height_pct or 200)
            pic = add_svg_to_slide(slide, svg_bytes, x, y, w, h)
            if rotation != 0 and pic is not None:
                xfrm = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rotation * 60000)))
            # Label for QR (reuse inline label logic)
            if label and label_pos != "none":
                from pptx.enum.text import PP_ALIGN
                label_margin = int(h * 0.04)
                if label_pos == "bottom":
                    lbl_x, lbl_y, lbl_w = x, y + h + label_margin, w
                elif label_pos == "right":
                    lbl_x, lbl_y, lbl_w = x + w + Emu(100000), y + h // 3, self._px_to_emu(15)
                else:
                    return
                textbox = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, Emu(300000))
                tf = textbox.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
                self._apply_styled_text(p, label, default_font_size=label_size)
            return
        
        # Resolve image path
        if src.startswith("icons:") or src.startswith("assets:"):
            img_path = resolve_image_path(src, "dark" if self.is_dark else "light")
        else:
            img_path = Path(src).expanduser()
            if not img_path.is_absolute():
                img_path = self._base_dir / src
            if not img_path.exists():
                print(f"Warning: Image not found: {img_path}", file=sys.stderr)
                return
        
        is_svg = img_path.suffix.lower() == ".svg"
        
        # Prepare SVG bytes (with optional recolor)
        svg_bytes = None
        if is_svg:
            svg_bytes = img_path.read_bytes()
            if src and is_recolor_protected(src):
                if icon_color:
                    print(f"Warning: iconColor ignored (recolor-protected asset): {src}", file=sys.stderr)
            else:
                effective_icon_color = icon_color or self.theme_colors["text"]
                recolored = _recolor_svg(svg_bytes, effective_icon_color)
                if recolored:
                    svg_bytes = recolored
        elif icon_color:
            print(f"Warning: iconColor ignored (not SVG): {img_path.name}", file=sys.stderr)
        
        # Calculate dimensions
        x = self._px_to_emu(x_pct)
        y = self._px_to_emu(y_pct)
        
        if width_pct:
            width = self._px_to_emu(width_pct)
            if height_pct:
                height = self._px_to_emu(height_pct)
            else:
                # Maintain original aspect ratio
                if is_svg:
                    img_w, img_h = get_svg_dimensions(img_path)
                else:
                    from PIL import Image
                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                    except Exception:
                        img_w, img_h = 1, 1
                height = int(width * img_h / img_w) if img_w > 0 else width
        
        if is_svg:
            if not width_pct:
                img_w, img_h = get_svg_dimensions(img_path)
                width = self._px_to_emu(img_w)
                height = self._px_to_emu(img_h)
            pic = add_svg_to_slide(slide, svg_bytes, x, y, width, height)
        else:
            if width_pct:
                pic = slide.shapes.add_picture(str(img_path), x, y, width=width, height=height)
            else:
                pic = slide.shapes.add_picture(str(img_path), x, y)
                width = pic.width
                height = pic.height
        
        # Apply rotation
        if rotation != 0:
            if is_svg:
                # Set rotation directly on XML xfrm element
                xfrm = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rotation * 60000)))
            else:
                pic.rotation = rotation
        
        # Add hyperlink if specified
        if link:
            if is_svg:
                pass  # Hyperlinks on SVG not supported via direct XML (rare use case)
            else:
                pic.click_action.hyperlink.address = link
        
        # Add label if specified
        if label and label_pos != "none":
            # Scale margin proportionally to icon size (base: 4% of height)
            label_margin = int(height * 0.04)
            if label_pos == "bottom":
                lbl_x = x
                lbl_y = y + height + label_margin
                lbl_w = width
            elif label_pos == "right":
                lbl_x = x + width + label_margin * 2
                lbl_y = y + height // 3
                lbl_w = self._px_to_emu(15)
            else:
                return
            
            textbox = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, Emu(300000))
            tf = textbox.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
            label = label.replace("\\n", "\n")
            label = _expand_styled_newlines(label)
            lines = label.split("\n")
            align = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
            for li, line in enumerate(lines):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = align
                self._apply_styled_text(p, line, default_font_size=label_size)
        
        # Apply visual effects to image
        if pic is not None:
            pic_el = pic._element if hasattr(pic, '_element') else pic
            apply_image_effects(pic_el, elem)
            apply_effects(pic_el, elem, self.EMU_PER_PX)
    

    def _add_video(self, slide, elem):
        """Add video element to slide."""
        src = elem.get("src", "")
        poster = elem.get("poster", "")
        if not src:
            return
        
        video_path = self._base_dir / src
        poster_path = self._base_dir / poster if poster else None
        
        if not video_path or not Path(video_path).exists():
            print(f"Warning: Video not found: {src}", file=sys.stderr)
            return
        
        x = Emu(elem.get("_xEmu") or self._px_to_emu(elem.get("x", 0)))
        y = Emu(elem.get("_yEmu") or self._px_to_emu(elem.get("y", 0)))
        w = Emu(elem.get("_widthEmu") or self._px_to_emu(elem.get("width", 100)))
        h = Emu(elem.get("_heightEmu") or self._px_to_emu(elem.get("height", 100)))
        
        # Determine MIME type
        ext = Path(video_path).suffix.lower()
        mime_map = {'.mp4': 'video/mp4', '.avi': 'video/avi', '.wmv': 'video/x-ms-wmv', '.mov': 'video/quicktime'}
        mime = mime_map.get(ext, 'video/mp4')
        
        slide.shapes.add_movie(
            str(video_path),
            x, y, w, h,
            poster_frame_image=str(poster_path) if poster_path and Path(poster_path).exists() else None,
            mime_type=mime,
        )
