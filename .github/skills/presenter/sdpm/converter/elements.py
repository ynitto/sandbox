# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Element extraction (shape, textbox, line, freeform, picture, group)."""
import json
import sys
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from .constants import _NS, EMU_PER_PX, _base_element, _add_flip, _serialize_lstStyle, _hex
from .color import _resolve_color_with_transforms
from .xml_helpers import (extract_line_dash, _resolve_line_from_style, _extract_fill_from_xml,
                          _extract_line_from_xml, _extract_effects_from_xml, _extract_visual_effects)
from .text import _extract_styled_text, _detect_font_size, _get_alignment, _extract_shape_text

_SHAPE_MAP = {
    MSO_SHAPE.RECTANGLE: "rectangle", MSO_SHAPE.ROUNDED_RECTANGLE: "rounded_rectangle",
    MSO_SHAPE.OVAL: "oval", MSO_SHAPE.RIGHT_ARROW: "arrow_right",
    MSO_SHAPE.LEFT_ARROW: "arrow_left", MSO_SHAPE.UP_ARROW: "arrow_up",
    MSO_SHAPE.DOWN_ARROW: "arrow_down", MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
    MSO_SHAPE.DIAMOND: "diamond", MSO_SHAPE.PENTAGON: "pentagon",
    MSO_SHAPE.HEXAGON: "hexagon", MSO_SHAPE.CHEVRON: "chevron",
    MSO_SHAPE.RIGHT_BRACE: "right_brace", MSO_SHAPE.LEFT_BRACE: "left_brace",
    60: "arrow_circular",
}
_PRESET_MAP = {
    'roundRect': 'rounded_rectangle', 'rect': 'rectangle', 'ellipse': 'oval',
    'triangle': 'triangle', 'diamond': 'diamond', 'pentagon': 'pentagon',
    'hexagon': 'hexagon', 'chevron': 'chevron', 'homePlate': 'pentagon',
    'heart': 'heart', 'cloud': 'cloud', 'lightningBolt': 'lightning_bolt',
    'star5': 'star_5_point', 'noSmoking': 'no_symbol', 'cross': 'cross', 'plus': 'cross',
    'trapezoid': 'trapezoid', 'parallelogram': 'parallelogram',
    'donut': 'donut', 'arc': 'arc', 'blockArc': 'block_arc', 'chord': 'chord',
    'pie': 'pie', 'pieWedge': 'pie_wedge',
    'leftRightArrow': 'arrow_left_right', 'upDownArrow': 'arrow_up_down',
    'curvedRightArrow': 'arrow_curved_right', 'curvedLeftArrow': 'arrow_curved_left',
    'curvedUpArrow': 'arrow_curved_up', 'curvedDownArrow': 'arrow_curved_down',
    'circularArrow': 'arrow_circular', 'leftCircularArrow': 'arrow_circular_left',
    'leftRightCircularArrow': 'arrow_circular_left_right',
    'calloutRoundRect': 'callout_rounded_rectangle', 'wedgeRoundRectCallout': 'callout_rounded_rectangle',
    'calloutRect': 'callout_rectangle', 'wedgeRectCallout': 'callout_rectangle',
    'calloutEllipse': 'callout_oval', 'wedgeEllipseCallout': 'callout_oval',
    'flowChartProcess': 'flowchart_process', 'flowChartDecision': 'flowchart_decision',
    'flowChartTerminator': 'flowchart_terminator',
    'leftBracket': 'left_bracket', 'rightBracket': 'right_bracket',
    'can': 'cylinder', 'mathNotEqual': 'math_not_equal',
}

def _resolve_shape_name(shape):
    """Resolve shape preset name from python-pptx or XML."""
    if shape.shape_type == 5:  # MSO_SHAPE_TYPE.FREEFORM
        return "rounded_rectangle"
    name = None
    try:
        if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type in _SHAPE_MAP:
            name = _SHAPE_MAP[shape.auto_shape_type]
    except Exception:
        pass
    if not name:
        try:
            prst = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst is not None:
                prst_val = prst.get('prst')
                name = _PRESET_MAP.get(prst_val, prst_val)  # Use raw prst value as fallback
        except Exception:
            pass
    if name == "oval" and shape.width == shape.height:
        return "circle"
    return name or "rounded_rectangle"

def extract_line_element(shape, theme_colors=None, color_mapping=None, theme_styles=None):
    """Extract line/connector as element dict."""
    try:
        # Build x1/y1/x2/y2 from bounding box + flip
        x = round(shape.left / EMU_PER_PX)
        y = round(shape.top / EMU_PER_PX)
        w = round(shape.width / EMU_PER_PX)
        h = round(shape.height / EMU_PER_PX)
        x1, y1, x2, y2 = x, y, x + w, y + h

        # Absorb flip into coordinates
        try:
            xfrm = shape._element.spPr.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                if xfrm.get('flipH') == '1':
                    x1, x2 = x2, x1
                if xfrm.get('flipV') == '1':
                    y1, y2 = y2, y1
        except Exception:
            pass

        elem = {"type": "line", "x1": x1, "y1": y1, "x2": x2, "y2": y2}
        
        # Extract connector type from XML
        try:
            sp_pr = shape._element.spPr
            prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst_geom is not None:
                prst = prst_geom.get('prst')
                if prst:
                    # Save exact preset type
                    elem["preset"] = prst
                    
                    # Map to general connector type
                    if 'straight' in prst.lower():
                        elem["connectorType"] = "straight"
                    elif 'bent' in prst.lower():
                        elem["connectorType"] = "elbow"
                    elif 'curved' in prst.lower():
                        elem["connectorType"] = "curved"
                
                # Extract adjustments
                av_lst = prst_geom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if av_lst is not None:
                    adjustments = []
                    for gd in av_lst.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adj_val = int(fmla.split()[1])
                            adjustments.append(adj_val / 100000.0)
                    if adjustments:
                        elem["adjustments"] = adjustments
        except Exception:
            elem["connectorType"] = "straight"  # default
        
        # Extract arrow heads from XML
        try:
            ln = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                if head_end is not None:
                    head_type = head_end.get('type')
                    if head_type:
                        elem["arrowStart"] = head_type
                
                tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                if tail_end is not None:
                    tail_type = tail_end.get('type')
                    if tail_type:
                        elem["arrowEnd"] = tail_type
        except Exception:
            pass
        
        # Extract line color or gradient (use XML helper)
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        line_info = _extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping)
        if "line" in line_info and line_info["line"] != "none":
            elem["color"] = line_info["line"]
        elif line_info.get("line") == "none":
            elem["color"] = "none"
        else:
            # Resolve from style reference
            style_info = _resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles)
            if style_info.get("line"):
                elem["color"] = style_info["line"]
            if style_info.get("lineWidth"):
                elem["lineWidth"] = style_info["lineWidth"]
        if "lineGradient" in line_info:
            elem["lineGradient"] = line_info["lineGradient"]
        if "lineWidth" in line_info:
            elem["lineWidth"] = line_info["lineWidth"]
        
        # Extract dash style
        dash = extract_line_dash(shape)
        if dash:
            elem["dashStyle"] = dash
        
        return elem
    except Exception as e:
        print(f"Warning: Failed to extract line: {e}", file=sys.stderr)
        return None

def extract_freeform_element(shape, theme_colors=None, color_mapping=None, builder_text_color=None):
    """Extract freeform/curve shape as element dict with path commands in px."""
    try:
        sp_pr = shape._element.spPr
        cust_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}custGeom')
        if cust_geom is None:
            return None

        elem = _base_element(shape, "freeform")
        _add_flip(elem, shape)
        # Preserve exact EMU size for freeform roundtrip fidelity
        elem["_widthEmu"] = shape.width
        elem["_heightEmu"] = shape.height

        # Extract all paths
        path_elements = cust_geom.findall('.//a:pathLst/a:path', _NS)
        if not path_elements:
            return None

        def _extract_path_commands(path_el):
            """Extract commands from a single path element."""
            path_w = int(path_el.get('w', shape.width))
            path_h = int(path_el.get('h', shape.height))
            sx = shape.width / path_w if path_w else 1
            sy = shape.height / path_h if path_h else 1

            def to_px(x, y):
                return round(int(x) * sx / EMU_PER_PX, 1), round(int(y) * sy / EMU_PER_PX, 1)

            path = []
            for child in path_el:
                tag = child.tag.split('}')[-1]
                pts = child.findall('a:pt', _NS)
                if tag == 'moveTo' and pts:
                    px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                    path.append({"cmd": "M", "x": px, "y": py})
                elif tag == 'lnTo' and pts:
                    px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                    path.append({"cmd": "L", "x": px, "y": py})
                elif tag == 'cubicBezTo' and len(pts) == 3:
                    coords = [to_px(p.get('x'), p.get('y')) for p in pts]
                    path.append({"cmd": "C", "pts": [[c[0], c[1]] for c in coords]})
                elif tag == 'quadBezTo' and len(pts) == 2:
                    coords = [to_px(p.get('x'), p.get('y')) for p in pts]
                    path.append({"cmd": "Q", "pts": [[c[0], c[1]] for c in coords]})
                elif tag == 'arcTo':
                    wR = int(child.get('wR', 0))
                    hR = int(child.get('hR', 0))
                    stAng = int(child.get('stAng', 0))
                    swAng = int(child.get('swAng', 0))
                    path.append({
                        "cmd": "A",
                        "wR": round(wR * sx / EMU_PER_PX, 1),
                        "hR": round(hR * sy / EMU_PER_PX, 1),
                        "stAng": round(stAng / 60000, 2),
                        "swAng": round(swAng / 60000, 2),
                    })
                elif tag == 'close':
                    path.append({"cmd": "Z"})
            return path

        if len(path_elements) == 1:
            # Single path → "path" key (backward compatible)
            elem["path"] = _extract_path_commands(path_elements[0])
            fill_attr = path_elements[0].get('fill')
            if fill_attr and fill_attr != 'norm':
                elem["pathFill"] = fill_attr
        else:
            # Multiple paths → "paths" key
            paths = []
            for pe in path_elements:
                p = {"commands": _extract_path_commands(pe)}
                fill_attr = pe.get('fill')
                if fill_attr and fill_attr != 'norm':
                    p["fill"] = fill_attr
                paths.append(p)
            elem["paths"] = paths

        # Preserve raw pathLst XML for lossless roundtrip
        path_el_first = path_elements[0]
        path_w = int(path_el_first.get('w', shape.width))
        path_h = int(path_el_first.get('h', shape.height))
        if path_w == shape.width and path_h == shape.height:
            from lxml import etree as _et
            pathLst = cust_geom.find('.//a:pathLst', _NS)
            if pathLst is not None:
                elem["_pathLstXml"] = _et.tostring(pathLst, encoding='unicode')

        # Fill
        elem.update(_extract_fill_from_xml(sp_pr, theme_colors, color_mapping))

        # Line
        line_info = _extract_line_from_xml(sp_pr, theme_colors, color_mapping)
        elem.update(line_info)

        # Line opacity
        ln = sp_pr.find(f'.//{{{_NS["a"]}}}ln')
        if ln is not None:
            solid = ln.find(f'{{{_NS["a"]}}}solidFill')
            if solid is not None:
                for clr_tag in ('srgbClr', 'schemeClr'):
                    clr = solid.find(f'{{{_NS["a"]}}}{clr_tag}')
                    if clr is not None:
                        alpha = clr.find(f'{{{_NS["a"]}}}alpha')
                        if alpha is not None:
                            elem["lineOpacity"] = round(int(alpha.get('val')) / 100000, 2)
                        break

        # Arrow heads
        try:
            if ln is not None:
                for attr, tag in [("headEnd", "headEnd"), ("tailEnd", "tailEnd")]:
                    el = ln.find(f'{{{_NS["a"]}}}{tag}')
                    if el is not None and el.get('type'):
                        elem[attr] = el.get('type')
        except Exception:
            pass

        # Effects
        elem.update(_extract_visual_effects(sp_pr, theme_colors, color_mapping))

        # Text (if freeform contains text)
        if shape.has_text_frame and shape.text_frame.text.strip():
            _extract_shape_text(shape, elem, theme_colors, color_mapping, builder_text_color=builder_text_color)

        return elem
    except Exception as e:
        print(f"Warning: Failed to extract freeform: {e}", file=sys.stderr)
        return None

def extract_shape_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, builder_text_color=None):
    """Extract shape as element dict."""
    try:
        elem = {
            "type": "shape",
            "x": round(shape.left / EMU_PER_PX),
            "y": round(shape.top / EMU_PER_PX),
            "width": round(shape.width / EMU_PER_PX),
            "height": round(shape.height / EMU_PER_PX),
            "shape": _resolve_shape_name(shape)
        }
        _add_flip(elem, shape)
        
        # Style references
        style_fill_idx = None
        style_fill_color = None
        try:
            style = shape._element.find(f'{{{_NS["p"]}}}style')
            if style is not None:
                fill_ref = style.find(f'{{{_NS["a"]}}}fillRef')
                if fill_ref is not None:
                    style_fill_idx = int(fill_ref.get('idx', 0))
                    sc = fill_ref.find(f'{{{_NS["a"]}}}schemeClr')
                    if sc is not None:
                        style_fill_color = sc.get('val')
        except Exception:
            pass
        
        # Rotation
        if shape.rotation != 0:
            elem["rotation"] = round(shape.rotation, 1)
        
        # Adjustments (only if explicitly set in XML avLst)
        try:
            sp_pr = shape._element.spPr
            prst_geom = sp_pr.find(f'{{{_NS["a"]}}}prstGeom')
            if prst_geom is not None:
                av_lst = prst_geom.find(f'{{{_NS["a"]}}}avLst')
                if av_lst is not None and len(av_lst) > 0:
                    adjs = []
                    for gd in av_lst.findall(f'{{{_NS["a"]}}}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adjs.append(round(int(fmla.split()[1]) / 100000, 5))
                    if adjs:
                        elem["adjustments"] = adjs
        except Exception:
            pass
        
        # Extract fill and line from XML
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))
        
        # Fill (XML first, python-pptx API fallback for style references)
        try:
            fill_info = _extract_fill_from_xml(sp_pr_xml, theme_colors, color_mapping)
            # Check if spPr has explicit <a:noFill>
            has_explicit_no_fill = sp_pr_xml is not None and sp_pr_xml.find(f'{{{_NS["a"]}}}noFill') is not None
            if fill_info.get("fill") != "none" or "gradient" in fill_info or "patternFill" in fill_info:
                elem.update(fill_info)
            elif has_explicit_no_fill:
                elem["fill"] = "none"
            else:
                if shape.fill.type == 1:  # SOLID
                    if shape.fill.fore_color.type == 1:  # RGB
                        rgb = shape.fill.fore_color.rgb
                        elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                    elif shape.fill.fore_color.type == 2:  # SCHEME
                        theme_color = shape.fill.fore_color.theme_color
                        if theme_colors and theme_color in theme_colors:
                            elem["fill"] = theme_colors[theme_color]
                    alpha_el = sp_pr_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha') if sp_pr_xml is not None else None
                    if alpha_el is not None:
                        elem["opacity"] = round(int(alpha_el.get('val', 100000)) / 1000, 1)
                elif shape.fill.type == 3:  # GRADIENT
                    try:
                        stops = []
                        for stop in shape.fill.gradient_stops:
                            s = {"position": round(stop.position, 3)}
                            if stop.color.type == 1:
                                rgb = stop.color.rgb
                                s["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                            elif stop.color.type == 2 and theme_colors and stop.color.theme_color in theme_colors:
                                s["color"] = theme_colors[stop.color.theme_color]
                            if "color" in s:
                                stops.append(s)
                        if stops:
                            angle = 0
                            try:
                                # python-pptx returns CCW angle; convert to CW (PowerPoint UI)
                                ccw = round(shape.fill.gradient_angle, 1)
                                angle = round((360 - ccw) % 360, 1)
                            except Exception:
                                pass
                            elem["gradient"] = {"stops": stops, "angle": angle}
                    except Exception:
                        pass
                elif shape.fill.type is None or shape.fill.type == 0 or shape.fill.type == 5:
                    # Resolve from style fillRef (unless useBgFill=1)
                    use_bg = shape._element.get('useBgFill') == '1'
                    if not use_bg and style_fill_idx and style_fill_idx > 0 and style_fill_color and theme_styles and theme_styles.get("fill"):
                        fill_idx = style_fill_idx - 1
                        if 0 <= fill_idx < len(theme_styles["fill"]):
                            from lxml import etree as _et
                            fill_xml = _et.fromstring(theme_styles["fill"][fill_idx])
                            scheme = fill_xml.find(f'.//{{{_NS["a"]}}}schemeClr')
                            if scheme is not None and scheme.get('val') == 'phClr':
                                resolved = _resolve_color_with_transforms(scheme, theme_colors, color_mapping, override_scheme=style_fill_color)
                                if resolved:
                                    elem["fill"] = resolved
                    if "fill" not in elem:
                        elem["fill"] = "none"
        except Exception:
            pass
        
        # Line (XML first, style reference fallback)
        try:
            line_info = _extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping)
            ln_xml = sp_pr_xml.find('a:ln', _NS) if sp_pr_xml is not None else None
            if line_info.get("line") not in (None, "none") or "lineGradient" in line_info:
                # If lineWidth missing, try style reference
                if "lineWidth" not in line_info:
                    style_info = _resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles)
                    if style_info.get("lineWidth"):
                        line_info["lineWidth"] = style_info["lineWidth"]
                elem.update(line_info)
            elif ln_xml is not None and len(ln_xml) > 0:
                elem.update(line_info)  # ln exists with noFill or explicit content
            else:
                elem.update(_resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles))
            dash = extract_line_dash(shape)
            if dash:
                elem["dashStyle"] = dash
            # Arrow heads
            if ln_xml is not None:
                for attr, tag in [("headEnd", "headEnd"), ("tailEnd", "tailEnd")]:
                    el = ln_xml.find(f'{{{_NS["a"]}}}{tag}')
                    if el is not None and el.get('type') and el.get('type') != 'none':
                        elem[attr] = el.get('type')
        except Exception:
            if "line" not in elem and "lineGradient" not in elem:
                elem["line"] = "none"
        if "line" not in elem and "lineGradient" not in elem:
            elem["line"] = "none"
        
        # Extract text with styles
        if shape.has_text_frame and shape.text.strip():
            _extract_shape_text(shape, elem, theme_colors, color_mapping, builder_text_color=builder_text_color)
        
        # Extract hyperlink
        try:
            if hasattr(shape, 'click_action') and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                elem["link"] = shape.click_action.hyperlink.address
            else:
                # Remove null link
                if "link" in elem and elem["link"] is None:
                    del elem["link"]
        except Exception:
            pass
        
        # Extract visual effects
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))
        
        # Preserve lstStyle for roundtrip fidelity (non-placeholder shapes)
        _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
        if _lst:
            elem["_lstStyle"] = _lst
        
        return elem
    except Exception as e:
        print(f"Warning: Failed to extract shape details: {e}", file=sys.stderr)
        return None

def extract_textbox_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, is_placeholder=False, builder_text_color=None):
    """Extract textbox as element dict."""
    # Check if it's actually a shape with preset geometry (not a plain textbox)
    try:
        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prst_geom is not None:
            prst = prst_geom.get('prst')
            # If it has any preset geometry (not just 'rect'), treat as shape
            if prst and prst != 'rect':
                # This is a shape with text, not a plain textbox
                return extract_shape_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    except Exception:
        pass
    
    elem = {
        "type": "textbox",
        "x": round(shape.left / EMU_PER_PX),  # px (1920x1080 basis)
        "y": round(shape.top / EMU_PER_PX),
        "width": round(shape.width / EMU_PER_PX),
    }
    
    # Extract height (for TEXT_TO_FIT_SHAPE auto-shrink)
    if shape.height:
        h_px = round(shape.height / EMU_PER_PX)
        if h_px > 10:
            elem["height"] = h_px
    # Extract rotation
    if shape.rotation != 0:
        elem["rotation"] = round(shape.rotation, 1)
    
    # Extract flip
    _add_flip(elem, shape)
    
    # Extract autoWidth
    try:
        body_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if body_pr is not None:
            if body_pr.get('wrap') == 'none':
                elem["autoWidth"] = True
            vert = body_pr.get('vert')
            if vert:
                elem["textDirection"] = vert
    except Exception:
        pass
    
    # Extract margins (EMU → px)
    tf = shape.text_frame
    if tf.margin_left is not None and tf.margin_left != 91440:
        elem["marginLeft"] = round(tf.margin_left / EMU_PER_PX)
    if tf.margin_top is not None and tf.margin_top != 45720:
        elem["marginTop"] = round(tf.margin_top / EMU_PER_PX)
    if tf.margin_right is not None and tf.margin_right != 91440:
        elem["marginRight"] = round(tf.margin_right / EMU_PER_PX)
    if tf.margin_bottom is not None and tf.margin_bottom != 45720:
        elem["marginBottom"] = round(tf.margin_bottom / EMU_PER_PX)
    
    # Extract fill and line using XML helpers
    try:
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        elem.update(_extract_fill_from_xml(sp_pr_xml, theme_colors, color_mapping))
        elem.update(_extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping))
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))
    except Exception:
        pass
    
    # Extract textGradient from runs with gradFill
    try:
        grad_runs = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rpr = run._r.find(f'{{{_NS["a"]}}}rPr')
                if rpr is not None:
                    grad = rpr.find(f'{{{_NS["a"]}}}gradFill')
                    if grad is not None:
                        stops = []
                        for gs in grad.findall(f'.//{{{_NS["a"]}}}gs'):
                            pos = round(int(gs.get('pos', '0')) / 100000, 2)
                            srgb = gs.find(f'{{{_NS["a"]}}}srgbClr')
                            if srgb is not None:
                                stops.append({"position": pos, "color": _hex(srgb)})
                        if stops:
                            angle = 0
                            lin = grad.find(f'{{{_NS["a"]}}}lin')
                            if lin is not None:
                                angle = round(int(lin.get('ang', '0')) / 60000)
                            grad_runs.append({"text": run.text, "gradient": {"angle": angle, "stops": stops}})
        if grad_runs:
            # Count total runs with text
            total_runs = sum(1 for p in shape.text_frame.paragraphs for r in p.runs if r.text.strip())
            grads = [json.dumps(gr["gradient"], sort_keys=True) for gr in grad_runs]
            # Promote to textGradient only if ALL runs have the same gradient
            if len(set(grads)) == 1 and len(grad_runs) >= total_runs:
                elem["textGradient"] = grad_runs[0]["gradient"]
            else:
                elem["_textGradientRuns"] = grad_runs
    except Exception:
        pass
    
    # Detect cap=none and bold=off overrides (when lstStyle has cap=all / b=1)
    try:
        _all_runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        if _all_runs:
            if all(r._r.find(f'{{{_NS["a"]}}}rPr') is not None and
                   r._r.find(f'{{{_NS["a"]}}}rPr').get('cap') == 'none'
                   for r in _all_runs):
                elem["_capNone"] = True
            if all(r._r.find(f'{{{_NS["a"]}}}rPr') is not None and
                   r._r.find(f'{{{_NS["a"]}}}rPr').get('b') == '0'
                   for r in _all_runs):
                elem["_boldOff"] = True
    except Exception:
        pass

    # Extract text with styles
    text_parts = []
    default_font_size = None
    
    # Determine default text color (must match builder's theme_colors["text"])
    # For placeholders, don't set default_text_color — lstStyle defines the actual default
    default_text_color = None
    if not is_placeholder:
        default_text_color = builder_text_color
        if not default_text_color and color_mapping and theme_colors:
            tx1_mapped = color_mapping.get('tx1', 'dk1')
            default_text_color = theme_colors.get(tx1_mapped)
    
    # Check if multiple paragraphs (should be items array)
    paragraphs_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    all_paragraphs = list(shape.text_frame.paragraphs)
    has_lstStyle = _serialize_lstStyle(shape) is not None
    
    if len(all_paragraphs) > 1:
        # Multiple paragraphs - extract as paragraphs with bullet info
        default_font_size = None if (is_placeholder or has_lstStyle) else _detect_font_size(all_paragraphs)
        paragraphs = []
        for paragraph in all_paragraphs:
            
            # Empty paragraph
            if not paragraph.text.strip():
                paragraphs.append({"text": ""})
                continue
            
            # Check for bullet or numbering
            has_bullet = False
            numbering_type = None
            bu_font = None
            mar_l = None
            indent = None
            space_after = None
            space_before = None
            line_spacing = None
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    bu_auto_num = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    bu_char = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    
                    if bu_auto_num is not None:
                        numbering_type = bu_auto_num.get('type', 'arabicPeriod')
                    elif bu_char is not None:
                        has_bullet = True
                    
                    bu_font_elem = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                    if bu_font_elem is not None:
                        bu_font = bu_font_elem.get('typeface')
                    mar_l = pPr.get('marL')
                    indent = pPr.get('indent')
                    spc_aft = pPr.find('.//a:spcAft/a:spcPts', _NS)
                    if spc_aft is not None:
                        space_after = spc_aft.get('val')
                    spc_bef = pPr.find('.//a:spcBef/a:spcPts', _NS)
                    if spc_bef is not None:
                        space_before = spc_bef.get('val')
                    ln_spc = pPr.find('.//a:lnSpc/a:spcPts', _NS)
                    if ln_spc is not None:
                        line_spacing = ('pts', ln_spc.get('val'))
                    else:
                        ln_spc_pct = pPr.find('.//a:lnSpc/a:spcPct', _NS)
                        if ln_spc_pct is not None:
                            line_spacing = ('pct', ln_spc_pct.get('val'))
            except Exception:
                pass
            
            item_text = _extract_styled_text(paragraph.runs, theme_colors, color_mapping, default_font_size=default_font_size, default_text_color=default_text_color, is_placeholder=is_placeholder, paragraph=paragraph)
            para_info = {"text": item_text}
            if has_bullet or numbering_type:
                list_def = {}
                if numbering_type:
                    list_def["type"] = numbering_type
                else:
                    list_def["type"] = "disc"
                level = paragraph.level if paragraph.level else 0
                if level > 0:
                    list_def["level"] = level
                para_info["list"] = list_def
            if bu_font:
                para_info["buFont"] = bu_font
            if mar_l is not None:
                para_info["marL"] = int(mar_l)
            if indent is not None:
                para_info["indent"] = int(indent)
            if space_after is not None:
                para_info["spaceAfter"] = int(space_after)
            if space_before is not None:
                para_info["spaceBefore"] = int(space_before)
            if line_spacing:
                if line_spacing[0] == 'pct':
                    para_info["lineSpacingPct"] = int(line_spacing[1])
                else:
                    para_info["lineSpacing"] = int(line_spacing[1])
            
            # Paragraph level (for sub-bullets)
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    lvl = pPr.get('lvl')
                    if lvl and lvl != '0':
                        para_info["level"] = int(lvl)
            except Exception:
                pass
            
            paragraphs.append(para_info)
        
        if paragraphs:
            elem["paragraphs"] = paragraphs
            
            # Add fontSize if not default
            if default_font_size and default_font_size != 18:
                elem["fontSize"] = default_font_size
            
            # Get alignment - per paragraph if mixed, top-level if uniform
            aligns = [_get_alignment(p) for p in paragraphs_with_text]
            unique = set(a for a in aligns if a)
            if len(unique) <= 1:
                align = aligns[0] if aligns else None
                if align and align != "left":
                    elem["align"] = align
            else:
                # Mixed alignment: set per-paragraph
                for para_info, paragraph in zip(paragraphs, shape.text_frame.paragraphs):
                    a = _get_alignment(paragraph)
                    if a:
                        para_info["align"] = a
            
            # Preserve lstStyle for roundtrip fidelity
            _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
            if _lst:
                elem["_lstStyle"] = _lst
            
            # Extract character spacing
            _spc_vals = set()
            for _p in shape.text_frame.paragraphs:
                for _r in _p.runs:
                    _rPr = _r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    _s = _rPr.get('spc') if _rPr is not None else None
                    if _s:
                        _spc_vals.add(int(_s))
            if len(_spc_vals) == 1:
                elem["_spc"] = _spc_vals.pop()
            
            return elem
    
    # Single paragraph - extract as text
    default_font_size = None if (is_placeholder or has_lstStyle) else _detect_font_size(shape.text_frame.paragraphs)
    for paragraph in shape.text_frame.paragraphs:
        text_parts.append(_extract_styled_text(paragraph.runs, theme_colors, color_mapping, default_font_size=default_font_size, default_text_color=default_text_color, is_placeholder=is_placeholder, paragraph=paragraph))
    
    elem["text"] = ''.join(text_parts)
    
    # Extract indent/marL from first paragraph
    if shape.text_frame.paragraphs:
        from pptx.oxml.ns import qn as _qn
        pPr = shape.text_frame.paragraphs[0]._element.find(_qn('a:pPr'))
        if pPr is not None:
            _indent = pPr.get('indent')
            if _indent is not None:
                elem["indent"] = int(_indent)
            _marL = pPr.get('marL')
            if _marL is not None:
                elem["marL"] = int(_marL)
    
    # Add fontSize if consistent
    if default_font_size:
        elem["fontSize"] = default_font_size
    
    # Detect alignment
    if shape.text_frame.paragraphs:
        align = _get_alignment(shape.text_frame.paragraphs[0])
        if align:
            elem["align"] = align
        # Line spacing from first paragraph
        pPr = shape.text_frame.paragraphs[0]._element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is not None:
            lnSpc_pct = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc/{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
            if lnSpc_pct is not None:
                elem["lineSpacingPct"] = int(lnSpc_pct.get('val'))
    
    # Extract visual effects
    try:
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if sp_pr_xml is None:
            sp_pr_xml = shape._element.spPr if hasattr(shape._element, 'spPr') else None
        elem.update(_extract_effects_from_xml(sp_pr_xml, theme_colors, color_mapping))
    except Exception:
        pass
    
    # Preserve lstStyle for roundtrip fidelity
    _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
    if _lst:
        elem["_lstStyle"] = _lst
    
    # Extract character spacing (spc) if uniform across all runs
    if shape.has_text_frame:
        spc_values = set()
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                spc = rPr.get('spc') if rPr is not None else None
                if spc:
                    spc_values.add(int(spc))
        if len(spc_values) == 1:
            elem["_spc"] = spc_values.pop()
    
    return elem
    """Extract SVG bytes from asvg:svgBlip if present. Returns bytes or None."""
    ASVG_NS = 'http://schemas.microsoft.com/office/drawing/2016/SVG/main'
    svg_blip = shape._element.find(f'.//{{{ASVG_NS}}}svgBlip')
    if svg_blip is None:
        return None
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    r_embed = svg_blip.get(f'{{{r_ns}}}embed')
    if not r_embed:
        return None
    try:
        part = shape.part.rels[r_embed].target_part
        return part.blob
    except (KeyError, Exception):
        return None

def extract_video_element(shape, output_dir=None, slide_idx=0, img_idx=0):
    """Extract video as element dict, saving video file and poster image."""
    from pptx.oxml.ns import qn as _qn
    elem = _base_element(shape, "video")
    try:
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        nvPr = shape._element.find(f'{_qn("p:nvPicPr")}/{_qn("p:nvPr")}')
        if nvPr is None:
            return None
        videoFile = nvPr.find(_qn('a:videoFile'))
        if videoFile is None:
            return None

        # Save video file
        r_link = videoFile.get(f'{{{r_ns}}}link')
        if r_link and output_dir:
            slide_part = shape.part
            rel = slide_part.rels[r_link]
            ext = rel.target_ref.split('.')[-1] or 'mp4'
            video_name = f"slide{slide_idx+1}_video{img_idx+1}.{ext}"
            media_dir = Path(output_dir) / "media"
            media_dir.mkdir(exist_ok=True)
            # Get blob via p14:media embed (more reliable)
            p14_ns = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
            media_el = nvPr.find(f'.//{{{p14_ns}}}media')
            if media_el is not None:
                r_embed = media_el.get(f'{{{r_ns}}}embed')
                if r_embed:
                    (media_dir / video_name).write_bytes(slide_part.rels[r_embed].target_part.blob)
            elem["src"] = f"media/{video_name}"

        # Save poster image
        blip = shape._element.find(f'{_qn("p:blipFill")}/{_qn("a:blip")}')
        if blip is not None and output_dir:
            r_embed = blip.get(f'{{{r_ns}}}embed')
            if r_embed:
                poster_part = shape.part.rels[r_embed].target_part
                poster_ext = poster_part.content_type.split('/')[-1].replace('jpeg', 'jpg')
                poster_name = f"slide{slide_idx+1}_poster{img_idx+1}.{poster_ext}"
                images_dir = Path(output_dir) / "images"
                images_dir.mkdir(exist_ok=True)
                (images_dir / poster_name).write_bytes(poster_part.blob)
                elem["poster"] = f"images/{poster_name}"
    except Exception as e:
        print(f"Warning: Failed to extract video: {e}", file=sys.stderr)
    return elem


def _extract_svg_blob(shape):
    """Extract SVG bytes from asvg:svgBlip if present. Returns bytes or None."""
    ASVG_NS = 'http://schemas.microsoft.com/office/drawing/2016/SVG/main'
    svg_blip = shape._element.find(f'.//{{{ASVG_NS}}}svgBlip')
    if svg_blip is None:
        return None
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    r_embed = svg_blip.get(f'{{{r_ns}}}embed')
    if not r_embed:
        return None
    try:
        part = shape.part.rels[r_embed].target_part
        return part.blob
    except (KeyError, Exception):
        return None


def extract_picture_element(shape, output_dir=None, slide_idx=0, img_idx=0, theme_colors=None, color_mapping=None):
    """Extract picture as element dict and save image file."""
    elem = _base_element(shape, "image")
    
    # Check for SVG (asvg:svgBlip)
    svg_bytes = _extract_svg_blob(shape)
    if svg_bytes is not None:
        if output_dir:
            images_dir = Path(output_dir) / "images"
            images_dir.mkdir(exist_ok=True)
            filename = f"slide{slide_idx + 1}_image{img_idx + 1}.svg"
            (images_dir / filename).write_bytes(svg_bytes)
            elem["src"] = f"images/{filename}"
        return elem
    
    # Save image to file
    if output_dir:
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Determine format
            ext = shape.image.ext or "png"
            
            # Create images directory
            images_dir = Path(output_dir) / "images"
            images_dir.mkdir(exist_ok=True)
            
            # Save image
            image_filename = f"slide{slide_idx + 1}_image{img_idx + 1}.{ext}"
            image_path = images_dir / image_filename
            
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            # Store relative path
            elem["src"] = f"images/{image_filename}"
        except Exception as e:
            print(f"Warning: Failed to save image: {e}", file=sys.stderr)
    
    # Extract hyperlink
    if hasattr(shape, 'click_action') and shape.click_action.hyperlink:
        elem["link"] = shape.click_action.hyperlink.address
    
    # Extract image effects into _originalEffects (underscore-prefixed so builder
    # ignores them by default).  When reusing images in new slides, agents should
    # NOT copy _originalEffects — this prevents unintended mask/crop/color changes.
    # To faithfully reproduce the original slide, spread _originalEffects into the
    # element: { ...elem, ...elem._originalEffects }.
    try:
        effects: dict = {}
        pic_el = shape._element
        sp_pr = pic_el.find(f'{{{_NS["p"]}}}spPr')
        if sp_pr is not None:
            # Mask (prstGeom != rect)
            prst_geom = sp_pr.find(f'{{{_NS["a"]}}}prstGeom')
            if prst_geom is not None:
                prst = prst_geom.get('prst')
                mask_rmap = {"ellipse": "circle", "roundRect": "rounded_rectangle", "hexagon": "hexagon", "diamond": "diamond", "triangle": "triangle", "pentagon": "pentagon", "star5": "star_5_point", "heart": "heart", "trapezoid": "trapezoid"}
                if prst and prst != 'rect' and prst in mask_rmap:
                    effects["mask"] = mask_rmap[prst]
            # Visual effects
            effects.update(_extract_visual_effects(sp_pr, theme_colors, color_mapping))

        blip_fill = pic_el.find(f'{{{_NS["p"]}}}blipFill')
        if blip_fill is not None:
            # Crop
            src_rect = blip_fill.find(f'{{{_NS["a"]}}}srcRect')
            if src_rect is not None:
                crop = {}
                for side in ('l', 't', 'r', 'b'):
                    v = src_rect.get(side)
                    if v and int(v) != 0:
                        key = {"l": "left", "t": "top", "r": "right", "b": "bottom"}[side]
                        crop[key] = int(v) / 1000
                if crop:
                    effects["crop"] = crop
            # Brightness/Contrast/Saturation
            blip = blip_fill.find(f'{{{_NS["a"]}}}blip')
            if blip is not None:
                lum = blip.find(f'{{{_NS["a"]}}}lum')
                if lum is not None:
                    b = lum.get('bright')
                    c = lum.get('contrast')
                    if b:
                        effects["brightness"] = round(int(b) / 1000)
                    if c:
                        effects["contrast"] = round(int(c) / 1000)
                sat = blip.find(f'{{{_NS["a"]}}}hsl')
                if sat is not None:
                    v = sat.get('sat')
                    if v:
                        effects["saturation"] = round(int(v) / 1000)
                duo = blip.find(f'{{{_NS["a"]}}}duotone')
                if duo is not None:
                    colors = []
                    for srgb in duo.findall(f'{{{_NS["a"]}}}srgbClr'):
                        colors.append(_hex(srgb))
                    if len(colors) >= 2:
                        effects["duotone"] = colors[:2]
                # Preserve blip effects XML for lossless roundtrip (biLevel, etc.)
                from lxml import etree as _et
                blip_effects = []
                for child in blip:
                    tag = child.tag.split('}')[-1]
                    if tag in ('biLevel', 'grayscl', 'clrChange', 'clrRepl'):
                        blip_effects.append(_et.tostring(child, encoding='unicode'))
                if blip_effects:
                    effects["_blipEffects"] = blip_effects
        if effects:
            elem["_originalEffects"] = effects
    except Exception:
        pass
    
    return elem

def _dispatch_shape(shape, theme_colors=None, color_mapping=None, theme_styles=None, output_dir=None, slide_idx=0, img_counter=0, builder_text_color=None, pptx_path=None):
    """Dispatch shape extraction by type. Returns (elem, img_counter)."""
    elem = None
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        elem, img_counter = extract_group_element(shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter, builder_text_color=builder_text_color)
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        elem = extract_textbox_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        elem = extract_picture_element(shape, output_dir, slide_idx, img_counter, theme_colors, color_mapping)
        if elem:
            img_counter += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape._element.tag.endswith('}pic'):
        elem = extract_picture_element(shape, output_dir, slide_idx, img_counter, theme_colors, color_mapping)
        if elem:
            img_counter += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        from .table import extract_table_element
        elem = extract_table_element(shape, theme_colors, color_mapping, pptx_path)
    elif hasattr(shape, 'has_chart') and shape.has_chart:
        from .chart import extract_chart_element
        elem = extract_chart_element(shape, theme_colors, color_mapping)
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        elem = extract_line_element(shape, theme_colors, color_mapping, theme_styles)
    elif shape.shape_type == 16:  # MEDIA (video)
        elem = extract_video_element(shape, output_dir, slide_idx, img_counter)
        if elem:
            img_counter += 1
    elif shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            elem = extract_freeform_element(shape, theme_colors, color_mapping, builder_text_color=builder_text_color)
        if not elem:
            # Check if AUTO_SHAPE is actually a line (prst=line)
            try:
                prst = shape._element.spPr.find(f'.//{{{_NS["a"]}}}prstGeom')
                if prst is not None and prst.get('prst') == 'line':
                    elem = extract_line_element(shape, theme_colors, color_mapping, theme_styles)
            except Exception:
                pass
        if not elem:
            elem = extract_shape_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    return elem, img_counter

def extract_group_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, output_dir=None, slide_idx=0, img_counter=0, builder_text_color=None):
    """Extract group as element dict with nested elements.
    
    Note: python-pptx returns absolute slide coordinates for grouped shapes.
    """
    elem = _base_element(shape, "group", elements=[])
    # Move rotation after elements for consistent key order
    rot = elem.pop("rotation", None)
    if rot is not None:
        elem["rotation"] = rot

    # Save raw XML for groups that can't be losslessly flattened
    # (rotated groups, or groups with many freeforms like SVG icons)
    has_freeforms = any(child._element.tag.endswith('}sp') and
                        child.shape_type == 5 for child in shape.shapes)
    if rot is not None or has_freeforms:
        try:
            from lxml import etree as _et
            elem["_groupXml"] = _et.tostring(shape._element, encoding='unicode')
        except Exception:
            pass
    
    # Extract group fill (for grpFill inheritance)
    grp_fill_color = None
    grp_fill_gradient = None
    grp_sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
    if grp_sp_pr is not None:
        sf = grp_sp_pr.find(f'{{{_NS["a"]}}}solidFill')
        gf = grp_sp_pr.find(f'{{{_NS["a"]}}}gradFill')
        if sf is not None:
            srgb = sf.find(f'{{{_NS["a"]}}}srgbClr')
            scheme = sf.find(f'{{{_NS["a"]}}}schemeClr')
            if srgb is not None:
                grp_fill_color = _hex(srgb)
            elif scheme is not None:
                from .color import _resolve_color_with_transforms
                grp_fill_color = _resolve_color_with_transforms(scheme, theme_colors, color_mapping)
        elif gf is not None:
            fill_info = _extract_fill_from_xml(grp_sp_pr, theme_colors, color_mapping)
            grp_fill_gradient = fill_info.get("gradient")
            # For gradient grpFill, preserve entire group XML for lossless roundtrip
            from lxml import etree as _et
            elem["_groupXml"] = _et.tostring(shape._element, encoding='unicode')
    
    # Extract each shape in the group
    for sub_shape in shape.shapes:
        try:
            sub_elem = None
            
            # Handle nested groups recursively
            sub_elem, img_counter = _dispatch_shape(sub_shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter, builder_text_color=builder_text_color)
            
            if sub_elem:
                # Resolve grpFill: if sub-element has fill=none but XML has grpFill, use group fill
                if sub_elem.get("fill") in (None, "none"):
                    sub_sp = sub_shape._element.find(f'{{{_NS["p"]}}}spPr')
                    if sub_sp is not None and sub_sp.find(f'{{{_NS["a"]}}}grpFill') is not None:
                        if grp_fill_color:
                            sub_elem["fill"] = grp_fill_color
                        elif grp_fill_gradient:
                            sub_elem["gradient"] = grp_fill_gradient
                # Propagate grpFill to nested group children
                if sub_elem.get("type") == "group":
                    grp_sp = sub_shape._element.find(f'{{{_NS["p"]}}}grpSpPr')
                    if grp_sp is not None and grp_sp.find(f'{{{_NS["a"]}}}grpFill') is not None:
                        for child_el in sub_elem.get("elements", []):
                            if child_el.get("fill") in (None, "none"):
                                if grp_fill_color:
                                    child_el["fill"] = grp_fill_color
                                elif grp_fill_gradient:
                                    child_el["gradient"] = grp_fill_gradient

                # Transform coordinates from child coordinate system to slide coordinates
                grp_sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
                xfrm = grp_sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm') if grp_sp_pr else None
                
                if xfrm is not None:
                    off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                    ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                    ch_off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chOff')
                    ch_ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chExt')
                    
                    if off is not None and ch_off is not None and ext is not None and ch_ext is not None:
                        group_off_x = int(off.get('x'))
                        group_off_y = int(off.get('y'))
                        group_ext_cx = int(ext.get('cx'))
                        group_ext_cy = int(ext.get('cy'))
                        ch_off_x = int(ch_off.get('x'))
                        ch_off_y = int(ch_off.get('y'))
                        ch_ext_cx = int(ch_ext.get('cx'))
                        ch_ext_cy = int(ch_ext.get('cy'))
                        
                        # Transform: abs = group_off + (child - chOff) * (group_ext / ch_ext)
                        child_x = sub_shape.left
                        child_y = sub_shape.top
                        
                        scale_x = group_ext_cx / ch_ext_cx if ch_ext_cx != 0 else 1
                        scale_y = group_ext_cy / ch_ext_cy if ch_ext_cy != 0 else 1
                        
                        abs_x = group_off_x + (child_x - ch_off_x) * scale_x
                        abs_y = group_off_y + (child_y - ch_off_y) * scale_y
                        
                        sub_elem["x"] = round(abs_x / EMU_PER_PX)
                        sub_elem["y"] = round(abs_y / EMU_PER_PX)
                        sub_elem["width"] = round(sub_shape.width * scale_x / EMU_PER_PX)
                        sub_elem["height"] = round(sub_shape.height * scale_y / EMU_PER_PX)
                        # For freeform in group: drop raw path XML, let builder reconstruct
                        # from px coords (which match the group-scaled shape size)
                        if sub_elem.get("type") == "freeform":
                            sub_elem["_xEmu"] = round(abs_x)
                            sub_elem["_yEmu"] = round(abs_y)
                            sub_elem["_widthEmu"] = round(sub_shape.width * scale_x)
                            sub_elem["_heightEmu"] = round(sub_shape.height * scale_y)
                            sub_elem.pop("_pathLstXml", None)
                        
                        # For nested groups, also transform all children recursively
                        if sub_elem.get("type") == "group" and sub_elem.get("elements"):
                            def _apply_group_transform(elements, gox, goy, gcx, gcy, sx, sy):
                                for el in elements:
                                    if "x" in el and "y" in el:
                                        old_x = el["x"] * EMU_PER_PX
                                        old_y = el["y"] * EMU_PER_PX
                                        new_x = gox + (old_x - gcx) * sx
                                        new_y = goy + (old_y - gcy) * sy
                                        el["x"] = round(new_x / EMU_PER_PX)
                                        el["y"] = round(new_y / EMU_PER_PX)
                                        if el.get("type") == "freeform":
                                            el["_xEmu"] = round(new_x)
                                            el["_yEmu"] = round(new_y)
                                    if "width" in el:
                                        el["width"] = round(el["width"] * sx)
                                    if "height" in el:
                                        el["height"] = round(el["height"] * sy)
                                    if el.get("type") == "freeform":
                                        if el.get("_widthEmu"):
                                            el["_widthEmu"] = round(el["_widthEmu"] * sx)
                                        else:
                                            el["_widthEmu"] = round(el["width"] * EMU_PER_PX)
                                        if el.get("_heightEmu"):
                                            el["_heightEmu"] = round(el["_heightEmu"] * sy)
                                        else:
                                            el["_heightEmu"] = round(el["height"] * EMU_PER_PX)
                                        el.pop("_pathLstXml", None)
                                    if el.get("type") == "group" and el.get("elements"):
                                        _apply_group_transform(el["elements"], gox, goy, gcx, gcy, sx, sy)
                            _apply_group_transform(sub_elem["elements"], group_off_x, group_off_y, ch_off_x, ch_off_y, scale_x, scale_y)
                    else:
                        # Fallback: use python-pptx coordinates as-is
                        pass
                
                # Propagate group rotation to child elements
                if shape.rotation != 0:
                    child_rot = sub_elem.get("rotation", 0)
                    sub_elem["rotation"] = round(child_rot + shape.rotation, 1)
                
                elem["elements"].append(sub_elem)
        except Exception as e:
            print(f"Warning: Failed to extract grouped shape: {e}", file=sys.stderr)
    
    return elem, img_counter
