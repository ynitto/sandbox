# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Reference: unified access to design patterns and example datasets.

Reference data is used to guide AI-generated slide content.
See docs/internal/DATASET_COMPLIANCE.md for dataset provenance.

Reference: unified access to design patterns, components, schemas, rules, and reviews."""
from pathlib import Path

from sdpm.reference.providers import FileProvider, ReferenceProvider  # noqa: F401


def _get_description(path: Path) -> str:
    """Extract description from first non-empty line of md or pptx speaker notes."""
    if path.suffix == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            if prs.slides:
                slide = prs.slides[0]
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    text = slide.notes_slide.notes_text_frame.text.replace('\x0B', '\n')
                    for line in text.splitlines():
                        if line.strip():
                            return line.strip()
        except Exception:
            pass
        return ""
    # md / html
    try:
        text = path.read_text(encoding="utf-8")
        # HTML — extract <title>
        if path.suffix == '.html':
            import re
            m = re.search(r'<title>(.*?)</title>', text, re.IGNORECASE)
            if m:
                return m.group(1).strip()
            return ""
        lines = text.splitlines()
        # Parse YAML frontmatter for description
        if lines and lines[0].strip() == '---':
            for i, line in enumerate(lines[1:], 1):
                if line.strip() == '---':
                    # Extract description from frontmatter
                    for fm_line in lines[1:i]:
                        if fm_line.startswith('description:'):
                            val = fm_line[len('description:'):].strip().strip('"').strip("'")
                            if val:
                                return val
                    # No description in frontmatter — fall through to first content line
                    lines = lines[i + 1:]
                    break
        for line in lines:
            if line.strip():
                return line.strip()
    except Exception:
        pass
    return ""


def generate_styles_index(styles_dir: Path) -> Path:
    """Generate an index HTML with sidebar + iframe preview for style HTMLs.

    Returns the path to the generated index file.
    """
    import html as html_mod
    import tempfile

    styles = []
    for f in sorted(styles_dir.iterdir()):
        if f.suffix == '.html' and not f.name.startswith('.'):
            desc = _get_description(f)
            # Split "name — description" from title
            name = f.stem
            subtitle = ""
            if " — " in desc:
                _, subtitle = desc.split(" — ", 1)
            elif desc:
                subtitle = desc
            styles.append((name, subtitle, f.resolve()))

    if not styles:
        return Path()

    first_src = styles[0][2].as_uri()

    items_html = ""
    for i, (name, subtitle, path) in enumerate(styles):
        active = ' active' if i == 0 else ""
        esc_name = html_mod.escape(name)
        esc_sub = html_mod.escape(subtitle)
        src = html_mod.escape(path.as_uri())
        items_html += (
            f'<div class="item{active}" onclick="select(this,\'{src}\')">'
            f'<div class="name">{esc_name}'
            f'<button class="copy" onclick="event.stopPropagation();copy(this,\'{esc_name}\')" title="Copy name">⧉</button>'
            f'</div>'
            f'<div class="sub">{esc_sub}</div>'
            f'</div>\n'
        )

    index_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>Style Gallery</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{display:flex;height:100vh;font-family:-apple-system,system-ui,"Segoe UI",sans-serif;background:#111;color:#ccc}}
#sidebar{{width:300px;min-width:300px;background:#161616;display:flex;flex-direction:column}}
#sidebar header{{padding:24px 24px 20px;border-bottom:1px solid #222}}
#sidebar header h1{{font-size:13px;font-weight:500;letter-spacing:.08em;text-transform:uppercase;color:#555}}
#list{{flex:1;overflow-y:auto;padding:8px 12px}}
#list::-webkit-scrollbar{{width:4px}}
#list::-webkit-scrollbar-thumb{{background:#333;border-radius:2px}}
.item{{padding:12px 14px;border-radius:8px;cursor:pointer;transition:all .2s ease-out;margin-bottom:2px}}
.item:hover{{background:#1e1e1e}}
.item.active{{background:#1a2332}}
.item .name{{font-size:15px;font-weight:600;color:#fff;margin-bottom:4px;transition:color .2s;display:flex;align-items:center;gap:8px}}
.item.active .name{{color:#79b8ff}}
.copy{{background:none;border:1px solid #555;color:#aaa;font-size:14px;border-radius:4px;cursor:pointer;padding:2px 6px;opacity:0;transition:opacity .15s,color .15s,border-color .15s;flex-shrink:0;line-height:1}}
.item:hover .copy{{opacity:1}}
.copy:hover{{color:#fff;border-color:#888}}
.copy.ok{{color:#56d364;border-color:#56d364}}
.item .sub{{font-size:12px;color:#aaa;line-height:1.5;transition:color .2s}}
.item.active .sub{{color:#a5cdee}}
#main{{flex:1;overflow:hidden;position:relative;background:#0a0a0a}}
iframe{{position:absolute;top:0;left:0;width:1920px;border:none;background:#fff;transform-origin:top left}}
</style></head><body>
<div id="sidebar">
<header><h1>Style Gallery</h1></header>
<div id="list">{items_html}</div>
</div>
<div id="main"><iframe id="preview" src="{html_mod.escape(first_src)}"></iframe></div>
<script>
function select(el,src){{
  document.querySelector('#list .active')?.classList.remove('active');
  el.classList.add('active');
  document.getElementById('preview').src=src;
}}
function copy(btn,name){{
  navigator.clipboard.writeText(name);
  btn.textContent='✓';btn.classList.add('ok');
  setTimeout(function(){{btn.textContent='⧉';btn.classList.remove('ok')}},1200);
}}
function fit(){{
  var m=document.getElementById('main'),f=document.getElementById('preview');
  var s=m.clientWidth/1920;
  f.style.transform='scale('+s+')';
  f.style.height=(m.clientHeight/s)+'px';
}}
window.addEventListener('resize',fit);
window.addEventListener('load',fit);
new ResizeObserver(fit).observe(document.getElementById('main'));
</script></body></html>"""

    out = Path(tempfile.gettempdir()) / "sdpm-styles-index.html"
    out.write_text(index_html, encoding="utf-8")
    return out


def search_patterns(query: str, limit: int = 0) -> list[dict]:
    """Search pattern notes by keywords.

    Searches speaker notes across all pptx files in references/examples/
    and md files in references/examples/styles/.
    Returns matching entries with their 1st-line description.
    limit=0 means no limit (return all matches).
    """
    import re

    from pptx import Presentation

    examples_dir = Path(__file__).parent.parent.parent / "references" / "examples"
    if not examples_dir.exists():
        return []

    queries = query.lower().split()
    pats = [re.compile(r'\b' + re.escape(q) + r'\b') for q in queries]
    results: list[tuple[int, str, int, str]] = []

    # Search pptx files (patterns only)
    patterns_pptx = examples_dir / "patterns.pptx"
    if patterns_pptx.exists():
        try:
            prs = Presentation(str(patterns_pptx))
        except Exception:
            prs = None
        if prs:
            for si, slide in enumerate(prs.slides):
                if not slide.has_notes_slide:
                    continue
                notes = slide.notes_slide.notes_text_frame.text.replace('\x0B', '\n')
                if not notes.strip():
                    continue
                notes_lower = notes.lower()
                match_count = sum(1 for p in pats if p.search(notes_lower))
                if match_count == 0:
                    continue
                desc = ""
                for line in notes.splitlines():
                    if line.strip():
                        desc = line.strip()
                        break
                results.append((match_count, "patterns", si + 1, desc))

    results.sort(key=lambda x: (-x[0], x[1], x[2]))
    out = []
    for r in results:
        entry: dict = {"path": r[1], "description": r[3]}
        if r[2] > 0:
            entry["page"] = r[2]
        out.append(entry)
    return out[:limit] if limit else out


def list_pptx_descriptions(pptx_path):
    """List all slide descriptions (speaker notes line 1) from a pptx file."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    descriptions = []
    for slide_idx, slide in enumerate(prs.slides):
        desc = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text.replace('\x0B', '\n')
            for line in text.splitlines():
                if line.strip():
                    desc = line.strip()
                    break
        descriptions.append((slide_idx + 1, desc or f"(slide {slide_idx + 1})"))
    return descriptions


def get_pptx_notes(pptx_path, pages=None):
    """Extract speaker notes from pptx slides.

    Args:
        pptx_path: Path to the pptx file.
        pages: Optional list of 1-based page numbers. None = all pages.
    Returns:
        List of (page_number, notes_text) tuples.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    results = []
    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1
        if pages and page_num not in pages:
            continue
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip().replace('\x0B', '\n')
        if notes:
            results.append((page_num, notes))
    return results


# ---------------------------------------------------------------------------
# High-level reference helpers (local filesystem)
# ---------------------------------------------------------------------------


def open_styles_gallery(styles_dir: Path) -> Path | None:
    """Generate styles gallery HTML and open in browser. Returns index path or None."""
    import webbrowser

    index_path = generate_styles_index(styles_dir)
    if index_path and index_path.exists():
        webbrowser.open(index_path.as_uri())
        return index_path
    return None


def _collect_files(directory: Path) -> dict[str, Path]:
    """Collect md/pptx files in directory. md takes priority over pptx for same stem."""
    files: dict[str, Path] = {}
    for f in sorted(directory.glob("*.pptx")):
        files[f.stem] = f
    for f in sorted(directory.glob("*.md")):
        files[f.stem] = f  # md overwrites pptx
    return files


def _strip_frontmatter(text: str) -> str:
    """Strip YAML frontmatter from markdown text."""
    lines = text.splitlines(True)
    if lines and lines[0].strip() == "---":
        for i, line in enumerate(lines[1:], 1):
            if line.strip() == "---":
                return "".join(lines[i + 1:]).lstrip("\n")
    return text


def list_category(category_dir: Path) -> list[dict[str, str]]:
    """List documents in a category directory with descriptions.

    Supports flat and subdirectory layouts (matching CLI _list_or_show).
    md takes priority over pptx for same stem.
    """
    if not category_dir.exists():
        return []
    subdirs = sorted(d for d in category_dir.iterdir() if d.is_dir())
    flat = not subdirs
    items: list[dict[str, str]] = []
    if flat:
        for stem, f in sorted(_collect_files(category_dir).items()):
            items.append({"name": stem, "description": _get_description(f)})
    else:
        for sd in subdirs:
            for stem, f in sorted(_collect_files(sd).items()):
                items.append({"name": f"{sd.name}/{stem}", "description": _get_description(f)})
    return items


def read_docs(category_dir: Path, names: list[str]) -> list[dict]:
    """Read documents by name from a category directory.

    Supports .md (with frontmatter strip), .pptx (notes via get_pptx_notes).
    Names can include page specifiers for pptx: "name/3" or "name/all".
    """
    results = []
    for name in names:
        pages = None
        has_page_specifier = False
        parts = name.rsplit("/", 1)
        file_name = name
        if len(parts) == 2 and (parts[1].isdigit() or parts[1] == "all"):
            file_name = parts[0]
            has_page_specifier = True
            pages = None if parts[1] == "all" else [int(parts[1])]

        md_path = category_dir / f"{file_name}.md"
        pptx_path = category_dir / f"{file_name}.pptx"

        if md_path.exists():
            text = _strip_frontmatter(md_path.read_text(encoding="utf-8"))
            results.append({"name": file_name, "content": text})
        elif pptx_path.exists():
            if not has_page_specifier:
                descriptions = list_pptx_descriptions(str(pptx_path))
                content = "\n".join(f"  {page:>3}  {desc}" for page, desc in descriptions)
            else:
                notes = get_pptx_notes(pptx_path, pages=pages)
                content = "\n".join(f"## Page {pn}\n\n{text}\n" for pn, text in notes)
            results.append({"name": name, "content": content})
        else:
            # Collect available names for error message
            available = sorted(_collect_files(category_dir).keys())
            raise FileNotFoundError(
                f"'{file_name}' not found in {category_dir.name}/. Available: {', '.join(available)}"
            )
    return results


def list_styles(styles_dir: Path) -> list[dict[str, str]]:
    """List available styles from HTML files in a styles directory.

    Returns list of dicts with name and description.
    """
    if not styles_dir.exists():
        return []
    styles: list[dict[str, str]] = []
    for f in sorted(styles_dir.iterdir()):
        if f.suffix != ".html" or f.name.startswith("."):
            continue
        desc = _get_description(f)
        styles.append({"name": f.stem, "description": desc})
    return styles
