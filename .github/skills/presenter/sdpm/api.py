# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""High-level API for sdpm — single entry points for generate, measure, preview, init, code_block.

These functions encapsulate the full workflow that the CLI (pptx_builder.py) performs.
mcp-local and other consumers should call these instead of assembling low-level APIs.
"""

import os
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any


def _resolve_template(data: dict, input_path: str | Path | None, templates_dir: Path) -> tuple[Path, bool]:
    """Resolve template path from presentation data.

    Returns (template_path, custom_template) or raises FileNotFoundError.
    """
    if data.get("template"):
        base_dir = Path(input_path).parent if input_path else Path(".")
        template = base_dir / data["template"]
        if template.exists():
            return template, True
        name = data["template"]
        named = templates_dir / (name if name.endswith(".pptx") else name + ".pptx")
        if named.exists():
            return named, True
    raise FileNotFoundError('No template specified. Set "template" in presentation JSON.')


def _get_output_base_dir() -> Path:
    """Get output base directory from config, with WSL fallback."""
    env_dir = os.environ.get("SDPM_OUTPUT_DIR")
    if env_dir:
        return Path(env_dir)
    try:
        from sdpm.config import get_output_dir

        return get_output_dir()
    except Exception:
        pass
    from sdpm.preview.backend import _is_wsl

    if _is_wsl():
        import subprocess

        try:
            result = subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                [
                    "/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe",
                    "-Command",
                    "[Environment]::GetFolderPath('MyDocuments')",
                ],
                capture_output=True,
                timeout=10,
            )
            win_path = result.stdout.decode("cp932", errors="replace").strip()
            if win_path:
                wsl = subprocess.run(["wslpath", win_path], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                if wsl.returncode == 0:
                    return Path(wsl.stdout.strip()) / "SDPM-Presentations"
        except Exception:
            pass
    return Path.home() / "Documents" / "SDPM-Presentations"


def init(
    name: str,
    template: str | Path | None = None,
    output_dir: str | Path | None = None,
) -> dict[str, Any]:
    """Initialize a presentation workspace.

    Creates output directory with presentation.json and specs/.

    Args:
        name: Presentation name (used in directory name).
        template: Template name or path. If provided, extracts fonts.
        output_dir: Explicit output directory. Auto-generated if None.

    Returns:
        Dict with output_dir, json_path, template, fonts, workspace.
    """
    from sdpm.analyzer import extract_fonts
    from sdpm.utils.io import write_json

    if output_dir:
        out_dir = Path(output_dir).expanduser()
    else:
        ts = datetime.now().strftime("%Y%m%d-%H%M")
        dir_name = f"{ts}-{name}" if name else ts
        out_dir = _get_output_base_dir() / dir_name
    out_dir.mkdir(parents=True, exist_ok=True)

    pres_data: dict[str, Any] = {"fonts": {"fullwidth": None, "halfwidth": None}, "slides": []}

    if template:
        templates_dir = Path(__file__).parent.parent / "templates"  # skill/templates/
        template_src = Path(template).expanduser()
        if not template_src.exists():
            candidate = templates_dir / (str(template) if str(template).endswith(".pptx") else f"{template}.pptx")
            if candidate.exists():
                template_src = candidate
        if template_src.exists():
            template_src = template_src.resolve()
            pres_data["template"] = template_src.name
            pres_data["fonts"] = extract_fonts(template_src)

    json_path = out_dir / "presentation.json"
    write_json(json_path, pres_data, suffix="\n")

    specs_dir = out_dir / "specs"
    specs_dir.mkdir(exist_ok=True)
    spec_files = ("brief.md", "outline.md")
    for spec_name in spec_files:
        (specs_dir / spec_name).touch()

    return {
        "output_dir": str(out_dir),
        "json_path": str(json_path),
        "template": pres_data.get("template", ""),
        "fonts": pres_data.get("fonts", {}),
        "workspace": ["presentation.json"] + [f"specs/{s}" for s in spec_files],
    }


@dataclass
class BuildConfig:
    """Resolved configuration for PPTX build."""

    template_path: Path
    custom_template: bool
    fonts: dict
    default_text_color: str
    slides: list[dict] = field(default_factory=list)  # override解決済み
    base_dir: Path = field(default_factory=lambda: Path("."))
    warnings: list[str] = field(default_factory=list)
    lint_diagnostics: list = field(default_factory=list)


def _resolve_config(json_path: str | Path) -> BuildConfig:
    """Resolve template, fonts, icons, overrides from JSON.

    Raises FileNotFoundError, ValueError on missing template/icons.
    """
    from sdpm.builder import PPTXBuilder, resolve_override, validate_icons_in_json
    from sdpm.utils.io import read_json

    input_path = Path(json_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Slides JSON not found: {json_path}")

    data = read_json(input_path)
    templates_dir = Path(__file__).parent.parent / "templates"
    warnings: list[str] = []

    template_file, custom = _resolve_template(data, str(input_path), templates_dir)

    # Auto-fill fonts
    from sdpm.analyzer import extract_fonts as _extract_fonts

    fonts = data.get("fonts")
    if not fonts or not fonts.get("fullwidth"):
        fonts = _extract_fonts(template_file)
        warnings.append("fonts auto-detected from template")

    # Auto-fill defaultTextColor
    dtc = data.get("defaultTextColor")
    if not dtc:
        _, is_dark = PPTXBuilder._extract_theme_colors(template_file)
        dtc = "#FFFFFF" if is_dark else "#333333"
        warnings.append(f"defaultTextColor auto-set to {dtc}")

    # Lint
    from sdpm.schema.lint import lint as lint_slides

    lint_diagnostics = lint_slides(data)

    # Validate icons
    missing = validate_icons_in_json(data)
    if missing:
        raise ValueError(f"Missing assets ({len(missing)}): {', '.join(sorted(missing)[:10])}")

    # Resolve overrides
    slides = data.get("slides", [])
    id_map = {}
    for s in slides:
        if "id" in s:
            id_map[s["id"]] = s
    resolved_slides = [resolve_override(s, id_map) for s in slides]

    return BuildConfig(
        template_path=template_file,
        custom_template=custom,
        fonts=fonts,
        default_text_color=dtc,
        slides=resolved_slides,
        base_dir=input_path.parent,
        warnings=warnings,
        lint_diagnostics=lint_diagnostics,
    )


def _build(config: BuildConfig, output_path: Path) -> Path:
    """Build PPTX from resolved config. Returns output path."""
    from sdpm.builder import PPTXBuilder

    builder = PPTXBuilder(
        config.template_path,
        custom_template=config.custom_template,
        fonts=config.fonts,
        base_dir=config.base_dir,
        default_text_color=config.default_text_color,
    )
    for s in config.slides:
        builder.add_slide(s)
    builder.save(output_path)
    return output_path


def generate(
    json_path: str | Path,
    output_path: str | Path | None = None,
) -> dict[str, Any]:
    """Generate PPTX from JSON.

    Includes: template resolution, icon validation, build, check_layout_imbalance.

    Args:
        json_path: Path to the slides JSON file.
        output_path: Output .pptx path. Auto-generated if None.

    Returns:
        Dict with output_path, slide_count, slides summary, warnings.
    """
    from sdpm.preview import check_layout_imbalance_data

    config = _resolve_config(json_path)

    # Output path
    input_path = Path(json_path)
    if not output_path:
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = input_path.with_suffix(".pptx")
        out = p.with_stem(f"{p.stem}_{ts}")
    else:
        out = Path(output_path)

    _build(config, out)

    imbalance = check_layout_imbalance_data(out, config.slides)
    if imbalance:
        for a in imbalance:
            config.warnings.append(f"page{a['slide']:02d} ({a['layout']}) offset: {a['offset']} ({a['direction']})")

    # Summary
    summary = []
    for i, s in enumerate(config.slides, 1):
        title = s.get("title", "(no title)")
        if isinstance(title, dict):
            title = title.get("text", "(no title)")
        summary.append(f"page{i:02d} - {title}")

    result: dict[str, Any] = {
        "output_path": str(out),
        "slide_count": len(config.slides),
        "slides": summary,
        "warnings": config.warnings,
    }
    if config.lint_diagnostics:
        result["errors"] = {"lintDiagnostics": config.lint_diagnostics}
    return result


def measure(
    json_path: str | Path,
    slides: list[int] | None = None,
) -> str:
    """Build PPTX from JSON, convert to SVG, extract text bboxes.

    Args:
        json_path: Path to the slides JSON file.
        slides: Slide numbers to measure (1-based). None for all.

    Returns:
        Text report of bbox measurements.
    """
    import tempfile

    from sdpm.preview.backend import LibreOfficeBackend
    from sdpm.preview.measure import format_measure_report, measure_from_svg

    config = _resolve_config(json_path)

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_pptx = Path(tmp_dir) / "measure.pptx"
        _build(config, tmp_pptx)

        backend = LibreOfficeBackend()
        svg_path = backend.export_svg(tmp_pptx)
        if svg_path is None:
            raise RuntimeError("SVG export failed. Is LibreOffice (soffice) installed?")

        try:
            results = measure_from_svg(svg_path, slides)
            return format_measure_report(results)
        finally:
            import shutil
            if svg_path:
                shutil.rmtree(svg_path.parent, ignore_errors=True)


def preview(
    json_path: str | Path,
    output_path: str | Path | None = None,
    pages: list[int] | None = None,
    grid: bool = False,
) -> dict[str, Any]:
    """Build PPTX from JSON and export slides as PNG images.

    Args:
        json_path: Path to the slides JSON file.
        output_path: Output .pptx path. Auto-generated if None.
        pages: Page numbers to export. None for all.
        grid: Add grid overlay to PNGs.

    Returns:
        Dict with preview_dir, files list, and output_path.
    """
    import glob
    import re
    import subprocess

    from pptx import Presentation

    from sdpm.preview import export_pdf

    config = _resolve_config(json_path)

    # Output path
    input_path = Path(json_path)
    if not output_path:
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = input_path.with_suffix(".pptx")
        out = p.with_stem(f"{p.stem}_{ts}")
    else:
        out = Path(output_path)

    _build(config, out)

    # Preview dir
    out_dir = Path("/tmp/pptx-preview")
    out_dir.mkdir(parents=True, exist_ok=True)

    pages_set = set(pages) if pages else None

    # PDF + pdftoppm pipeline
    pdf = out_dir / "slides.pdf"
    if not export_pdf(out, pdf):
        raise RuntimeError("PDF export failed. Is LibreOffice (soffice) installed?")

    cmd = ["pdftoppm", "-png", "-scale-to", "1280", str(pdf), str(out_dir / "page")]
    result = subprocess.run(cmd, capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    if result.returncode != 0:
        raise RuntimeError(f"PNG conversion failed. Is poppler (pdftoppm) installed? {result.stderr}")

    # Rename with slide titles
    prs = Presentation(str(out))
    titles = _extract_slide_titles(prs)

    generated = []
    for png in sorted(glob.glob(str(out_dir / "page-*.png"))):
        match = re.match(r"page-(\d+)\.png", Path(png).name)
        if match:
            num = int(match.group(1))
            if pages_set and num not in pages_set:
                Path(png).unlink()
                continue
            new_name = f"page{num:02d}-{titles.get(num, 'notitle')}.png"
            new_path = out_dir / new_name
            Path(png).rename(new_path)
            generated.append(str(new_path))

    pdf.unlink(missing_ok=True)

    if grid:
        _apply_grid_overlay(generated)

    return {"preview_dir": str(out_dir), "files": generated, "output_path": str(out)}


def _extract_slide_titles(prs) -> dict[int, str]:
    """Extract sanitized slide titles from a Presentation object."""
    import re

    titles = {}
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
        title = re.sub(r'[\\/:*?"<>|]', "", title)
        titles[i] = title or "notitle"
    return titles


def _apply_grid_overlay(png_paths: list[str]) -> None:
    """Add coordinate grid overlay to PNG files."""
    from PIL import Image, ImageDraw, ImageFont

    color = (255, 0, 0, 128)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except Exception:
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except Exception:
            font = ImageFont.load_default()

    for png_path in png_paths:
        img = Image.open(png_path).convert("RGBA")
        w, h = img.size
        overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        for pct in range(5, 100, 5):
            x, y = int(w * pct / 100), int(h * pct / 100)
            px_x, px_y = int(1920 * pct / 100), int(1080 * pct / 100)
            draw.line([(x, 0), (x, h)], fill=color, width=1)
            draw.line([(0, y), (w, y)], fill=color, width=1)
            if pct % 10 == 0:
                draw.text((x + 4, 4), f"{px_x}px ({pct}%)", fill=color, font=font)
                draw.text((4, y + 4), f"{px_y}px ({pct}%)", fill=color, font=font)
        Image.alpha_composite(img, overlay).convert("RGB").save(png_path)


def code_block(
    code: str,
    language: str = "python",
    theme: str = "dark",
    x: int = 0,
    y: int = 0,
    width: int = 800,
    height: int = 300,
    font_size: int = 12,
    show_label: bool = True,
) -> list[dict[str, Any]]:
    """Generate slide elements for a syntax-highlighted code block.

    Args:
        code: Source code text.
        language: Programming language for highlighting.
        theme: Color theme ("dark" or "light").
        x, y, width, height: Position and size in pixels.
        font_size: Code font size in pt.
        show_label: Show language label bar.

    Returns:
        List of element dicts for slide JSON.
    """
    from sdpm.builder.constants import CODE_COLORS
    from sdpm.utils.text import highlight_code

    colors = CODE_COLORS.get(theme, CODE_COLORS["dark"])
    bg = colors["background"]
    inverse_theme = "light" if theme == "dark" else "dark"
    inverse_bg = CODE_COLORS[inverse_theme]["background"]
    label_fg = "000000" if theme == "dark" else "FFFFFF"
    label_height = 22

    label_map = {"typescript": "TypeScript", "javascript": "JavaScript", "csharp": "C#", "cpp": "C++"}
    label_text = label_map.get(language, language.capitalize())

    elements: list[dict[str, Any]] = []
    if show_label:
        elements.append(
            {
                "type": "textbox",
                "x": x,
                "y": y,
                "width": width,
                "height": label_height,
                "fontSize": 8,
                "align": "left",
                "fill": inverse_bg,
                "text": f"{{{{#{label_fg}:{label_text}}}}}",
                "marginLeft": 50000,
                "marginTop": 0,
                "marginRight": 0,
                "marginBottom": 0,
                "autoWidth": True,
            }
        )
        code_y = y + label_height
        code_height = height - label_height
    else:
        code_y = y
        code_height = height

    spans = highlight_code(code, language, theme)
    elements.append(
        {
            "type": "textbox",
            "x": x,
            "y": code_y,
            "width": width,
            "height": code_height,
            "fontSize": font_size,
            "align": "left",
            "fill": bg,
            "text": spans,
        }
    )

    return elements
